from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "knowledge_base_architecture_diagrams.pptx"

W = 13.333
H = 7.5
FONT = "Noto Sans CJK TC"
DISPLAY = "Noto Serif CJK TC"

C = {
    "navy": "0B1F33",
    "navy2": "12314B",
    "blue": "2F6FED",
    "blue2": "5B8DEF",
    "blue_soft": "EAF1FF",
    "teal": "00A6A6",
    "teal2": "007F7F",
    "teal_soft": "E4F7F6",
    "green": "2E9D68",
    "green_soft": "E6F6EE",
    "amber": "F4B740",
    "amber2": "A96B00",
    "amber_soft": "FFF4D6",
    "red": "D94F4F",
    "red_soft": "FDECEC",
    "ink": "172B3A",
    "slate": "52677A",
    "muted": "71859A",
    "line": "CCD8E3",
    "panel": "F4F7FA",
    "white": "FFFFFF",
    "bg": "F8FAFC",
}


def rgb(value):
    return RGBColor.from_string(value.replace("#", "").upper())


def rect(slide, x, y, w, h, fill, line=None, radius=True, transparency=0, width=1):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()
    return shape


def text(
    slide,
    x,
    y,
    w,
    h,
    value,
    size=10,
    color=None,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    font=None,
    fit=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if fit else MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = value
    run.font.name = font or FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color or C["ink"])
    return box


def pill(slide, x, y, w, value, fill, color, h=0.26, size=7.5, line=None):
    shape = rect(slide, x, y, w, h, fill, line=line or fill, radius=True)
    shape.text_frame.clear()
    shape.text_frame.word_wrap = False
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    shape.text_frame.margin_left = Pt(3)
    shape.text_frame.margin_right = Pt(3)
    shape.text_frame.margin_top = Pt(0)
    shape.text_frame.margin_bottom = Pt(0)
    p = shape.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    return shape


def background(slide):
    rect(slide, 0, 0, W, H, C["bg"], radius=False)
    rect(slide, 0, 0, W, 0.07, C["navy"], radius=False)
    rect(slide, 0, 0, 0.16, H, C["blue"], radius=False)


def header(slide, section, title, subtitle, number):
    text(slide, 0.48, 0.25, 4.8, 0.18, section.upper(), 8.5, C["blue"], True)
    text(slide, 0.48, 0.48, 10.9, 0.43, title, 22, C["navy"], True, font=DISPLAY, fit=True)
    text(slide, 0.49, 0.92, 11.45, 0.25, subtitle, 9.5, C["slate"], fit=True)
    pill(slide, 11.98, 0.31, 0.72, f"{number:02d}", C["navy"], C["white"], h=0.34, size=9)


def footer(slide, note):
    rect(slide, 0.48, 7.02, 12.22, 0.01, C["line"], radius=False)
    text(slide, 0.48, 7.08, 10.4, 0.17, f"架構重點｜{note}", 7.6, C["slate"], True, fit=True)
    text(slide, 10.95, 7.08, 1.75, 0.17, "AS-IS · 2026.07.20", 7.0, C["muted"], align=PP_ALIGN.RIGHT, fit=True)


def zone(slide, x, y, w, h, label, accent, fill=None, dashed=False):
    shape = rect(slide, x, y, w, h, fill or C["white"], line=accent, radius=True, width=1.3)
    if dashed:
        shape.line.dash_style = 2
    pill(slide, x + 0.12, y + 0.10, min(1.55, w - 0.24), label, accent, C["white"], h=0.26, size=7.0)
    return shape


def node(
    slide,
    x,
    y,
    w,
    h,
    title,
    subtitle,
    accent,
    icon=None,
    fill=None,
    title_size=9.5,
    subtitle_size=7.2,
):
    rect(slide, x, y, w, h, fill or C["white"], line=C["line"], radius=True)
    rect(slide, x, y, 0.06, h, accent, radius=False)
    tx = x + 0.16
    tw = w - 0.27
    show_icon = bool(icon) and (w >= 1.70 or str(icon).isdigit())
    if show_icon:
        rect(slide, x + 0.15, y + 0.13, 0.38, 0.38, accent, radius=True)
        text(slide, x + 0.15, y + 0.13, 0.38, 0.38, icon, 7.5, C["white"], True, PP_ALIGN.CENTER, MSO_VERTICAL_ANCHOR.MIDDLE)
        tx = x + 0.62
        tw = w - 0.76
    text(slide, tx, y + 0.11, tw, 0.28, title, title_size, accent, True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
    text(slide, x + 0.15, y + 0.49, w - 0.28, h - 0.58, subtitle, subtitle_size, C["slate"], valign=MSO_VERTICAL_ANCHOR.TOP, fit=True)


def store(slide, x, y, w, h, title, subtitle, accent, icon):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CAN, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(C["white"])
    shape.line.color.rgb = rgb(accent)
    shape.line.width = Pt(1.2)
    if w < 1.55:
        text(slide, x + 0.10, y + 0.12, w - 0.20, 0.27, title, 8.1, accent, True, align=PP_ALIGN.CENTER, fit=True)
    else:
        pill(slide, x + 0.12, y + 0.14, 0.44, icon, accent, C["white"], h=0.27, size=7.0)
        text(slide, x + 0.66, y + 0.12, w - 0.80, 0.27, title, 9.2, accent, True, fit=True)
    text(slide, x + 0.15, y + 0.50, w - 0.30, h - 0.60, subtitle, 7.0, C["slate"], align=PP_ALIGN.CENTER, fit=True)


def _arrow_end(connector):
    line_xml = connector._element.spPr.get_or_add_ln()
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    line_xml.append(tail)


def edge(slide, x1, y1, x2, y2, color=None, width=1.6, elbow=False, label=None, label_x=None, label_y=None):
    connector_type = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    line = slide.shapes.add_connector(connector_type, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color or C["muted"])
    line.line.width = Pt(width)
    _arrow_end(line)
    if label:
        lx = label_x if label_x is not None else (x1 + x2) / 2 - 0.35
        ly = label_y if label_y is not None else (y1 + y2) / 2 - 0.12
        pill(slide, lx, ly, 0.76, label, C["white"], color or C["muted"], h=0.22, size=6.1, line=C["line"])
    return line


def right_arrow(slide, x, y, w=0.26, h=0.32, color=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color or C["muted"])
    shape.line.fill.background()
    return shape


def down_arrow(slide, x, y, w=0.28, h=0.30, color=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color or C["muted"])
    shape.line.fill.background()
    return shape


def add_system_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    header(
        slide,
        "Architecture Diagram 1/5 · End-to-End",
        "Knowledge Base 系統完整端到端架構圖",
        "從瀏覽器入口、FastAPI/Celery 編排、知識處理與 AI 推論，到 Qdrant、Neo4j 與檔案資產的完整連線。",
        1,
    )

    zone(slide, 0.40, 1.34, 1.35, 5.48, "USERS", C["blue"], C["blue_soft"])
    zone(slide, 1.91, 1.34, 2.10, 5.48, "EXPERIENCE", C["blue2"], "F1F5FF")
    zone(slide, 4.18, 1.34, 3.05, 5.48, "APPLICATION CORE", C["teal"], C["teal_soft"])
    zone(slide, 7.40, 1.34, 2.30, 5.48, "INTELLIGENCE", C["amber2"], C["amber_soft"])
    zone(slide, 9.87, 1.34, 3.00, 5.48, "DATA & AI SERVICES", C["green"], C["green_soft"])

    # Connections are placed before nodes so boxes remain visually dominant.
    edge(slide, 1.56, 2.39, 2.08, 2.39, C["blue"], label="HTTPS", label_x=1.52, label_y=2.08)
    edge(slide, 1.56, 4.68, 2.08, 4.68, C["blue"], label="HTTPS", label_x=1.52, label_y=4.36)
    edge(slide, 3.77, 2.23, 4.38, 2.23, C["blue2"], label="REST", label_x=3.72, label_y=1.93)
    edge(slide, 3.77, 3.22, 4.38, 3.22, C["blue2"], label="WSS", label_x=3.72, label_y=2.92)
    edge(slide, 3.77, 4.58, 4.38, 4.58, C["blue2"], label="REST", label_x=3.72, label_y=4.28)

    edge(slide, 5.40, 2.69, 5.40, 3.82, C["teal"], label="task", label_x=5.03, label_y=3.16)
    edge(slide, 5.96, 4.33, 5.96, 5.07, C["teal"], label="queue", label_x=5.58, label_y=4.57)
    edge(slide, 5.17, 5.55, 7.60, 2.35, C["teal"], elbow=True, label="search", label_x=6.64, label_y=2.74)
    edge(slide, 6.43, 5.55, 7.60, 4.10, C["teal"], elbow=True, label="ingest", label_x=6.82, label_y=4.63)
    edge(slide, 6.56, 3.22, 7.60, 5.48, C["teal"], elbow=True, label="proxy", label_x=6.89, label_y=5.02)

    edge(slide, 9.47, 2.35, 10.09, 2.10, C["amber2"], elbow=True, label="vector", label_x=9.47, label_y=1.72)
    edge(slide, 9.47, 2.55, 10.09, 3.28, C["amber2"], elbow=True, label="graph", label_x=9.47, label_y=2.88)
    edge(slide, 9.47, 2.80, 10.09, 6.08, C["amber2"], elbow=True, label="LLM", label_x=9.47, label_y=5.66)
    edge(slide, 9.47, 4.10, 10.09, 4.48, C["amber2"], elbow=True, label="write", label_x=9.47, label_y=4.12)
    edge(slide, 9.47, 4.25, 10.09, 3.52, C["amber2"], elbow=True)
    edge(slide, 9.47, 5.48, 10.09, 6.08, C["amber2"], elbow=True, label="agent", label_x=9.47, label_y=5.56)
    edge(slide, 11.34, 5.55, 8.70, 5.48, C["green"], elbow=True, label="identity", label_x=9.78, label_y=5.21)

    node(slide, 0.57, 1.91, 1.01, 0.94, "知識使用者", "搜尋、問答、引用追溯", C["blue"], "KU", title_size=8.5, subtitle_size=6.6)
    node(slide, 0.57, 4.19, 1.01, 0.98, "管理/維護者", "上傳、Chunk、排程、日誌", C["blue"], "AD", title_size=8.1, subtitle_size=6.5)

    node(slide, 2.08, 1.70, 1.69, 0.73, "Nginx TLS Gateway", "HTTPS 靜態頁 + API/WS proxy", C["blue2"], "NG", title_size=8.6, subtitle_size=6.5)
    node(slide, 2.08, 2.62, 1.69, 0.78, "Vue Search / Admin", "搜尋、上傳、管理、Chunk、Skills", C["blue2"], "UI", title_size=8.6, subtitle_size=6.4)
    node(slide, 2.08, 3.64, 1.69, 0.78, "Chat UI", "KB context、來源、WebSocket", C["blue2"], "CH", title_size=8.8, subtitle_size=6.6)
    node(slide, 2.08, 4.60, 1.69, 0.84, "Upload / Watch", "API 上傳與 watch folder", C["blue2"], "UP", title_size=8.6, subtitle_size=6.5)

    node(slide, 4.38, 1.73, 2.18, 0.96, "FastAPI / Uvicorn", "REST API、Pydantic、task status、admin API", C["teal"], "FA", title_size=9.4, subtitle_size=7.0)
    node(slide, 4.38, 2.89, 2.18, 0.82, "WebSocket Proxy", "chat queue、session lock、upstream events", C["teal"], "WS", title_size=9.1, subtitle_size=6.8)
    store(slide, 4.38, 3.91, 2.18, 0.92, "Redis", "Celery broker/result、cache、task/chat state", C["teal"], "RD")
    node(slide, 4.38, 5.07, 1.02, 0.90, "Search Worker", "queue=search", C["teal"], "CS", title_size=7.8, subtitle_size=6.4)
    node(slide, 5.54, 5.07, 1.02, 0.90, "Ingest Worker", "queue=ingest", C["teal"], "CI", title_size=7.8, subtitle_size=6.4)
    node(slide, 4.38, 6.14, 2.18, 0.50, "Celery Beat · watch folder scheduler", "", C["teal"], title_size=7.6)

    node(slide, 7.60, 1.78, 1.87, 1.03, "SearchEngine", "Basic / Deep / Vector / Hybrid / Report Graph", C["amber2"], "SE", title_size=9.2, subtitle_size=6.7)
    node(slide, 7.60, 3.48, 1.87, 1.02, "Document Pipeline", "MarkItDown、OCR、assets、chunk、metadata", C["amber2"], "DP", title_size=8.9, subtitle_size=6.7)
    node(slide, 7.60, 5.04, 1.87, 0.94, "OpenClaw Gateway", "agent events、tools、session", C["amber2"], "OC", title_size=8.8, subtitle_size=6.7)

    store(slide, 10.09, 1.67, 1.26, 0.92, "Qdrant", "768D chunks", C["green"], "QD")
    store(slide, 11.47, 1.67, 1.20, 0.92, "Neo4j", "graph data", C["green"], "N4")
    store(slide, 10.09, 2.88, 2.58, 0.92, "File Store", "raw / watch / uploads / processed / assets", C["green"], "FS")
    store(slide, 10.09, 4.16, 2.58, 0.92, "OpenClaw Runtime", "identity / workspace / skills / memory", C["green"], "RT")
    node(slide, 10.09, 5.68, 2.58, 0.84, "Ollama LLM Service", "qwen3-coder-next · host :11434", C["green"], "LLM", title_size=9.2, subtitle_size=6.8)

    pill(slide, 0.50, 6.86, 1.30, "SOLID ARROW", C["navy"], C["white"], h=0.24, size=6.5)
    text(slide, 1.93, 6.87, 4.10, 0.20, "表示主要 request / task / data flow", 7.0, C["slate"])
    footer(slide, "互動入口與背景任務分離；Search 與 Ingest 共用 Redis 編排，但分別進入檢索與文件處理管線。")


def add_query_chat_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    header(
        slide,
        "Architecture Diagram 2/5 · Runtime",
        "查詢與 OpenClaw 聊天架構圖",
        "同一個前端同時使用 REST 任務式搜尋與 WebSocket 串流聊天；兩條路徑在答案呈現時整合 KB 來源。",
        2,
    )

    zone(slide, 0.42, 1.40, 12.40, 2.38, "PATH A · KB SEARCH / REST + CELERY", C["blue"], C["blue_soft"])
    zone(slide, 0.42, 4.03, 12.40, 2.42, "PATH B · AI CHAT / WEBSOCKET", C["teal"], C["teal_soft"])

    top = [
        (0.68, "Browser", "提出問題", C["blue"], "01"),
        (2.17, "Nginx", "POST /search", C["blue"], "02"),
        (3.66, "FastAPI", "建立 task_id", C["blue"], "03"),
        (5.15, "Redis / Celery", "search queue", C["blue"], "04"),
        (6.82, "SearchEngine", "路由與雙引擎", C["blue"], "05"),
        (8.49, "Qdrant + Neo4j", "向量 + 圖譜", C["blue"], "06"),
        (10.16, "Ollama", "答案生成", C["blue"], "07"),
        (11.48, "Task Result", "poll + citations", C["blue"], "08"),
    ]
    for i, (x, title_value, sub, accent, icon) in enumerate(top):
        w = 1.26 if i in {0, 1, 2, 7} else 1.43
        node(slide, x, 2.01, w, 1.10, title_value, sub, accent, icon, title_size=8.4, subtitle_size=6.8)
        if i < len(top) - 1:
            next_x = top[i + 1][0]
            right_arrow(slide, x + w + 0.07, 2.42, max(0.18, next_x - (x + w) - 0.12), 0.30, C["blue"])
    text(slide, 0.72, 3.30, 11.70, 0.22, "回傳：answer · sources · citation_distribution · mode · queue_position", 7.8, C["blue"], True, align=PP_ALIGN.CENTER)

    bottom = [
        (0.68, "Chat UI", "載入 chat-config", C["teal"], "01"),
        (2.30, "KB Context", "依題型背景檢索", C["teal"], "02"),
        (3.92, "FastAPI /ws", "auth + WSS", C["teal"], "03"),
        (5.54, "Queue / Locks", "browser + global slot", C["teal"], "04"),
        (7.16, "OpenClaw GW", "agent / tools", C["teal"], "05"),
        (8.78, "Ollama", "LLM inference", C["teal"], "06"),
        (10.40, "Agent Events", "stream delta / lifecycle", C["teal"], "07"),
        (11.69, "Rendered Reply", "來源 + fallback", C["teal"], "08"),
    ]
    for i, (x, title_value, sub, accent, icon) in enumerate(bottom):
        w = 1.26 if i in {7} else 1.35
        node(slide, x, 4.66, w, 1.10, title_value, sub, accent, icon, title_size=8.2, subtitle_size=6.6)
        if i < len(bottom) - 1:
            next_x = bottom[i + 1][0]
            right_arrow(slide, x + w + 0.06, 5.07, max(0.18, next_x - (x + w) - 0.11), 0.30, C["teal"])
    edge(slide, 2.98, 4.66, 6.98, 3.11, C["teal2"], elbow=True, label="KB result", label_x=4.50, label_y=3.52)
    text(slide, 0.72, 5.99, 11.70, 0.23, "可靠度：idempotency key · WebSocket reconnect · active-slot TTL · queued request cleanup · deterministic KB excerpt fallback", 7.6, C["teal2"], True, align=PP_ALIGN.CENTER, fit=True)
    footer(slide, "REST 路徑確保完整知識庫結果；WebSocket 路徑優先串流體感，並用 KB context 與引用補強可信度。")


def add_ingestion_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    header(
        slide,
        "Architecture Diagram 3/5 · Ingestion",
        "文件攝入與知識建立架構圖",
        "輸入文件先進入非同步任務，再經轉檔、資產保存、分塊與模式路由，最後寫入 Qdrant、Neo4j 與可回溯檔案。",
        3,
    )

    zone(slide, 0.42, 1.38, 2.02, 5.30, "INPUT", C["blue"], C["blue_soft"])
    zone(slide, 2.62, 1.38, 2.06, 5.30, "ORCHESTRATION", C["teal"], C["teal_soft"])
    zone(slide, 4.86, 1.38, 2.30, 5.30, "NORMALIZATION", C["blue2"], "F1F5FF")
    zone(slide, 7.34, 1.38, 2.28, 5.30, "MODE ROUTER", C["amber2"], C["amber_soft"])
    zone(slide, 9.80, 1.38, 3.02, 5.30, "PERSISTENCE", C["green"], C["green_soft"])

    node(slide, 0.64, 1.90, 1.58, 0.87, "API Upload", "/api/upload + task_id", C["blue"], "UP")
    node(slide, 0.64, 3.04, 1.58, 0.87, "Watch Folder", "Celery Beat 定時掃描", C["blue"], "WF")
    node(slide, 0.64, 4.18, 1.58, 0.87, "Raw / Manual", "批次或維運匯入", C["blue"], "RM")

    edge(slide, 2.23, 2.34, 2.84, 2.34, C["blue"], label="file", label_x=2.20, label_y=2.02)
    edge(slide, 2.23, 3.48, 2.84, 3.48, C["blue"])
    edge(slide, 2.23, 4.62, 2.84, 4.62, C["blue"])

    store(slide, 2.84, 1.95, 1.62, 0.92, "Redis", "task state / file hash / queue", C["teal"], "RD")
    node(slide, 2.84, 3.31, 1.62, 1.00, "Celery Ingest", "processing lock、狀態、失敗可見", C["teal"], "CI")
    node(slide, 2.84, 4.75, 1.62, 0.76, "Celery Beat", "watch scheduler", C["teal"], "CB")
    edge(slide, 3.65, 2.87, 3.65, 3.31, C["teal"], label="task", label_x=3.28, label_y=2.91)
    edge(slide, 3.65, 4.75, 3.65, 4.31, C["teal"], label="trigger", label_x=3.26, label_y=4.43)

    edge(slide, 4.46, 3.80, 5.08, 3.80, C["teal"], label="process", label_x=4.40, label_y=3.49)
    node(slide, 5.08, 1.86, 1.84, 0.92, "FileConverter", "MarkItDown + PDF/Excel 增強", C["blue2"], "CV")
    node(slide, 5.08, 3.05, 1.84, 0.92, "Assets / OCR", "頁面快照、圖片、表格文字", C["blue2"], "AS")
    node(slide, 5.08, 4.24, 1.84, 0.92, "Chunker", "標題感知 500 chars + metadata", C["blue2"], "CK")
    edge(slide, 6.00, 2.78, 6.00, 3.05, C["blue2"])
    edge(slide, 6.00, 3.97, 6.00, 4.24, C["blue2"])

    edge(slide, 6.92, 4.70, 7.56, 4.70, C["blue2"], label="chunks", label_x=6.91, label_y=4.38)
    node(slide, 7.56, 1.86, 1.82, 0.88, "REPORT", "Report graph + vector", C["amber2"], "R")
    node(slide, 7.56, 3.01, 1.82, 0.88, "VECTOR-ONLY", "Lab / Project / Automation", C["amber2"], "V")
    node(slide, 7.56, 4.16, 1.82, 1.02, "SEMANTIC", "4G/5G / WiFi + Ollama entity extraction", C["amber2"], "S")
    node(slide, 7.56, 5.46, 1.82, 0.76, "Routing Rules", "filename + extraction mode", C["amber2"], "RT", title_size=8.2)

    edge(slide, 9.38, 2.30, 10.02, 2.18, C["amber2"], elbow=True)
    edge(slide, 9.38, 2.45, 11.42, 3.34, C["amber2"], elbow=True)
    edge(slide, 9.38, 3.45, 10.02, 2.18, C["amber2"], elbow=True)
    edge(slide, 9.38, 4.60, 10.02, 2.18, C["amber2"], elbow=True)
    edge(slide, 9.38, 4.76, 11.42, 3.34, C["amber2"], elbow=True)
    edge(slide, 9.38, 4.92, 10.02, 5.42, C["amber2"], elbow=True)

    store(slide, 10.02, 1.75, 1.20, 0.92, "Qdrant", "chunks + vectors", C["green"], "QD")
    store(slide, 11.42, 2.91, 1.18, 0.92, "Neo4j", "documents + graph", C["green"], "N4")
    store(slide, 10.02, 4.99, 2.58, 1.02, "Processed File Store", "Markdown / source.json / assets / originals", C["green"], "FS")
    node(slide, 10.02, 6.16, 2.58, 0.38, "Index refresh · task completed", "", C["green"], title_size=7.4)

    pill(slide, 0.62, 6.37, 0.94, "CONTROLS", C["navy"], C["white"], h=0.24, size=6.5)
    text(slide, 1.70, 6.37, 5.14, 0.20, "SHA-256 去重 · cleanup-before-write · processing locks · task TTL", 7.0, C["slate"], fit=True)
    footer(slide, "文件型態決定寫入策略；所有路徑都保留 processed 文件與來源 metadata，確保可追溯與可重攝入。")


def add_retrieval_data_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    header(
        slide,
        "Architecture Diagram 4/5 · Retrieval & Data",
        "混合檢索、資料融合與引用架構圖",
        "Query Router 依題型選擇向量、圖譜與檔案內容，經業務重排與引用封裝後，再交由 deterministic 或 AI 回答路徑。",
        4,
    )

    zone(slide, 0.40, 1.38, 2.05, 5.35, "QUERY ROUTING", C["blue"], C["blue_soft"])
    zone(slide, 2.63, 1.38, 3.00, 5.35, "KNOWLEDGE SOURCES", C["green"], C["green_soft"])
    zone(slide, 5.81, 1.38, 2.60, 5.35, "FUSION", C["teal"], C["teal_soft"])
    zone(slide, 8.59, 1.38, 2.00, 5.35, "ANSWER ENGINE", C["amber2"], C["amber_soft"])
    zone(slide, 10.77, 1.38, 2.08, 5.35, "DELIVERY", C["blue2"], "F1F5FF")

    node(slide, 0.65, 1.88, 1.55, 0.78, "Query Analyzer", "題型、關鍵字、專案/case hints", C["blue"], "QA")
    modes = ["AUTO", "BASIC", "DEEP", "VECTOR", "HYBRID", "HYBRID+", "REPORT GRAPH"]
    y = 2.91
    for idx, mode in enumerate(modes):
        fill = C["blue"] if idx in {0, 4, 5} else C["white"]
        color = C["white"] if idx in {0, 4, 5} else C["blue"]
        pill(slide, 0.70, y, 1.45, mode, fill, color, h=0.31, size=6.8, line=C["blue"])
        y += 0.45

    edge(slide, 2.20, 2.28, 2.87, 2.12, C["blue"], elbow=True, label="semantic", label_x=2.18, label_y=1.78)
    edge(slide, 2.20, 2.42, 2.87, 3.41, C["blue"], elbow=True, label="graph", label_x=2.18, label_y=3.05)
    edge(slide, 2.20, 2.56, 2.87, 4.75, C["blue"], elbow=True, label="content", label_x=2.18, label_y=4.38)

    store(slide, 2.87, 1.70, 2.50, 1.08, "Qdrant Vector Store", "collection knowledge_base · BGE 768D · Cosine · metadata", C["green"], "QD")
    store(slide, 2.87, 3.01, 2.50, 1.08, "Neo4j Knowledge Graph", "Document / Entity / Report / Section / TestItem / Case / Metric", C["green"], "N4")
    store(slide, 2.87, 4.36, 2.50, 1.08, "Processed Content", "Markdown、source metadata、image refs、deterministic extracts", C["green"], "FS")

    edge(slide, 5.37, 2.24, 6.05, 2.24, C["green"])
    edge(slide, 5.37, 3.55, 6.05, 3.55, C["green"])
    edge(slide, 5.37, 4.90, 6.05, 4.90, C["green"])
    fusion_steps = [
        ("Parallel Recall", "向量 raw + 圖譜 raw + processed"),
        ("Business Rerank", "report/case/numeric/project preference"),
        ("Dedup & Classify", "doc/chunk/source type/storage category"),
        ("Citation Package", "sources + citation_distribution + excerpts"),
    ]
    y = 1.83
    for idx, (title_value, sub) in enumerate(fusion_steps, start=1):
        node(slide, 6.05, y, 2.12, 0.78, title_value, sub, C["teal"], str(idx), title_size=8.4, subtitle_size=6.3)
        if idx < len(fusion_steps):
            down_arrow(slide, 6.97, y + 0.80, 0.28, 0.25, C["teal"])
        y += 1.02

    edge(slide, 8.17, 3.55, 8.82, 3.55, C["teal"], label="context", label_x=8.12, label_y=3.23)
    node(slide, 8.82, 1.82, 1.54, 0.98, "Deterministic", "報告/WiFi/比較題直接輸出", C["amber2"], "D")
    node(slide, 8.82, 3.18, 1.54, 0.98, "Ollama Synthesis", "根據 context 生成答案", C["amber2"], "LLM")
    node(slide, 8.82, 4.54, 1.54, 0.98, "OpenClaw", "agent / tools / streaming", C["amber2"], "OC")
    edge(slide, 10.36, 2.31, 11.00, 2.31, C["amber2"])
    edge(slide, 10.36, 3.67, 11.00, 3.67, C["amber2"])
    edge(slide, 10.36, 5.03, 11.00, 5.03, C["amber2"])

    node(slide, 11.00, 1.82, 1.60, 0.98, "Search Result", "answer + sources + mode", C["blue2"], "SR")
    node(slide, 11.00, 3.18, 1.60, 0.98, "Citation UI", "文件、片段、分類占比", C["blue2"], "CT")
    node(slide, 11.00, 4.54, 1.60, 0.98, "Chat Reply", "stream + source fallback", C["blue2"], "CH")

    rect(slide, 2.87, 5.84, 7.49, 0.54, C["red_soft"], line="F1BABA", radius=True)
    pill(slide, 3.05, 5.98, 0.88, "GUARDRAIL", C["red"], C["white"], h=0.25, size=6.4)
    text(slide, 4.10, 5.93, 6.00, 0.31, "有來源時不得誤稱查無資料；片段不足要揭露限制並保留來源。", 7.4, C["ink"], True, fit=True)
    footer(slide, "平台的核心不是單一 LLM，而是 Query Router、雙資料庫召回、業務重排與 citation package 的整體設計。")


def add_deployment_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background(slide)
    header(
        slide,
        "Architecture Diagram 5/5 · Deployment",
        "現行站台與 on-prem release 部署架構圖",
        "相同應用程式在兩種拓撲中有不同的服務邊界：現行站台使用 host Qdrant；on-prem release 將 Qdrant 納入 Compose。",
        5,
    )

    # Current site panel
    zone(slide, 0.40, 1.35, 6.10, 5.48, "CURRENT SITE · 61.216.9.52", C["blue"], C["blue_soft"])
    node(slide, 0.64, 1.87, 1.16, 0.68, "Browser", "HTTPS :3030", C["blue"], "U", title_size=8.5, subtitle_size=6.4)
    node(slide, 2.12, 1.87, 1.22, 0.68, "Nginx", "TLS + static + proxy", C["blue"], "NG", title_size=8.5, subtitle_size=6.2)
    edge(slide, 1.80, 2.21, 2.12, 2.21, C["blue"])

    zone(slide, 0.66, 2.86, 5.56, 2.58, "DOCKER COMPOSE", C["blue"], C["white"])
    node(slide, 0.88, 3.35, 1.24, 0.78, "FastAPI Web", "Uvicorn :8000 ×4", C["blue"], "FA", title_size=8.2, subtitle_size=6.3)
    store(slide, 2.34, 3.35, 1.20, 0.78, "Redis", "broker/cache", C["blue"], "RD")
    store(slide, 3.76, 3.35, 1.20, 0.78, "Neo4j", "graph DB", C["blue"], "N4")
    node(slide, 5.16, 3.35, 0.82, 0.78, "Beat", "watch", C["blue"], "CB", title_size=7.2)
    node(slide, 0.88, 4.40, 1.56, 0.72, "Search Worker", "queue=search", C["blue"], "CS", title_size=7.8)
    node(slide, 2.68, 4.40, 1.56, 0.72, "Ingest Worker", "queue=ingest", C["blue"], "CI", title_size=7.8)
    node(slide, 4.48, 4.40, 1.50, 0.72, "Frontend Files", "Vue + chat.html", C["blue"], "UI", title_size=7.8)
    edge(slide, 2.73, 2.55, 1.50, 3.35, C["blue"], elbow=True, label="proxy", label_x=1.76, label_y=2.74)
    edge(slide, 2.73, 2.55, 5.20, 4.40, C["blue"], elbow=True, label="static", label_x=4.30, label_y=2.74)

    node(slide, 0.88, 5.77, 1.56, 0.70, "Host Qdrant", ":6335", C["green"], "QD", title_size=8.2)
    node(slide, 2.68, 5.77, 1.56, 0.70, "Host Ollama", ":11434", C["green"], "LLM", title_size=8.2)
    node(slide, 4.48, 5.77, 1.50, 0.70, "OpenClaw GW", "WebSocket", C["green"], "OC", title_size=7.8)
    edge(slide, 1.50, 4.13, 1.50, 5.77, C["green"], label="vector", label_x=1.12, label_y=5.20)
    edge(slide, 1.50, 4.13, 3.46, 5.77, C["green"], elbow=True, label="LLM", label_x=2.54, label_y=5.20)
    edge(slide, 1.50, 4.13, 5.23, 5.77, C["green"], elbow=True, label="WS", label_x=4.36, label_y=5.20)

    # Release panel
    zone(slide, 6.72, 1.35, 6.12, 5.48, "ON-PREM RELEASE · CONFIGURABLE", C["teal"], C["teal_soft"])
    node(slide, 6.96, 1.87, 1.16, 0.68, "Browser", "HTTPS port", C["teal"], "U", title_size=8.5, subtitle_size=6.4)
    node(slide, 8.44, 1.87, 1.22, 0.68, "Nginx", "TLS + static + proxy", C["teal"], "NG", title_size=8.5, subtitle_size=6.2)
    edge(slide, 8.12, 2.21, 8.44, 2.21, C["teal"])

    zone(slide, 6.98, 2.86, 5.56, 2.58, "RELEASE COMPOSE", C["teal"], C["white"])
    node(slide, 7.20, 3.35, 1.16, 0.78, "FastAPI", "release workers", C["teal"], "FA", title_size=8.0)
    store(slide, 8.56, 3.35, 1.08, 0.78, "Redis", "volume", C["teal"], "RD")
    store(slide, 9.84, 3.35, 1.08, 0.78, "Neo4j", "volumes", C["teal"], "N4")
    store(slide, 11.12, 3.35, 1.16, 0.78, "Qdrant", "volume", C["teal"], "QD")
    node(slide, 7.20, 4.40, 1.48, 0.72, "Search Worker", "configurable", C["teal"], "CS", title_size=7.7)
    node(slide, 8.90, 4.40, 1.48, 0.72, "Ingest Worker", "isolated queue", C["teal"], "CI", title_size=7.7)
    node(slide, 10.60, 4.40, 1.68, 0.72, "Beat + Frontend", "scheduler + static", C["teal"], "RT", title_size=7.5)
    edge(slide, 9.05, 2.55, 7.78, 3.35, C["teal"], elbow=True, label="proxy", label_x=8.05, label_y=2.74)
    edge(slide, 9.05, 2.55, 11.44, 4.40, C["teal"], elbow=True, label="static", label_x=10.52, label_y=2.74)

    node(slide, 7.20, 5.77, 1.54, 0.70, "Host Ollama", ":11434", C["green"], "LLM", title_size=8.1)
    node(slide, 8.98, 5.77, 1.54, 0.70, "Host OpenClaw", "gateway", C["green"], "OC", title_size=8.0)
    node(slide, 10.76, 5.77, 1.52, 0.70, "Installer", "preflight / upgrade", C["green"], "IN", title_size=8.1)
    edge(slide, 7.78, 4.13, 7.97, 5.77, C["green"], elbow=True, label="LLM", label_x=7.45, label_y=5.20)
    edge(slide, 7.78, 4.13, 9.75, 5.77, C["green"], elbow=True, label="WS", label_x=8.86, label_y=5.20)

    pill(slide, 6.98, 6.52, 1.14, "PERSISTENCE", C["navy"], C["white"], h=0.24, size=6.4)
    text(slide, 8.24, 6.52, 4.00, 0.20, "app/data · runtime/openclaw · named volumes · release metadata", 6.9, C["slate"], fit=True)
    footer(slide, "現行站台與 release 的主要差異在 Qdrant 邊界與持久卷封裝；Ollama/OpenClaw 仍需視 host runtime 配置。")


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "Knowledge Base Architecture Diagrams"
    prs.core_properties.subject = "End-to-end, runtime, ingestion, retrieval and deployment diagrams"
    prs.core_properties.author = "Knowledge Base Architecture Team"
    prs.core_properties.comments = "Diagram-first revision generated from the current repository on 2026-07-20."

    add_system_overview(prs)
    add_query_chat_architecture(prs)
    add_ingestion_architecture(prs)
    add_retrieval_data_architecture(prs)
    add_deployment_architecture(prs)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(f"saved {build()}")
