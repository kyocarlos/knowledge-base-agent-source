from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("/home/da40_ai_gb10/knowledge-base")
NEO4J_OUT = OUT_DIR / "neo4j_customer_intro.pptx"
QDRANT_OUT = OUT_DIR / "qdrant_customer_intro.pptx"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0.0, radius=False):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def set_text(frame, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Sans CJK TC"):
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)


def textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, text, size, color, bold=bold, align=align)
    return box


def badge(slide, left, top, width, text, fill, color):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        Inches(0.30),
        fill=fill,
        line=fill,
        radius=True,
    )
    set_text(shape.text_frame, text, 8.5, color, align=PP_ALIGN.CENTER)
    return shape


def card(slide, left, top, width, height, title, body, accent, fill="FFFFFF"):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=fill,
        line="CBD5E1",
        radius=True,
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(7)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Noto Sans CJK TC"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = rgb(accent)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Noto Sans CJK TC"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = rgb("475569")
    return shape


def add_background(slide, accent_a="DBEAFE", accent_b="D1FAE5", accent_c="FEF3C7"):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    for left, top, width, height, color in [
        (Inches(-0.4), Inches(-0.3), Inches(3.8), Inches(3.8), accent_a),
        (Inches(9.3), Inches(-0.3), Inches(4.3), Inches(3.8), accent_b),
        (Inches(8.7), Inches(4.8), Inches(4.4), Inches(2.7), accent_c),
    ]:
        deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        deco.fill.solid()
        deco.fill.fore_color.rgb = rgb(color)
        deco.fill.transparency = 0.80
        deco.line.fill.background()

    border = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.22),
        Inches(0.18),
        Inches(12.89),
        Inches(7.14),
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    border.fill.transparency = 1


def add_header(slide, theme_name, subtitle, title_color, badge_fill, badge_color):
    textbox(slide, Inches(0.52), Inches(0.32), Inches(2.8), Inches(0.20), "Knowledge Base / Backend Database", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.56), Inches(4.5), Inches(0.42), theme_name, 24, title_color, bold=True)
    textbox(slide, Inches(0.52), Inches(0.98), Inches(8.5), Inches(0.32), subtitle, 12, "475569")
    badge(slide, Inches(10.78), Inches(0.50), Inches(1.90), "客戶說明用", badge_fill, badge_color)


def add_neo4j_slide(slide):
    add_background(slide, accent_a="DBEAFE", accent_b="D1FAE5", accent_c="E0E7FF")
    add_header(
        slide,
        "Neo4j",
        "用來保存資料之間的關係，讓系統知道文件、專案與章節彼此如何連結。",
        "0F172A",
        "EFF6FF",
        "1D4ED8",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(3.30), Inches(1.35), "它負責什麼", "Neo4j 主要保存「關係」，例如哪份文件屬於哪個專案、哪個章節和哪個測試項目相關。", "2563EB")
    card(slide, Inches(0.72), Inches(3.25), Inches(3.30), Inches(1.35), "系統怎麼用", "當使用者查詢跨文件、跨專案或需要上下文脈絡的內容時，Neo4j 可以提供關聯線索。", "16A34A")
    card(slide, Inches(0.72), Inches(4.72), Inches(3.30), Inches(1.35), "客戶怎麼理解", "它像一張關係地圖，重點不是存全文，而是告訴系統「誰和誰有關」。", "D97706")

    # simple node diagram
    center = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(5.10), Inches(2.05), Inches(1.30), Inches(1.30), fill="DBEAFE", line="BFDBFE")
    set_text(center.text_frame, "Neo4j", 16, "1D4ED8", bold=True, align=PP_ALIGN.CENTER)
    node1 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.00), Inches(1.15), Inches(1.05), Inches(1.05), fill="FFFFFF", line="CBD5E1")
    set_text(node1.text_frame, "文件", 13, "0F172A", bold=True, align=PP_ALIGN.CENTER)
    node2 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(6.45), Inches(1.10), Inches(1.05), Inches(1.05), fill="FFFFFF", line="CBD5E1")
    set_text(node2.text_frame, "專案", 13, "0F172A", bold=True, align=PP_ALIGN.CENTER)
    node3 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.05), Inches(4.05), Inches(1.05), Inches(1.05), fill="FFFFFF", line="CBD5E1")
    set_text(node3.text_frame, "章節", 13, "0F172A", bold=True, align=PP_ALIGN.CENTER)
    node4 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(6.52), Inches(4.10), Inches(1.05), Inches(1.05), fill="FFFFFF", line="CBD5E1")
    set_text(node4.text_frame, "關聯", 13, "0F172A", bold=True, align=PP_ALIGN.CENTER)

    # connecting lines
    for x1, y1, x2, y2 in [
        (Inches(5.05), Inches(1.68), Inches(5.60), Inches(2.05)),
        (Inches(6.60), Inches(1.65), Inches(6.10), Inches(2.05)),
        (Inches(5.05), Inches(4.58), Inches(5.55), Inches(3.35)),
        (Inches(6.55), Inches(4.58), Inches(6.10), Inches(3.35)),
    ]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = rgb("94A3B8")
        line.line.width = Pt(1.5)

    card(slide, Inches(8.25), Inches(1.85), Inches(4.00), Inches(1.20), "一句話說明", "Neo4j 讓系統能夠把不同文件與章節串起來，回答「這些資料彼此有什麼關係」。", "1D4ED8", fill="FFFFFF")
    card(slide, Inches(8.25), Inches(3.25), Inches(4.00), Inches(1.20), "最適合的查詢", "跨文件比較、專案關聯、章節脈絡、測試項目對照。", "16A34A", fill="FFFFFF")
    card(slide, Inches(8.25), Inches(4.65), Inches(4.00), Inches(1.20), "對客戶的說法", "Neo4j 是關係型資料庫，用來保存知識圖譜的結構與上下文。", "D97706", fill="FFFFFF")

    section = slide.shapes.add_textbox(Inches(0.74), Inches(6.20), Inches(2.5), Inches(0.20))
    set_text(section.text_frame, "實際範例", 11, "0F172A", bold=True)
    ex1 = card(
        slide,
        Inches(0.72),
        Inches(6.42),
        Inches(3.55),
        Inches(0.60),
        "SCU2060 ↔ SCU2140",
        "兩份報告都含 Throughput / Latency 測試，可用 Neo4j 串出共同主題與差異。",
        "2563EB",
    )
    ex2 = card(
        slide,
        Inches(4.47),
        Inches(6.42),
        Inches(3.55),
        Inches(0.60),
        "SCU2050 ↔ SCU2060",
        "兩份報告都可回到 Handover / Performance 相關脈絡，方便交叉比對。",
        "16A34A",
    )
    ex3 = card(
        slide,
        Inches(8.22),
        Inches(6.42),
        Inches(4.00),
        Inches(0.60),
        "Neo4j 看到的是關係",
        "不是只看單一文件，而是把文件彼此之間的連結整理出來。",
        "D97706",
    )

    footer = slide.shapes.add_textbox(Inches(0.74), Inches(7.10), Inches(11.9), Inches(0.20))
    set_text(footer.text_frame, "重點：Neo4j 管關係，讓系統理解文件與文件之間的連結。", 9.5, "64748B", align=PP_ALIGN.RIGHT)


def add_qdrant_slide(slide):
    add_background(slide, accent_a="DCFCE7", accent_b="DBEAFE", accent_c="FEF3C7")
    add_header(
        slide,
        "Qdrant",
        "用來保存內容向量，讓系統能用語意方式快速找到最相關的段落。",
        "0F172A",
        "ECFDF5",
        "047857",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(3.30), Inches(1.35), "它負責什麼", "Qdrant 主要保存文件內容的向量，讓系統可以依照意思去找相近內容，而不是只看關鍵字。", "16A34A")
    card(slide, Inches(0.72), Inches(3.25), Inches(3.30), Inches(1.35), "系統怎麼用", "當使用者問法比較口語、簡略，或需要找相似段落時，Qdrant 可以快速召回最接近的內容。", "2563EB")
    card(slide, Inches(0.72), Inches(4.72), Inches(3.30), Inches(1.35), "客戶怎麼理解", "它像一個智慧索引，幫系統找到「意思最接近」的資料片段。", "D97706")

    center = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.02), Inches(2.05), Inches(1.40), Inches(1.25), fill="DCFCE7", line="BBF7D0", radius=True)
    set_text(center.text_frame, "Qdrant", 16, "047857", bold=True, align=PP_ALIGN.CENTER)
    ring1 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.00), Inches(1.25), Inches(3.44), Inches(3.44), fill="FFFFFF", line="C8E6D7")
    ring1.fill.transparency = 1
    ring1.fill.transparency = 1
    ring2 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.43), Inches(1.68), Inches(2.58), Inches(2.58), fill="FFFFFF", line="C8E6D7")
    ring2.fill.transparency = 1
    ring2.fill.transparency = 1
    # Bring center to top by re-adding? simpler: add after rings by order? keep fine.
    labels = [
        (Inches(4.00), Inches(1.05), "相似內容"),
        (Inches(6.15), Inches(1.80), "向量檢索"),
        (Inches(4.18), Inches(4.08), "語意召回"),
        (Inches(6.10), Inches(4.15), "內容索引"),
    ]
    for left, top, text in labels:
        badge(slide, left, top, Inches(1.40), text, "ECFDF5", "047857")

    card(slide, Inches(8.25), Inches(1.85), Inches(4.00), Inches(1.20), "一句話說明", "Qdrant 讓系統能理解「這段內容跟我問的意思像不像」。", "047857", fill="FFFFFF")
    card(slide, Inches(8.25), Inches(3.25), Inches(4.00), Inches(1.20), "最適合的查詢", "語意搜尋、相似段落召回、口語化提問、找相關內容。", "2563EB", fill="FFFFFF")
    card(slide, Inches(8.25), Inches(4.65), Inches(4.00), Inches(1.20), "對客戶的說法", "Qdrant 是向量資料庫，用來做智慧搜尋與內容比對。", "D97706", fill="FFFFFF")

    section = slide.shapes.add_textbox(Inches(0.74), Inches(6.20), Inches(2.5), Inches(0.20))
    set_text(section.text_frame, "實際範例", 11, "0F172A", bold=True)
    ex1 = card(
        slide,
        Inches(0.72),
        Inches(6.42),
        Inches(3.72),
        Inches(0.60),
        "範例報告切 chunk",
        "SIT-TR-SC-NR-Throughput-SCU2060-n79-EV-V13.8.xlsx 可拆成多個內容區塊。",
        "047857",
    )
    ex2 = card(
        slide,
        Inches(4.65),
        Inches(6.42),
        Inches(3.72),
        Inches(0.60),
        "Chunk 1 / 2 / 3",
        "2. Introduction、3. Test Result Summary、4. Performance Test 各自成為獨立向量。",
        "2563EB",
    )
    ex3 = card(
        slide,
        Inches(8.58),
        Inches(6.42),
        Inches(3.64),
        Inches(0.60),
        "Qdrant 看到的是內容",
        "查詢時先比對哪個 chunk 最像，再把結果交給後端整理。",
        "D97706",
    )

    footer = slide.shapes.add_textbox(Inches(0.74), Inches(7.10), Inches(11.9), Inches(0.20))
    set_text(footer.text_frame, "重點：Qdrant 管內容相似度，讓系統更快找到最像的資料段落。", 9.5, "64748B", align=PP_ALIGN.RIGHT)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

add_neo4j_slide(prs.slides.add_slide(prs.slide_layouts[6]))
prs.save(NEO4J_OUT)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

add_qdrant_slide(prs.slides.add_slide(prs.slide_layouts[6]))
prs.save(QDRANT_OUT)

print(f"saved {NEO4J_OUT}")
print(f"saved {QDRANT_OUT}")
