from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "ANRITSU_AGENT_A2A_IMPLEMENTATION_GUIDE.pptx"
FONT = "Noto Sans CJK TC"
W, H = 13.333, 7.5
C = {
    "navy": "0B1F33", "ink": "193247", "muted": "607589", "line": "C8D5DF",
    "blue": "2F6FED", "blue_fill": "EAF1FF", "teal": "008F89", "teal_fill": "E5F7F5",
    "orange": "C87512", "orange_fill": "FFF1D7", "red": "B33A3A", "red_fill": "FCECEC",
    "green": "2E8555", "green_fill": "E8F6ED", "purple": "7653A6", "purple_fill": "F1ECFA",
    "gray_fill": "F4F7FA", "white": "FFFFFF",
}


def rgb(value):
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill, line=None, rounded=True, width=1.2):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()
    return shape


def text(slide, x, y, w, h, value, size=12, color=None, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = valign
    frame.margin_left = frame.margin_right = Pt(7)
    frame.margin_top = frame.margin_bottom = Pt(3)
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color or C["ink"])
    return box


def line(slide, x1, y1, x2, y2, color=None, width=1.8, arrow=False):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(color or C["line"])
    connector.line.width = Pt(width)
    if arrow:
        from pptx.oxml.xmlchemy import OxmlElement
        end = connector._element.spPr.get_or_add_ln()
        tail = OxmlElement("a:tailEnd")
        tail.set("type", "triangle")
        tail.set("w", "sm")
        tail.set("len", "sm")
        end.append(tail)
    return connector


def title(slide, heading, subtitle=None, number=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C["white"])
    rect(slide, 0, 0, W, 0.16, C["blue"], None, False)
    text(slide, 0.55, 0.34, 10.9, 0.48, heading, 24, C["navy"], True)
    if subtitle:
        text(slide, 0.58, 0.86, 11.5, 0.28, subtitle, 10, C["muted"])
    if number is not None:
        text(slide, 12.0, 0.38, 0.75, 0.30, f"{number:02d}", 10, C["blue"], True, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, heading, body, fill="gray_fill", accent="blue", heading_size=14, body_size=10):
    rect(slide, x, y, w, h, C[fill], C[accent], True, 1.4)
    text(slide, x + 0.10, y + 0.10, w - 0.20, 0.34, heading, heading_size, C[accent], True)
    text(slide, x + 0.10, y + 0.50, w - 0.20, h - 0.58, body, body_size, C["ink"])


def pill(slide, x, y, w, label, fill, color):
    rect(slide, x, y, w, 0.32, fill, None, True)
    text(slide, x, y + 0.01, w, 0.25, label, 9, color, True, PP_ALIGN.CENTER)


def footer(slide, value="KB / Anritsu A2A Implementation Guide"):
    text(slide, 0.58, 7.16, 8.5, 0.18, value, 7, C["muted"])
    text(slide, 10.4, 7.16, 2.35, 0.18, "Implementation handoff", 7, C["muted"], False, PP_ALIGN.RIGHT)


def new_slide(prs, heading, subtitle, number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, heading, subtitle, number)
    footer(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C["navy"])
    rect(slide, 0, 0, W, 0.20, C["blue"], None, False)
    text(slide, 0.78, 1.15, 11.4, 0.90, "Anritsu Agent × KB Agent", 34, C["white"], True)
    text(slide, 0.82, 2.10, 10.8, 0.56, "A2A 功能實作與整合指南", 23, "BFD7FF", True)
    text(slide, 0.84, 3.05, 10.2, 0.80,
         "讓 KB Agent 可以安全委派 Anritsu Agent 執行測試，\n"
         "並沿用既有 Excel 上傳與 Neo4j / Qdrant 攝入流程。", 15, "D7E7F6")
    rect(slide, 0.84, 5.15, 11.4, 0.92, "173A57", "2F6FED", True, 1.2)
    text(slide, 1.08, 5.33, 10.9, 0.48,
         "核心原則：A2A 只負責 Agent-to-Agent 委派；儀器控制留在 Anritsu；KB 核心服務不重寫。",
         13, C["white"], True, PP_ALIGN.CENTER)
    text(slide, 0.84, 6.72, 6.0, 0.22, "交付對象：Anritsu Agent 開發與整合人員", 9, "BFD7FF")
    text(slide, 10.2, 6.72, 2.05, 0.22, "2026-08", 9, "BFD7FF", False, PP_ALIGN.RIGHT)

    # 2. Executive decisions
    slide = new_slide(prs, "先看結論", "這份指南的決策與不可違反的邊界", 2)
    card(slide, 0.72, 1.45, 3.75, 1.75, "角色", "KB Agent = A2A Client\nAnritsu Agent = A2A Server\nAnritsu 端保留儀器控制權", "blue_fill", "blue", 15, 12)
    card(slide, 4.80, 1.45, 3.75, 1.75, "資料流程", "A2A 傳測試工作與狀態\nExcel 仍走既有 /api/upload/ingest\n不是新增第二套 uploader", "teal_fill", "teal", 15, 12)
    card(slide, 8.88, 1.45, 3.75, 1.75, "隔離", "不讓 Anritsu 連 Neo4j、Qdrant、Redis\n不修改 KB 搜尋、Celery 與資料庫\n只新增 A2A adapter", "orange_fill", "orange", 15, 12)
    text(slide, 0.82, 3.75, 11.8, 0.35, "Anritsu Agent 必須配合的範圍", 16, C["navy"], True)
    items = [
        ("新增", "A2A Server、Agent Card、Task executor、權限與狀態管理"),
        ("沿用", "既有儀器控制、MCP tools、Excel 產生器與 KB 上傳程式"),
        ("禁止", "LLM 直接生成任意 shell / 儀器命令；A2A 不可取得資料庫權限"),
    ]
    for i, (head, body) in enumerate(items):
        y = 4.35 + i * 0.72
        pill(slide, 0.88, y + 0.05, 0.95, head, C["blue_fill"] if i == 0 else C["orange_fill"] if i == 2 else C["green_fill"], C["blue"] if i == 0 else C["orange"] if i == 2 else C["green"])
        text(slide, 2.02, y, 10.4, 0.42, body, 12, C["ink"], True)

    # 3. Protocol roles
    slide = new_slide(prs, "A2A、MCP 與既有 HTTP 的分工", "不要用一種協定取代另一種協定", 3)
    card(slide, 0.72, 1.45, 3.75, 3.75, "A2A：Agent ↔ Agent", "用途\n• KB Agent 委派測試\n• 傳送任務與 context\n• 追蹤長時間 Task\n• 回傳狀態與結果 metadata\n\nAnritsu 角色：A2A Server", "blue_fill", "blue", 15, 11)
    card(slide, 4.80, 1.45, 3.75, 3.75, "MCP：Agent ↔ Tool", "用途\n• 儀器控制\n• iperf 執行\n• 讀取測試檔案\n• Excel 驗證\n\n既有 Anritsu MCP / tool adapter 可保留", "purple_fill", "purple", 15, 11)
    card(slide, 8.88, 1.45, 3.75, 3.75, "HTTP：Agent ↔ KB API", "用途\n• 上傳 Excel\n• 查詢攝入 task\n• Bearer token 驗證\n• strict headers / KM_Metadata\n\n既有 /api/upload/ingest 可沿用", "teal_fill", "teal", 15, 11)
    rect(slide, 1.45, 5.65, 10.35, 0.72, C["gray_fill"], C["line"], True, 1.0)
    text(slide, 1.65, 5.82, 9.95, 0.35, "推薦鏈路：KB Agent → A2A → Anritsu Agent → MCP / Instrument → Excel → 既有 KB HTTP ingest", 13, C["navy"], True, PP_ALIGN.CENTER)

    # 4. Architecture
    slide = new_slide(prs, "目標架構：只增加出站委派橋接層", "KB 核心服務保持原樣，Anritsu 保持儀器控制權", 4)
    # left
    card(slide, 0.55, 1.55, 2.25, 1.45, "使用者", "chat.html\n詢問或批准測試", "gray_fill", "navy", 14, 10)
    card(slide, 3.25, 1.55, 2.55, 1.45, "KB Agent", "OpenClaw / Agent\n決定是否委派", "blue_fill", "blue", 14, 10)
    card(slide, 6.35, 1.55, 2.55, 1.45, "Delegation Bridge", "獨立程序\nA2A Client", "purple_fill", "purple", 14, 10)
    card(slide, 9.45, 1.55, 3.25, 1.45, "Anritsu Agent", "Windows 11\nA2A Server + test executor", "orange_fill", "orange", 14, 10)
    line(slide, 2.80, 2.28, 3.25, 2.28, C["blue"], 2, True)
    line(slide, 5.80, 2.28, 6.35, 2.28, C["purple"], 2, True)
    line(slide, 8.90, 2.28, 9.45, 2.28, C["orange"], 2, True)
    text(slide, 3.00, 1.96, 0.65, 0.20, "委派", 8, C["muted"], True, PP_ALIGN.CENTER)
    text(slide, 5.92, 1.96, 0.72, 0.20, "A2A", 8, C["muted"], True, PP_ALIGN.CENTER)
    text(slide, 8.97, 1.96, 0.53, 0.20, "VPN", 8, C["muted"], True, PP_ALIGN.CENTER)
    # lower paths
    card(slide, 9.55, 3.65, 2.75, 1.25, "儀器 / iperf", "現有控制流程\n固定 profile 執行", "red_fill", "red", 13, 10)
    card(slide, 3.20, 5.25, 2.80, 1.05, "既有 KB ingest", "/api/upload/ingest\nCelery + strict contract", "teal_fill", "teal", 13, 9)
    card(slide, 7.20, 5.25, 2.20, 1.05, "Neo4j", "TestRun / Metric", "gray_fill", "navy", 13, 9)
    card(slide, 10.15, 5.25, 2.20, 1.05, "Qdrant", "chunks / citations", "gray_fill", "navy", 13, 9)
    line(slide, 10.82, 3.00, 10.82, 3.65, C["orange"], 2, True)
    line(slide, 10.82, 4.90, 4.60, 5.25, C["teal"], 2, True)
    line(slide, 6.00, 5.78, 7.20, 5.78, C["line"], 1.6, True)
    line(slide, 9.40, 5.78, 10.15, 5.78, C["line"], 1.6, True)
    text(slide, 0.72, 6.62, 11.8, 0.25, "部署原則：A2A Bridge 使用獨立程序與設定，不修改既有 KB Compose、port、資料庫或 Celery worker。", 10, C["navy"], True, PP_ALIGN.CENTER)

    # 5. Anritsu changes
    slide = new_slide(prs, "Anritsu Agent 要新增什麼", "保留既有 agent；新增 adapter 與可控的執行邊界", 5)
    modules = [
        ("a2a_server.py", "接收 A2A request\nAgent Card / JSON-RPC", "blue_fill", "blue"),
        ("agent_card.json", "名稱、版本、skills\nsecuritySchemes、capabilities", "purple_fill", "purple"),
        ("task_executor.py", "Task 狀態、timeout\n取消、重試、correlation", "orange_fill", "orange"),
        ("profile registry", "allowlisted 測試設定\n禁止任意命令", "red_fill", "red"),
        ("result handoff", "呼叫既有 uploader\n回傳 task_id / hash", "teal_fill", "teal"),
        ("audit / local state", "SQLite 或 JSON\n事件、run_id、錯誤", "gray_fill", "navy"),
    ]
    for i, (head, body, fill, accent) in enumerate(modules):
        x = 0.72 + (i % 3) * 4.15
        y = 1.55 + (i // 3) * 1.75
        card(slide, x, y, 3.55, 1.35, head, body, fill, accent, 13, 10)
    rect(slide, 0.82, 5.36, 11.65, 0.88, C["green_fill"], C["green"], True, 1.2)
    text(slide, 1.05, 5.55, 11.2, 0.46, "可沿用：儀器控制程式、iperf runner、Excel 產生器、既有 MCP tools、既有 KB upload client。", 13, C["green"], True, PP_ALIGN.CENTER)
    text(slide, 0.82, 6.50, 11.65, 0.28, "不要在 Anritsu 端安裝 Neo4j、Qdrant、Redis、Celery 或第二套 LLM。", 11, C["red"], True, PP_ALIGN.CENTER)

    # 6. Packages
    slide = new_slide(prs, "Windows 11 套件與部署", "只安裝 A2A Server 與必要的 HTTP / Excel 依賴", 6)
    card(slide, 0.72, 1.45, 5.75, 2.35, "必要環境", "Python 3.10+\n官方 A2A Python SDK\nA2A HTTP server / FastAPI extra\nhttpx：呼叫 KB ingest API\nopenpyxl：若現有 Excel parser 未提供", "blue_fill", "blue", 15, 12)
    card(slide, 6.85, 1.45, 5.75, 2.35, "可選依賴", "python-dotenv：受保護設定\ncryptography：mTLS / JWS 自行處理\nuvicorn：SDK 未自帶時補裝\nSQLite：Python 內建，不必另裝 Redis", "gray_fill", "navy", 15, 12)
    rect(slide, 0.72, 4.20, 11.88, 1.15, C["navy"], None, True)
    text(slide, 0.98, 4.42, 11.35, 0.72, 'py -3.11 -m venv .venv\n.\\.venv\\Scripts\\Activate.ps1\npip install "a2a-sdk[http-server]" httpx openpyxl python-dotenv', 12, C["white"], False, PP_ALIGN.LEFT)
    text(slide, 0.82, 5.72, 11.7, 0.65, "若現有 Anritsu agent 已有 MCP server，保留原 MCP；A2A adapter 只呼叫既有控制函式，不重做工具層。", 12, C["ink"], True, PP_ALIGN.CENTER)

    # 7. Agent Card
    slide = new_slide(prs, "Agent Card 與能力宣告", "讓 KB Bridge 能先發現能力，再選擇可用操作", 7)
    card(slide, 0.72, 1.45, 4.00, 3.75, "Endpoint", "https://<anritsu-host>/\n.well-known/agent-card.json\n\n建議 exposed interface：\nPOST /a2a\nJSON-RPC 2.0\nA2A protocolVersion 1.0\n\n僅限 VPN / HTTPS / mTLS", "blue_fill", "blue", 15, 11)
    rect(slide, 5.05, 1.45, 7.55, 3.75, C["navy"], None, True)
    text(slide, 5.30, 1.65, 7.05, 3.35,
         '{\n'
         '  "name": "Anritsu Test Agent",\n'
         '  "version": "1.0.0",\n'
         '  "supportedInterfaces": [{\n'
         '    "url": "https://anritsu/a2a",\n'
         '    "protocolVersion": "1.0"\n'
         '  }],\n'
         '  "capabilities": {"streaming": false},\n'
         '  "skills": ["run_iperf_test", "get_test_status", "cancel_test"],\n'
         '  "securitySchemes": {"a2aBearer": "..."}\n'
         '}', 10, C["white"], False, PP_ALIGN.LEFT)
    rect(slide, 0.90, 5.62, 11.55, 0.68, C["orange_fill"], C["orange"], True, 1.0)
    text(slide, 1.12, 5.78, 11.1, 0.34, "Agent Card 不得放 token、內部服務 URL、儀器密碼或主機檔案路徑。", 12, C["orange"], True, PP_ALIGN.CENTER)

    # 8. Task payload
    slide = new_slide(prs, "測試工作訊息：固定 JSON，不傳任意命令", "自然語言由 KB Agent 轉成結構化 job；Anritsu 只執行 allowlisted profile", 8)
    rect(slide, 0.72, 1.40, 5.75, 4.55, C["navy"], None, True)
    text(slide, 0.98, 1.66, 5.20, 3.95,
         '{\n'
         '  "job_type": "run_iperf_test",\n'
         '  "environment": "anritsu",\n'
         '  "profile_id":\n'
         '    "ncq2200b2v-throughput-v1",\n'
         '  "dut_model": "NCQ2200B2V",\n'
         '  "run_id": "run-20260806-001",\n'
         '  "test_cases": ["sa_dl_tcp", "sa_ul_tcp"],\n'
         '  "duration_seconds": 60,\n'
         '  "requested_by": "user-01"\n'
         '}', 11, C["white"], False, PP_ALIGN.LEFT)
    card(slide, 6.85, 1.40, 5.75, 1.35, "允許", "profile_id、run_id、test case、duration、DUT 與測試參數；由 schema 驗證型別與範圍。", "green_fill", "green", 14, 10)
    card(slide, 6.85, 3.00, 5.75, 1.35, "拒絕", "shell command、PowerShell、任意儀器 SCPI、任意路徑、未註冊 profile、超過時間上限。", "red_fill", "red", 14, 10)
    card(slide, 6.85, 4.60, 5.75, 1.35, "批准", "真實測試前可要求使用者確認；dry-run 先顯示測試摘要，再執行。", "orange_fill", "orange", 14, 10)

    # 9. Sequence
    slide = new_slide(prs, "端到端互動流程", "A2A 負責委派與回報；Excel 走既有 KB ingest", 9)
    steps = [
        ("1", "User\nrequest", "使用者要求\n執行測試"),
        ("2", "KB Agent\nplan", "建立 job JSON\n要求確認"),
        ("3", "A2A\nsubmit", "送至 Anritsu\n取得 taskId"),
        ("4", "Instrument\nrun", "儀器 + iperf\n產生 Excel"),
        ("5", "Existing\ningest", "既有 upload API\nCelery / DB"),
        ("6", "KB Agent\nanswer", "回報狀態\n引用結果"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.55 + i * 2.12
        rect(slide, x, 2.05, 1.70, 1.65, C["blue_fill"] if i < 3 else C["teal_fill"], C["blue"] if i < 3 else C["teal"], True, 1.4)
        pill(slide, x + 0.56, 1.72, 0.58, num, C["blue"] if i < 3 else C["teal"], C["white"])
        text(slide, x + 0.10, 2.25, 1.50, 0.48, head, 12, C["navy"], True, PP_ALIGN.CENTER)
        text(slide, x + 0.10, 2.86, 1.50, 0.56, body, 9, C["ink"], False, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            line(slide, x + 1.70, 2.87, x + 2.12, 2.87, C["line"], 2, True)
    rect(slide, 0.82, 4.55, 11.65, 1.10, C["gray_fill"], C["line"], True, 1.0)
    text(slide, 1.08, 4.78, 11.1, 0.64, "Correlation 必須保存：A2A taskId ↔ run_id ↔ KB ingest task_id ↔ file_hash。\nA2A 回應傳 metadata，不必把 Excel 二進位內容塞入 JSON-RPC。", 12, C["navy"], True, PP_ALIGN.CENTER)

    # 10. Task lifecycle
    slide = new_slide(prs, "Task 狀態與錯誤處理", "狀態要由普通程式碼維護，不交給 LLM 推測", 10)
    states = [
        ("submitted", "已接受", "建立 task / run_id"),
        ("working", "執行中", "控制儀器 / iperf"),
        ("input-required", "需補資料", "缺少參數或人工確認"),
        ("completed", "完成", "Excel 已上傳並攝入"),
        ("failed", "失敗", "保留錯誤與 log"),
        ("canceled", "取消", "停止可停止的工作"),
    ]
    for i, (code, name, body) in enumerate(states):
        x = 0.62 + (i % 3) * 4.18
        y = 1.45 + (i // 3) * 1.65
        accent = "green" if code == "completed" else "red" if code in {"failed", "canceled"} else "blue"
        fill = "green_fill" if accent == "green" else "red_fill" if accent == "red" else "blue_fill"
        card(slide, x, y, 3.55, 1.20, code, f"{name}\n{body}", fill, accent, 13, 10)
    rect(slide, 0.82, 5.10, 11.65, 1.10, C["orange_fill"], C["orange"], True, 1.0)
    text(slide, 1.08, 5.30, 11.1, 0.68, "重點：若 Excel 上傳成功但攝入失敗，Task 必須顯示明確失敗；不能回報整體 completed。\n可重試時沿用相同 run_id / idempotency 規則，避免重複測試或重複寫入。", 11, C["orange"], True, PP_ALIGN.CENTER)

    # 11. Existing Excel ingest
    slide = new_slide(prs, "Excel 結果流程：沿用既有功能", "A2A 不取代 KB 現有上傳與攝入", 11)
    card(slide, 0.72, 1.55, 2.35, 1.55, "Anritsu Agent", "測試完成\n產生 Excel\n計算 hash", "orange_fill", "orange", 14, 11)
    card(slide, 3.70, 1.55, 2.55, 1.55, "既有 uploader", "HTTP / MCP tool\nPOST /api/upload/ingest", "blue_fill", "blue", 14, 11)
    card(slide, 6.90, 1.55, 2.40, 1.55, "KB ingest", "Redis / Celery\nstrict contract", "teal_fill", "teal", 14, 11)
    card(slide, 9.95, 1.55, 2.50, 1.55, "Storage", "Neo4j + Qdrant\nsource assets", "gray_fill", "navy", 14, 11)
    line(slide, 3.07, 2.32, 3.70, 2.32, C["orange"], 2, True)
    line(slide, 6.25, 2.32, 6.90, 2.32, C["blue"], 2, True)
    line(slide, 9.30, 2.32, 9.95, 2.32, C["teal"], 2, True)
    text(slide, 0.82, 3.75, 11.65, 0.35, "Anritsu 端需要確認的欄位", 16, C["navy"], True)
    fields = "Authorization: Bearer <既有 ingest token>\nX-Agent-ID: anritsu-agent-01\nIdempotency-Key / X-KB-Source-System / X-KB-Environment-Id / X-KB-Run-Id\nX-KB-Artifact-Type / X-KB-Document-Id\nExcel KM_Metadata：sourceSystem、environmentId、projectId、runId、artifactType、reportSchema、documentId、idempotencyKey、generatedAt"
    rect(slide, 0.82, 4.25, 11.65, 1.62, C["navy"], None, True)
    text(slide, 1.05, 4.45, 11.2, 1.20, fields, 10, C["white"], False, PP_ALIGN.LEFT)
    text(slide, 0.82, 6.25, 11.65, 0.34, "目前 Anritsu ingest token 只用於上傳 KB；不要用同一 token 授予儀器控制權。", 11, C["red"], True, PP_ALIGN.CENTER)

    # 12. Credentials and security
    slide = new_slide(prs, "認證與安全邊界", "A2A 控制權限與 KB ingest 權限必須分離", 12)
    card(slide, 0.72, 1.45, 3.70, 3.85, "KB → Anritsu", "用途：委派與查詢 Task\n\n建議：\n• 獨立 A2A credential\n• VPN + HTTPS / mTLS\n• scope：test:run、test:status\n• IP / agent allowlist\n• audit log\n\n不可寫入 KB DB", "blue_fill", "blue", 15, 11)
    card(slide, 4.82, 1.45, 3.70, 3.85, "Anritsu → KB", "用途：上傳 Excel\n\n沿用：\n• anritsu-agent-01\n• Bearer ingest token\n• strict headers\n• KM_Metadata\n• idempotency / hash\n\n不可控制儀器", "teal_fill", "teal", 15, 11)
    card(slide, 8.92, 1.45, 3.70, 3.85, "禁止事項", "• 不接受任意 shell\n• 不公開儀器密碼\n• 不公開內部 URL\n• 不接受未註冊 profile\n• 不允許無限 duration\n• 不把 token 放 Excel / log\n• 不把資料庫掛到 Windows", "red_fill", "red", 15, 11)
    rect(slide, 0.82, 5.72, 11.65, 0.65, C["orange_fill"], C["orange"], True, 1.0)
    text(slide, 1.05, 5.87, 11.2, 0.34, "所有真實儀器操作都應可追溯到 user、run_id、profile_id、A2A taskId 與測試結果 hash。", 11, C["orange"], True, PP_ALIGN.CENTER)

    # 13. Failure and retry
    slide = new_slide(prs, "失敗、重試與冪等", "先設計失敗路徑，避免重複控制儀器或重複攝入", 13)
    rows = [
        ("A2A 連線失敗", "Bridge retry with backoff", "不建立第二個實體測試；保留 request 狀態"),
        ("Anritsu Agent 離線", "timeout → failed / retry", "使用者可稍後重新委派"),
        ("儀器執行失敗", "failed + error code", "保留 iperf log；不假裝 completed"),
        ("Excel 上傳失敗", "retry uploader", "沿用同一 run_id / idempotency key"),
        ("KB 攝入失敗", "查 task status / retry", "分開顯示 test completed 與 ingest failed"),
        ("重複 request", "return original task", "以 run_id + hash 判定 duplicate/conflict"),
    ]
    rect(slide, 0.72, 1.45, 11.88, 0.48, C["navy"], None, True)
    text(slide, 0.92, 1.56, 2.10, 0.24, "情境", 11, C["white"], True, PP_ALIGN.LEFT)
    text(slide, 3.35, 1.56, 3.25, 0.24, "程式行為", 11, C["white"], True, PP_ALIGN.LEFT)
    text(slide, 7.05, 1.56, 5.10, 0.24, "驗收要求", 11, C["white"], True, PP_ALIGN.LEFT)
    for i, row in enumerate(rows):
        y = 1.99 + i * 0.70
        rect(slide, 0.72, y, 11.88, 0.62, C["gray_fill"] if i % 2 == 0 else C["white"], C["line"], False, 0.6)
        text(slide, 0.92, y + 0.10, 2.10, 0.34, row[0], 10, C["ink"], True)
        text(slide, 3.35, y + 0.10, 3.25, 0.34, row[1], 10, C["blue"], True)
        text(slide, 7.05, y + 0.10, 5.10, 0.34, row[2], 10, C["ink"])
    text(slide, 0.82, 6.55, 11.65, 0.30, "A2A taskId、run_id、KB task_id、file_hash 必須一起寫入 audit / local state。", 11, C["navy"], True, PP_ALIGN.CENTER)

    # 14. Risk priority map
    slide = new_slide(prs, "風險總覽與處理優先順序", "先阻止會破壞既有測試或產生錯誤數據的風險", 14)
    risks = [
        ("P0", "儀器同時被控制", "結果失真 / 儀器狀態混亂", "instrument lock + 單一 owner", "red"),
        ("P0", "任意命令注入", "未授權操作 / 主機風險", "profile allowlist + schema", "red"),
        ("P0", "同進程或套件衝突", "原 agent 無法啟動", "sidecar + 獨立 venv", "red"),
        ("P1", "重試造成重複測試", "儀器重跑 / 重複報告", "run_id + idempotency", "orange"),
        ("P1", "Excel / ingest 部分成功", "KB 顯示錯誤完成狀態", "雙狀態 + reconciliation", "orange"),
        ("P1", "網路或服務重啟", "Task 遺失 / 無法回報", "journal + resume + outbox", "orange"),
        ("P2", "版本或 schema 漂移", "A2A / 報告不相容", "版本固定 + contract test", "blue"),
        ("P2", "Log 與監控不足", "事故無法追蹤", "correlation + audit", "blue"),
    ]
    rect(slide, 0.55, 1.30, 12.20, 0.44, C["navy"], None, True)
    for x, w, value in [(0.72, 0.65, "等級"), (1.55, 2.60, "風險"), (4.35, 3.15, "可能影響"), (7.75, 4.55, "必要解法")]:
        text(slide, x, 1.40, w, 0.22, value, 10, C["white"], True)
    for i, (priority, risk, impact, solution, accent) in enumerate(risks):
        y = 1.82 + i * 0.61
        rect(slide, 0.55, y, 12.20, 0.53, C[accent + "_fill"], C["line"], False, 0.5)
        text(slide, 0.72, y + 0.08, 0.65, 0.30, priority, 10, C[accent], True, PP_ALIGN.CENTER)
        text(slide, 1.55, y + 0.08, 2.60, 0.30, risk, 10, C["ink"], True)
        text(slide, 4.35, y + 0.08, 3.15, 0.30, impact, 9, C["ink"])
        text(slide, 7.75, y + 0.08, 4.55, 0.30, solution, 9, C[accent], True)
    text(slide, 0.82, 6.82, 11.65, 0.22, "Go-live 前：所有 P0 必須通過；P1 必須有自動恢復或明確人工處置；P2 必須納入維運。", 10, C["navy"], True, PP_ALIGN.CENTER)

    # 15. Process and dependency isolation
    slide = new_slide(prs, "風險一：程序、套件與設定衝突", "A2A 故障不得連帶停止既有 Anritsu 測試功能", 15)
    card(slide, 0.72, 1.40, 3.75, 3.80, "可能影響", "• SDK 升級破壞既有套件\n• A2A exception 終止原 agent\n• port / env 名稱互相覆蓋\n• log 或工作目錄混用\n• Windows service 啟停連動\n• rollback 時無法還原", "red_fill", "red", 15, 11)
    card(slide, 4.80, 1.40, 3.75, 3.80, "設計解法", "• A2A 採獨立 sidecar process\n• 使用獨立 Python venv\n• 獨立 port、.env、log、PID\n• 只透過 stable adapter 呼叫原功能\n• pin SDK / dependency versions\n• A2A 可單獨停用與移除", "green_fill", "green", 15, 11)
    card(slide, 8.88, 1.40, 3.75, 3.80, "驗收方式", "• 停止 A2A 後手動測試正常\n• A2A crash 不影響原 agent\n• 安裝前後套件清單可比較\n• port conflict 明確失敗\n• 舊啟動捷徑仍可使用\n• rollback 後功能完全恢復", "blue_fill", "blue", 15, 11)
    rect(slide, 0.82, 5.66, 11.65, 0.70, C["orange_fill"], C["orange"], True, 1.0)
    text(slide, 1.05, 5.82, 11.2, 0.36, "禁止把 a2a-sdk 直接 pip install 到原 agent 的 production venv，除非先完成 dependency lock 與完整回歸。", 11, C["orange"], True, PP_ALIGN.CENTER)

    # 16. Instrument concurrency and resource isolation
    slide = new_slide(prs, "風險二：儀器並發與主機資源競爭", "同一時間只能有一個受控 owner 操作同一台儀器", 16)
    card(slide, 0.72, 1.40, 3.75, 3.80, "可能影響", "• 人工與 A2A 同時下命令\n• 兩個 A2A task 同時執行\n• 儀器留在錯誤狀態\n• iperf port / process 衝突\n• CPU、NIC、磁碟競爭\n• cancel 後背景程序仍存在", "red_fill", "red", 15, 11)
    card(slide, 4.80, 1.40, 3.75, 3.80, "設計解法", "• 每台儀器建立 exclusive lock\n• lock owner 保存 taskId / run_id\n• 預設最大 concurrency = 1\n• profile 定義 port / duration 上限\n• finally block 執行儀器 cleanup\n• cancel 同時終止 iperf child process", "green_fill", "green", 15, 11)
    card(slide, 8.88, 1.40, 3.75, 3.80, "驗收方式", "• 人工測試中 A2A 回 busy\n• A2A 中人工入口顯示占用\n• 第二個 task 不得控制儀器\n• timeout 後 lock 必須釋放\n• iperf process 無殘留\n• 儀器回到已知 safe state", "blue_fill", "blue", 15, 11)
    rect(slide, 0.82, 5.66, 11.65, 0.70, C["navy"], None, True)
    text(slide, 1.05, 5.82, 11.2, 0.36, "Lock 不只是檔案旗標：必須有 owner、TTL、heartbeat 與受控釋放，避免錯誤 task 解鎖別人的工作。", 11, C["white"], True, PP_ALIGN.CENTER)

    # 17. Data and ingest integrity
    slide = new_slide(prs, "風險三：Excel、結果與 KB 攝入不一致", "測試完成不等於攝入完成；兩種狀態必須分開", 17)
    card(slide, 0.72, 1.40, 3.75, 3.80, "可能影響", "• Excel 尚未寫完就被上傳\n• 同名檔案被覆蓋\n• hash / KM_Metadata 不一致\n• 測試成功但 ingest 失敗\n• Neo4j 成功、Qdrant 失敗\n• 重送造成 duplicate/conflict", "red_fill", "red", 15, 11)
    card(slide, 4.80, 1.40, 3.75, 3.80, "設計解法", "• 先寫 .tmp，再 atomic rename\n• 每個 run 使用唯一目錄\n• 關閉 workbook 後才計算 hash\n• 分開 test_status / ingest_status\n• 沿用 strict headers / KM_Metadata\n• outbox + reconciliation + polling", "green_fill", "green", 15, 11)
    card(slide, 8.88, 1.40, 3.75, 3.80, "驗收方式", "• 上傳中的 Excel 可正常開啟\n• hash 與 server 回傳一致\n• 相同 run/hash 回原 task\n• 相同 run/不同 hash 回 conflict\n• Neo4j/Qdrant 都完成才 completed\n• ingest failed 可補傳且不重測", "blue_fill", "blue", 15, 11)
    rect(slide, 0.82, 5.66, 11.65, 0.70, C["orange_fill"], C["orange"], True, 1.0)
    text(slide, 1.05, 5.82, 11.2, 0.36, "A2A completion artifact 應回傳 run_id、report_name、file_hash、KB task_id、test_status、ingest_status。", 11, C["orange"], True, PP_ALIGN.CENTER)

    # 18. Network, authentication, and command safety
    slide = new_slide(prs, "風險四：網路、認證與命令安全", "把 A2A 視為儀器控制入口，而不是一般聊天 API", 18)
    card(slide, 0.72, 1.40, 3.75, 3.80, "可能影響", "• Windows port 暴露公網\n• token 被複製或寫入 log\n• ingest token 被用來控制儀器\n• Agent Card 洩漏內部資訊\n• LLM / prompt injection 產生命令\n• 未授權使用者啟動真實測試", "red_fill", "red", 15, 11)
    card(slide, 4.80, 1.40, 3.75, 3.80, "設計解法", "• VPN + HTTPS，正式採 mTLS\n• A2A 與 ingest credentials 分離\n• test:run / test:status scopes\n• profile allowlist + Pydantic schema\n• 真實測試要求 user approval\n• token 只放 secret store / env", "green_fill", "green", 15, 11)
    card(slide, 8.88, 1.40, 3.75, 3.80, "驗收方式", "• 無憑證與錯誤 scope 被拒絕\n• VPN 外無法連線\n• raw shell / SCPI 被拒絕\n• 超過 duration / case 上限被拒絕\n• Agent Card 無秘密與內部路徑\n• log 掃描找不到 token", "blue_fill", "blue", 15, 11)
    rect(slide, 0.82, 5.66, 11.65, 0.70, C["red_fill"], C["red"], True, 1.0)
    text(slide, 1.05, 5.82, 11.2, 0.36, "目前 anritsu-agent-01 token 僅供 KB Excel ingest；A2A 控制入口必須建立另一組 credential。", 11, C["red"], True, PP_ALIGN.CENTER)

    # 19. Recovery, compatibility, and regression
    slide = new_slide(prs, "風險五：任務恢復、版本相容與回歸", "系統重啟或版本更新後，不能遺失狀態或破壞手動流程", 19)
    card(slide, 0.72, 1.40, 3.75, 3.80, "可能影響", "• Agent 重啟後 task 消失\n• working 狀態永久卡住\n• A2A SDK / spec 版本不同\n• job schema 新舊不相容\n• Windows 時間不同步\n• 新功能破壞原本 Excel / 操作", "red_fill", "red", 15, 11)
    card(slide, 4.80, 1.40, 3.75, 3.80, "設計解法", "• SQLite task journal + startup recovery\n• heartbeat / stale-task policy\n• 固定 A2A 1.0 與 package lock\n• job_schema_version + report_schema\n• NTP / UTC timestamp\n• feature flag + one-command rollback", "green_fill", "green", 15, 11)
    card(slide, 8.88, 1.40, 3.75, 3.80, "驗收方式", "• working 時重啟並正確恢復\n• stale task 可轉 failed / reconcile\n• 不支援版本明確拒絕\n• 舊 job / Excel fixture 仍通過\n• 人工測試完整回歸\n• 關閉 feature flag 即恢復原狀", "blue_fill", "blue", 15, 11)
    rect(slide, 0.82, 5.66, 11.65, 0.70, C["gray_fill"], C["line"], True, 1.0)
    text(slide, 1.05, 5.82, 11.2, 0.36, "最低 audit 欄位：timestamp、caller、profile_id、run_id、A2A taskId、KB task_id、state、error_code、file_hash。", 10, C["navy"], True, PP_ALIGN.CENTER)

    # 20. Implementation phases
    slide = new_slide(prs, "實作分階段", "先驗證通訊，再逐步開放真實儀器操作", 20)
    phases = [
        ("Phase 0", "契約與 dry-run", "Agent Card\njob schema\n不控制儀器", "blue"),
        ("Phase 1", "單一 profile", "固定 NCQ profile\n實際執行 iperf", "orange"),
        ("Phase 2", "結果攝入", "沿用 Excel upload\n驗證 Neo4j/Qdrant", "teal"),
        ("Phase 3", "可靠性", "取消、timeout\nretry、audit", "purple"),
        ("Phase 4", "擴充", "多 profile\nAmarisoft 共用契約", "green"),
    ]
    for i, (phase, head, body, accent) in enumerate(phases):
        x = 0.58 + i * 2.55
        rect(slide, x, 1.75, 2.12, 3.20, C[accent + "_fill"], C[accent], True, 1.4)
        pill(slide, x + 0.53, 1.50, 1.05, phase, C[accent], C["white"])
        text(slide, x + 0.12, 2.10, 1.88, 0.58, head, 13, C[accent], True, PP_ALIGN.CENTER)
        text(slide, x + 0.15, 2.93, 1.82, 1.05, body, 11, C["ink"], False, PP_ALIGN.CENTER)
        if i < len(phases) - 1:
            line(slide, x + 2.12, 3.33, x + 2.55, 3.33, C["line"], 1.8, True)
    rect(slide, 0.82, 5.45, 11.65, 0.90, C["gray_fill"], C["line"], True, 1.0)
    text(slide, 1.08, 5.66, 11.1, 0.46, "每一階段都必須能獨立 rollback；任何階段失敗不得影響既有 KB chat、search、ingest 與資料庫。", 12, C["navy"], True, PP_ALIGN.CENTER)

    # 21. Acceptance checklist
    slide = new_slide(prs, "Anritsu Agent 驗收清單", "完成下列項目才算可交付 KB 整合", 21)
    checks = [
        "Agent Card 可透過 HTTPS 取得，沒有洩漏秘密或內部路徑",
        "KB Bridge 可發送 A2A job，Anritsu 回傳 taskId",
        "只接受 allowlisted profile，不接受任意 command",
        "dry-run、真實測試、取消、timeout、failed 都有明確狀態",
        "run_id、A2A taskId、KB task_id、file_hash 可相互追蹤",
        "Excel 仍由既有 /api/upload/ingest 上傳，不新增第二套 uploader",
        "KM_Metadata 與 strict headers 通過驗證，duplicate/conflict 行為正確",
        "Neo4j 與 Qdrant 都完成後才回報攝入 completed",
        "Windows agent 重啟後不會遺失未完成 task 或 audit 記錄",
        "KB 原有 chat、search、ingest 與服務 health check 無回歸",
    ]
    for i, item in enumerate(checks):
        x = 0.82 + (i % 2) * 6.00
        y = 1.43 + (i // 2) * 0.82
        rect(slide, x, y, 5.55, 0.58, C["green_fill"], C["green"], True, 0.8)
        text(slide, x + 0.10, y + 0.08, 0.32, 0.35, "✓", 14, C["green"], True, PP_ALIGN.CENTER)
        text(slide, x + 0.52, y + 0.08, 4.82, 0.35, item, 9, C["ink"], True)
    rect(slide, 0.82, 5.85, 11.65, 0.62, C["orange_fill"], C["orange"], True, 1.0)
    text(slide, 1.05, 6.00, 11.2, 0.34, "驗收入口：https://61.216.9.52:3030/chat.html；外部測試同時保存 Anritsu agent log 與 KB task response。", 10, C["orange"], True, PP_ALIGN.CENTER)

    # 22. Handoff / sources
    slide = new_slide(prs, "交付給 Anritsu Agent 的實作順序", "先完成 adapter，再逐步開放實際測試", 22)
    card(slide, 0.72, 1.45, 5.75, 3.75, "本次要做", "1. 建立 Python venv 與 A2A SDK\n2. 寫 Agent Card 與 A2A server\n3. 定義固定 job schema / profile\n4. 包裝既有儀器控制函式\n5. 建立 task state / audit\n6. 接回既有 Excel uploader\n7. 用 dry-run 驗證整條鏈路", "blue_fill", "blue", 15, 12)
    card(slide, 6.85, 1.45, 5.75, 3.75, "不可做", "1. 不把 KB DB 掛到 Windows\n2. 不新建第二套 Excel ingest\n3. 不共用 ingest token 作儀器控制\n4. 不讓 LLM 生成任意命令\n5. 不直接暴露 Windows 到公網\n6. 不將未完成測試回報 completed", "red_fill", "red", 15, 12)
    rect(slide, 0.82, 5.65, 11.65, 0.72, C["navy"], None, True)
    text(slide, 1.08, 5.82, 11.1, 0.35, "官方參考：a2a-protocol.org/latest/specification、github.com/a2aproject/a2a-python", 10, C["white"], True, PP_ALIGN.CENTER)
    text(slide, 0.82, 6.60, 11.65, 0.28, "完成後先提交 Agent Card、job schema、測試 log 與 dry-run 結果，再進入真實儀器測試。", 11, C["navy"], True, PP_ALIGN.CENTER)

    prs.save(OUT)
    print(f"wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
