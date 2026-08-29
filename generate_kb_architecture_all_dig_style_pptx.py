from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "knowledge_base_architecture_all_dig_style.pptx"

W, H = 13.333, 7.5
FONT = "Noto Sans CJK TC"

C = {
    "white": "FFFFFF",
    "ink": "263541",
    "muted": "66727C",
    "line": "727C84",
    "grid": "D8DEE3",
    "navy": "173B59",
    "navy_fill": "EAF1F6",
    "blue": "28729A",
    "blue_fill": "EDF7FB",
    "purple": "755595",
    "purple_fill": "F7F1FA",
    "orange": "C58A39",
    "orange_fill": "FFF4E2",
    "teal": "267B7C",
    "teal_fill": "EDF8F6",
    "green": "659462",
    "green_fill": "F0F7EE",
    "gray": "747B80",
    "gray_fill": "F3F5F6",
    "red": "A85F5F",
    "red_fill": "FAEEEE",
}


def rgb(value):
    return RGBColor.from_string(value)


def add_text(slide, x, y, w, h, value, size=9, color=None, bold=False,
             align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if fit else MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = frame.margin_right = Pt(2)
    frame.margin_top = frame.margin_bottom = Pt(1)
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color or C["ink"])
    return box


def shape(slide, x, y, w, h, fill, line, radius=True, width=1.2):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    obj = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    obj.fill.solid()
    obj.fill.fore_color.rgb = rgb(fill)
    obj.line.color.rgb = rgb(line)
    obj.line.width = Pt(width)
    return obj


def node(slide, x, y, w, h, title, detail, color, fill, title_size=8.5, detail_size=6.6):
    obj = shape(slide, x, y, w, h, fill, color, width=1.05)
    frame = obj.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Pt(5)
    frame.margin_top = frame.margin_bottom = Pt(3)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(2)
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    if detail:
        p2 = frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = p2.space_after = Pt(0)
        r2 = p2.add_run()
        r2.text = detail
        r2.font.name = FONT
        r2.font.size = Pt(detail_size)
        r2.font.color.rgb = rgb(C["muted"])
    return obj


def store(slide, x, y, w, h, title, detail, color, fill):
    obj = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CAN, Inches(x), Inches(y), Inches(w), Inches(h))
    obj.fill.solid()
    obj.fill.fore_color.rgb = rgb(fill)
    obj.line.color.rgb = rgb(color)
    obj.line.width = Pt(1.05)
    frame = obj.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Pt(4)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(8.2)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = detail
    r2.font.name = FONT
    r2.font.size = Pt(6.2)
    r2.font.color.rgb = rgb(C["muted"])
    return obj


def boundary(slide, x, y, w, h, title, color):
    obj = shape(slide, x, y, w, h, C["white"], color, width=1.35)
    label_w = min(max(1.8, len(title) * 0.12), w - 0.35)
    shape(slide, x + 0.18, y - 0.11, label_w, 0.28, C["white"], C["white"], radius=False, width=0)
    add_text(slide, x + 0.24, y - 0.10, label_w - 0.12, 0.25, title, 8.5, color, True,
             PP_ALIGN.LEFT, fit=True)
    return obj


def _arrow_end(connector):
    ln = connector._element.spPr.get_or_add_ln()
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    ln.append(tail)


def edge(slide, x1, y1, x2, y2, label=None, lx=None, ly=None, elbow=False,
         color=None, width=1.05, dashed=False):
    kind = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    conn = slide.shapes.add_connector(kind, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(color or C["line"])
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = 2
    _arrow_end(conn)
    if label:
        x = lx if lx is not None else (x1 + x2) / 2 - 0.55
        y = ly if ly is not None else (y1 + y2) / 2 - 0.10
        add_text(slide, x, y, 1.1, 0.20, label, 5.8, C["muted"], False, fit=True)
    return conn


def page_title(slide, number, title, subtitle):
    add_text(slide, 0.36, 0.15, 0.38, 0.25, f"{number:02d}", 7.5, C["muted"], True)
    add_text(slide, 0.78, 0.12, 5.4, 0.34, title, 14, C["ink"], True, PP_ALIGN.LEFT)
    add_text(slide, 6.15, 0.17, 6.75, 0.22, subtitle, 6.8, C["muted"], False, PP_ALIGN.RIGHT)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.36), Inches(0.52), Inches(12.93), Inches(0.52))
    line.line.color.rgb = rgb(C["grid"])
    line.line.width = Pt(0.7)


def footer(slide, note):
    add_text(slide, 0.38, 7.24, 8.8, 0.14, note, 5.6, C["muted"], False, PP_ALIGN.LEFT)
    add_text(slide, 10.2, 7.24, 2.72, 0.14, "Knowledge Base · AS-IS · 2026.07", 5.6, C["muted"], False, PP_ALIGN.RIGHT)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(C["white"])
    return slide


def overview(prs):
    s = blank_slide(prs)
    page_title(s, 1, "Knowledge Base 完整系統架構", "入口、應用服務、AI 推論與統一知識資料平台")

    # Boundaries and connections are intentionally sparse, matching all_dig.jpg.
    boundary(s, 0.35, 1.45, 3.05, 3.70, "存取與既有 Web 系統", C["blue"])
    boundary(s, 3.75, 1.45, 5.25, 3.70, "Knowledge Base AI 應用系統", C["purple"])
    boundary(s, 2.10, 5.55, 8.75, 1.35, "Unified Knowledge Data Platform", C["teal"])

    # Nginx is shared ingress; role-specific paths are traced by matching colors.
    edge(s, 1.86, 2.42, 1.18, 2.75, "HTTPS", 1.34, 2.47, elbow=True)
    edge(s, 1.86, 2.42, 2.57, 2.75, "HTTPS", 2.03, 2.47, elbow=True)
    edge(s, 1.86, 2.42, 1.18, 3.83, "WSS", 1.20, 3.12, elbow=True)
    edge(s, 1.86, 2.42, 2.57, 3.83, "HTTPS", 2.41, 3.12, elbow=True)
    edge(s, 1.82, 3.12, 4.12, 2.10, "Search API", 2.85, 2.38, elbow=True, color=C["blue"])
    edge(s, 3.11, 3.12, 4.12, 2.10, "Admin API", 3.36, 2.65, elbow=True, color=C["purple"])
    edge(s, 1.82, 4.20, 5.95, 2.10, "WebSocket /ws", 3.31, 3.31, elbow=True, color=C["blue"])
    edge(s, 3.11, 4.20, 6.84, 2.73, "Ingest Task", 4.48, 3.65, elbow=True, color=C["green"])
    edge(s, 5.02, 2.45, 5.02, 2.72)
    edge(s, 6.84, 2.45, 6.84, 2.72)
    edge(s, 5.46, 3.52, 5.46, 3.82)
    edge(s, 6.40, 4.22, 8.08, 4.22, "Task", 6.95, 4.02)
    edge(s, 10.82, 2.45, 10.82, 2.78)
    edge(s, 10.82, 3.55, 10.82, 3.88)
    edge(s, 8.82, 4.22, 9.72, 4.22, "Tool / Session", 8.84, 4.01)
    edge(s, 4.88, 4.65, 3.32, 5.85, "Vector", 3.75, 5.15, elbow=True)
    edge(s, 5.58, 4.65, 5.17, 5.85, "Graph", 5.12, 5.15, elbow=True)
    edge(s, 6.28, 4.65, 7.01, 5.85, "Files", 6.48, 5.15, elbow=True)
    edge(s, 8.47, 4.65, 8.85, 5.85, "State", 8.38, 5.15, elbow=True)

    node(s, 3.20, 0.64, 1.82, 0.68, "使用者", "Search / Chat → 查詢與對話", C["blue"], C["blue_fill"], 9.2, 6.1)
    node(s, 5.23, 0.64, 1.82, 0.68, "管理者", "Admin → 管理 API / 任務", C["purple"], C["purple_fill"], 9.2, 6.1)
    node(s, 7.26, 0.64, 1.82, 0.68, "知識擁護者", "Upload → Ingest / 建索引", C["green"], C["green_fill"], 9.2, 6.1)
    add_text(s, 9.34, 0.72, 3.04, 0.42,
             "角色路徑：依同色節點與 API 線追蹤\n灰色線：三個角色共用基礎設施",
             6.2, C["muted"], False, PP_ALIGN.LEFT)
    node(s, 0.84, 1.76, 2.04, 0.66, "Nginx TLS Gateway", "HTTPS :3030 / 靜態資源", C["blue"], C["blue_fill"])
    node(s, 0.54, 2.75, 1.28, 0.73, "Search UI", "查詢 / 引用", C["blue"], C["blue_fill"])
    node(s, 1.93, 2.75, 1.18, 0.73, "Admin UI", "管理 / 設定", C["purple"], C["purple_fill"])
    node(s, 0.54, 3.83, 1.28, 0.73, "Chat UI", "即時對話", C["blue"], C["blue_fill"])
    node(s, 1.93, 3.83, 1.18, 0.73, "Upload / Watch", "文件 / 索引", C["green"], C["green_fill"])

    node(s, 4.12, 1.76, 1.80, 0.69, "FastAPI / Uvicorn", "REST · Task · Admin", C["purple"], C["purple_fill"])
    node(s, 5.95, 1.76, 1.80, 0.69, "WebSocket Proxy", "/ws · event relay", C["purple"], C["purple_fill"])
    node(s, 4.12, 2.73, 1.80, 0.79, "SearchEngine", "Basic · Deep · Hybrid+", C["purple"], C["purple_fill"])
    node(s, 5.95, 2.73, 1.80, 0.79, "Document Pipeline", "Convert · OCR · Chunk", C["purple"], C["purple_fill"])
    node(s, 4.52, 3.82, 1.88, 0.83, "Celery Workers", "Search / Ingest / Beat", C["purple"], C["purple_fill"])
    node(s, 7.33, 3.82, 1.49, 0.83, "Redis", "Broker · Cache · Locks", C["purple"], C["purple_fill"])

    node(s, 9.72, 1.78, 2.20, 0.67, "OpenClaw Gateway", "Agent · Tools · Session", C["orange"], C["orange_fill"])
    node(s, 9.72, 2.78, 2.20, 0.77, "OpenClaw Runtime", "Identity · Workspace · Skills", C["orange"], C["orange_fill"])
    node(s, 9.72, 3.88, 2.20, 0.77, "Ollama LLM", "qwen3-coder-next · :11434", C["orange"], C["orange_fill"])

    store(s, 2.45, 5.86, 1.72, 0.78, "Qdrant", "BGE 768D vectors", C["teal"], C["teal_fill"])
    store(s, 4.29, 5.86, 1.72, 0.78, "Neo4j", "Knowledge graph", C["teal"], C["teal_fill"])
    store(s, 6.13, 5.86, 1.72, 0.78, "File Store", "Raw / Processed / Assets", C["teal"], C["teal_fill"])
    store(s, 7.97, 5.86, 1.72, 0.78, "Runtime State", "Task / Chat / Memory", C["teal"], C["teal_fill"])
    footer(s, "灰線為主要資料流；色彩只標示責任域，所有元件均可在 PowerPoint 編輯。")


def query_chat(prs):
    s = blank_slide(prs)
    page_title(s, 2, "查詢與 OpenClaw 對話架構", "同步 UI、非同步搜尋工作與 Agent 串流事件的雙路徑")
    boundary(s, 0.45, 1.15, 2.25, 5.35, "使用者介面", C["blue"])
    boundary(s, 3.00, 1.15, 3.25, 5.35, "FastAPI 與任務編排", C["purple"])
    boundary(s, 6.60, 1.15, 2.55, 5.35, "知識檢索", C["teal"])
    boundary(s, 9.50, 1.15, 3.35, 5.35, "AI Agent 與回答", C["orange"])

    edge(s, 2.26, 2.02, 3.42, 2.02, "POST /search", 2.38, 1.82)
    edge(s, 4.18, 2.43, 4.18, 3.00, "enqueue", 4.30, 2.60)
    edge(s, 5.72, 3.42, 7.02, 3.42, "query", 5.95, 3.22)
    edge(s, 7.86, 3.83, 7.86, 4.43, "recall", 7.98, 4.04)
    edge(s, 8.68, 4.83, 9.92, 4.83, "context", 8.86, 4.63)
    edge(s, 4.18, 3.84, 4.18, 4.62, "status", 4.30, 4.14)
    edge(s, 3.42, 5.03, 2.26, 5.03, "poll / result", 2.42, 4.83)

    edge(s, 2.26, 3.03, 3.42, 3.03, "WSS /ws", 2.43, 2.83)
    edge(s, 5.72, 5.68, 9.92, 5.68, "agent events", 7.48, 5.48)
    edge(s, 11.11, 4.08, 11.11, 3.50)
    edge(s, 11.11, 2.76, 11.11, 2.35)

    node(s, 0.84, 1.64, 1.42, 0.78, "Search UI", "問題 / 模式", C["blue"], C["blue_fill"])
    node(s, 0.84, 2.64, 1.42, 0.78, "Chat UI", "對話 / 引用", C["blue"], C["blue_fill"])
    node(s, 0.84, 4.62, 1.42, 0.82, "Result View", "答案 / Sources", C["blue"], C["blue_fill"])

    node(s, 3.42, 1.63, 2.30, 0.80, "REST API", "/search · /tasks/{id}", C["purple"], C["purple_fill"])
    node(s, 3.42, 3.00, 2.30, 0.84, "Redis + Celery Search", "Broker · Result · Lock", C["purple"], C["purple_fill"])
    node(s, 3.42, 4.62, 2.30, 0.82, "WebSocket Proxy", "Session queue · relay", C["purple"], C["purple_fill"])

    node(s, 7.02, 3.00, 1.68, 0.83, "SearchEngine", "Mode router", C["teal"], C["teal_fill"])
    store(s, 6.92, 4.43, 1.88, 0.83, "Qdrant + Neo4j", "Vector + Graph", C["teal"], C["teal_fill"])

    node(s, 9.92, 1.63, 2.38, 0.72, "Ollama LLM", "Answer synthesis", C["orange"], C["orange_fill"])
    node(s, 9.92, 2.76, 2.38, 0.74, "OpenClaw Gateway", "Agent / Tool / Session", C["orange"], C["orange_fill"])
    node(s, 9.92, 4.08, 2.38, 0.83, "Citation Package", "來源不足時保留揭露", C["green"], C["green_fill"])
    node(s, 9.92, 5.32, 2.38, 0.72, "Rendered Reply", "串流事件 / 最終回答", C["green"], C["green_fill"])
    footer(s, "搜尋路徑以 Celery 非同步執行；Chat 路徑由 WebSocket 代理 OpenClaw，兩者共用知識來源。")


def ingestion(prs):
    s = blank_slide(prs)
    page_title(s, 3, "文件攝入與知識建立架構", "來源接入、轉檔/OCR、分塊與三種索引策略")
    boundary(s, 0.42, 1.15, 2.15, 5.45, "文件來源", C["blue"])
    boundary(s, 2.92, 1.15, 3.10, 5.45, "攝入處理管線", C["purple"])
    boundary(s, 6.38, 1.15, 3.10, 5.45, "資料導向路由", C["orange"])
    boundary(s, 9.83, 1.15, 3.05, 5.45, "知識儲存", C["teal"])

    edge(s, 2.18, 2.03, 3.30, 2.03, "upload", 2.37, 1.83)
    edge(s, 2.18, 3.12, 3.30, 3.12, "watch", 2.39, 2.92)
    edge(s, 2.18, 4.22, 3.30, 4.22, "manual", 2.38, 4.02)
    edge(s, 4.44, 2.45, 4.44, 2.78)
    edge(s, 4.44, 3.60, 4.44, 3.94)
    edge(s, 5.58, 4.38, 6.80, 4.38, "chunks", 5.79, 4.18)
    edge(s, 7.97, 2.38, 10.20, 2.05, "Report", 8.78, 1.85, elbow=True)
    edge(s, 7.97, 3.35, 10.20, 3.45, "Vector", 8.79, 3.16, elbow=True)
    edge(s, 7.97, 4.38, 10.20, 4.80, "Semantic", 8.74, 4.51, elbow=True)
    edge(s, 7.97, 5.43, 10.20, 5.95, "archive", 8.80, 5.60, elbow=True)

    node(s, 0.79, 1.62, 1.39, 0.82, "API Upload", "PDF / DOCX / XLSX", C["blue"], C["blue_fill"])
    node(s, 0.79, 2.71, 1.39, 0.82, "Watch Folder", "自動偵測", C["blue"], C["blue_fill"])
    node(s, 0.79, 3.81, 1.39, 0.82, "Raw / Manual", "批次匯入", C["blue"], C["blue_fill"])
    node(s, 0.79, 5.05, 1.39, 0.82, "SHA-256", "重複檔案控制", C["green"], C["green_fill"])

    node(s, 3.30, 1.63, 2.28, 0.82, "Redis Task State", "任務 TTL / Lock", C["purple"], C["purple_fill"])
    node(s, 3.30, 2.78, 2.28, 0.82, "Celery Ingest Worker", "清理舊索引後寫入", C["purple"], C["purple_fill"])
    node(s, 3.30, 3.94, 2.28, 0.88, "FileConverter + OCR", "MarkItDown / Assets", C["purple"], C["purple_fill"])
    node(s, 3.30, 5.13, 2.28, 0.82, "Chunker + Metadata", "段落 / 頁碼 / 來源", C["purple"], C["purple_fill"])

    node(s, 6.80, 1.63, 1.17, 0.75, "REPORT", "報告結構", C["orange"], C["orange_fill"])
    node(s, 6.80, 2.97, 1.17, 0.75, "VECTOR", "Lab / Project", C["orange"], C["orange_fill"])
    node(s, 6.80, 4.00, 1.17, 0.75, "SEMANTIC", "4G / 5G / WiFi", C["orange"], C["orange_fill"])
    node(s, 6.80, 5.05, 1.17, 0.75, "PROCESSED", "原文保存", C["orange"], C["orange_fill"])
    node(s, 8.14, 1.63, 0.94, 0.75, "Graph", "nodes", C["orange"], C["orange_fill"])
    node(s, 8.14, 4.00, 0.94, 0.75, "Ollama", "entities", C["orange"], C["orange_fill"])

    store(s, 10.20, 1.63, 2.00, 0.82, "Neo4j", "Document / Entity / Report", C["teal"], C["teal_fill"])
    store(s, 10.20, 3.04, 2.00, 0.82, "Qdrant", "BGE 768D embeddings", C["teal"], C["teal_fill"])
    store(s, 10.20, 4.43, 2.00, 0.82, "Neo4j + Qdrant", "語意實體 + 向量", C["teal"], C["teal_fill"])
    store(s, 10.20, 5.68, 2.00, 0.72, "File Store", "processed / assets", C["teal"], C["teal_fill"])
    footer(s, "同一管線依文件類型選擇 Report Graph、Vector-only 或 Semantic Graph；避免一種索引策略套用所有資料。")


def retrieval(prs):
    s = blank_slide(prs)
    page_title(s, 4, "混合檢索、資料融合與引用架構", "以可追溯來源組裝答案，避免模型脫離知識庫內容")
    boundary(s, 0.42, 1.25, 2.40, 5.20, "查詢理解", C["blue"])
    boundary(s, 3.18, 1.25, 3.30, 5.20, "平行召回", C["teal"])
    boundary(s, 6.84, 1.25, 2.80, 5.20, "融合與治理", C["purple"])
    boundary(s, 10.00, 1.25, 2.88, 5.20, "回答與呈現", C["green"])

    edge(s, 2.37, 2.25, 3.62, 2.25, "vector", 2.56, 2.05)
    edge(s, 2.37, 3.34, 3.62, 3.34, "graph", 2.57, 3.14)
    edge(s, 2.37, 4.43, 3.62, 4.43, "file", 2.62, 4.23)
    edge(s, 5.94, 2.25, 7.25, 3.31, elbow=True)
    edge(s, 5.94, 3.34, 7.25, 3.31)
    edge(s, 5.94, 4.43, 7.25, 3.31, elbow=True)
    edge(s, 8.55, 3.73, 8.55, 4.20)
    edge(s, 9.22, 4.62, 10.42, 4.62, "context + refs", 9.34, 4.42)
    edge(s, 11.35, 4.20, 11.35, 3.62)
    edge(s, 11.35, 2.84, 11.35, 2.38)

    node(s, 0.87, 1.72, 1.50, 0.85, "Query Analyzer", "意圖 / 關鍵字 / 範圍", C["blue"], C["blue_fill"])
    node(s, 0.87, 2.93, 1.50, 0.82, "Mode Router", "Basic · Deep · Hybrid+", C["blue"], C["blue_fill"])
    node(s, 0.87, 4.16, 1.50, 0.82, "Filters", "類別 / 專案 / 時間", C["blue"], C["blue_fill"])

    store(s, 3.62, 1.82, 2.32, 0.86, "Qdrant", "向量相似度召回", C["teal"], C["teal_fill"])
    store(s, 3.62, 2.91, 2.32, 0.86, "Neo4j", "實體關聯與報告結構", C["teal"], C["teal_fill"])
    store(s, 3.62, 4.00, 2.32, 0.86, "Processed Files", "原文片段與資產", C["teal"], C["teal_fill"])

    node(s, 7.25, 2.90, 1.30, 0.83, "Fusion", "合併候選", C["purple"], C["purple_fill"])
    node(s, 7.89, 4.20, 1.33, 0.83, "Rerank", "業務排序 / 去重", C["purple"], C["purple_fill"])
    node(s, 7.52, 5.36, 1.94, 0.66, "Guardrail", "有來源不得宣稱無資料", C["red"], C["red_fill"])

    node(s, 10.42, 1.70, 1.86, 0.68, "Ollama / OpenClaw", "回答生成", C["green"], C["green_fill"])
    node(s, 10.42, 2.84, 1.86, 0.78, "Citation Package", "來源 / 頁碼 / 分數", C["green"], C["green_fill"])
    node(s, 10.42, 4.20, 1.86, 0.82, "Search / Chat Reply", "答案與可追溯引用", C["green"], C["green_fill"])
    footer(s, "召回與答案生成分離：先取得可驗證的 context/citations，再交由 LLM 或 Agent 組裝回答。")


def deployment(prs):
    s = blank_slide(prs)
    page_title(s, 5, "現行站台與 On-Prem 部署架構", "清楚區分容器內服務、主機 AI Runtime 與持久化資料")
    boundary(s, 0.42, 1.18, 6.08, 5.65, "現行站台：knowledge-base", C["blue"])
    boundary(s, 6.84, 1.18, 6.05, 5.65, "B2B / On-Prem Release", C["purple"])

    edge(s, 1.51, 2.16, 1.51, 2.55, "TLS", 1.63, 2.28)
    edge(s, 1.51, 3.27, 1.51, 3.68)
    edge(s, 2.52, 4.09, 3.56, 4.09, "jobs", 2.69, 3.89)
    edge(s, 2.52, 4.09, 3.56, 5.29, elbow=True)
    edge(s, 4.40, 4.50, 4.40, 4.88)
    edge(s, 5.03, 5.29, 5.55, 5.29, "host", 4.98, 5.08)

    edge(s, 7.93, 2.16, 7.93, 2.55, "TLS", 8.05, 2.28)
    edge(s, 7.93, 3.27, 7.93, 3.68)
    edge(s, 8.94, 4.09, 9.98, 4.09, "jobs", 9.10, 3.89)
    edge(s, 8.94, 4.09, 9.98, 5.29, elbow=True)
    edge(s, 10.82, 4.50, 10.82, 4.88)
    edge(s, 11.45, 5.29, 11.98, 5.29, "host", 11.41, 5.08)

    node(s, 0.84, 1.53, 1.34, 0.63, "Browser", "使用者", C["navy"], C["navy_fill"])
    node(s, 0.84, 2.55, 1.34, 0.72, "Nginx", ":3030 → TLS 443", C["blue"], C["blue_fill"])
    node(s, 0.84, 3.68, 1.68, 0.82, "FastAPI Web", ":8000 · 4 workers", C["blue"], C["blue_fill"])
    node(s, 3.56, 3.68, 1.68, 0.82, "Celery", "Search / Ingest / Beat", C["purple"], C["purple_fill"])
    store(s, 3.56, 4.88, 1.68, 0.82, "Redis + Neo4j", "容器持久化", C["teal"], C["teal_fill"])
    node(s, 5.55, 4.88, 0.66, 0.82, "Host AI", "Qdrant\nOllama\nOpenClaw", C["orange"], C["orange_fill"], 7.5, 5.8)
    node(s, 0.84, 5.28, 1.68, 0.76, "Frontend Files", "Vue / chat.html", C["gray"], C["gray_fill"])

    node(s, 7.26, 1.53, 1.34, 0.63, "Browser", "客戶端", C["navy"], C["navy_fill"])
    node(s, 7.26, 2.55, 1.34, 0.72, "Nginx", "可配置 HTTPS", C["blue"], C["blue_fill"])
    node(s, 7.26, 3.68, 1.68, 0.82, "FastAPI Web", "Release container", C["blue"], C["blue_fill"])
    node(s, 9.98, 3.68, 1.68, 0.82, "Celery", "Search / Ingest / Beat", C["purple"], C["purple_fill"])
    store(s, 9.98, 4.88, 1.68, 0.82, "Redis / Neo4j / Qdrant", "Compose volumes", C["teal"], C["teal_fill"])
    node(s, 11.98, 4.88, 0.66, 0.82, "Host AI", "Ollama\nOpenClaw", C["orange"], C["orange_fill"], 7.5, 5.8)
    node(s, 7.26, 5.28, 1.68, 0.76, "Installer / Preflight", "Config · Upgrade · Backup", C["gray"], C["gray_fill"])

    add_text(s, 0.84, 6.26, 5.25, 0.27, "現況：Qdrant、Ollama、OpenClaw 位於 host；應用與工作程序由 Compose 管理。", 6.3, C["blue"], False, PP_ALIGN.LEFT)
    add_text(s, 7.26, 6.26, 5.25, 0.27, "Release：Qdrant 隨 Compose 交付；Ollama、OpenClaw 保留為主機外部相依。", 6.3, C["purple"], False, PP_ALIGN.LEFT)
    footer(s, "兩種部署拓撲不可混為一談；安裝前應檢查主機 runtime、連線埠與持久化路徑。")


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "Knowledge Base Architecture - all_dig style"
    prs.core_properties.subject = "Enterprise architecture diagrams"
    prs.core_properties.author = "Knowledge Base Architecture Team"
    overview(prs)
    query_chat(prs)
    ingestion(prs)
    retrieval(prs)
    deployment(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
