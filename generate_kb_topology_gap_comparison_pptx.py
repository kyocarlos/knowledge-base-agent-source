from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "knowledge_base_topology_gap_comparison.pptx"

W, H = 13.333, 7.5
FONT = "Noto Sans CJK TC"

C = {
    "white": "FFFFFF",
    "ink": "263541",
    "muted": "66727C",
    "line": "7A848B",
    "grid": "D8DEE3",
    "blue": "28729A",
    "blue_fill": "EDF7FB",
    "purple": "755595",
    "purple_fill": "F7F1FA",
    "orange": "C58A39",
    "orange_fill": "FFF4E2",
    "teal": "267B7C",
    "teal_fill": "EDF8F6",
    "green": "4D8C57",
    "green_fill": "EAF5EC",
    "amber": "C17B19",
    "amber_fill": "FFF3D9",
    "red": "B45151",
    "red_fill": "FBEAEA",
    "gray": "747B80",
    "gray_fill": "F3F5F6",
}

STATUS = {
    "ready": ("已具備", C["green"], C["green_fill"]),
    "partial": ("部分具備", C["amber"], C["amber_fill"]),
    "missing": ("尚缺", C["red"], C["red_fill"]),
}


def rgb(value):
    return RGBColor.from_string(value)


def add_text(slide, x, y, w, h, value, size=9, color=None, bold=False,
             align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if fit else MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = frame.margin_right = Pt(2)
    frame.margin_top = frame.margin_bottom = Pt(1)
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color or C["ink"])
    return box


def shape(slide, x, y, w, h, fill, line, radius=True, width=1.2, dashed=False):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    obj = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    obj.fill.solid()
    obj.fill.fore_color.rgb = rgb(fill)
    obj.line.color.rgb = rgb(line)
    obj.line.width = Pt(width)
    if dashed:
        obj.line.dash_style = 2
    return obj


def boundary(slide, x, y, w, h, title, color):
    obj = shape(slide, x, y, w, h, C["white"], color, width=1.35)
    label_w = min(max(1.8, len(title) * 0.12), w - 0.35)
    shape(slide, x + 0.18, y - 0.11, label_w, 0.28, C["white"], C["white"], radius=False, width=0)
    add_text(slide, x + 0.24, y - 0.10, label_w - 0.12, 0.25, title, 8.5, color, True,
             PP_ALIGN.LEFT)
    return obj


def _arrow_end(connector):
    ln = connector._element.spPr.get_or_add_ln()
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    ln.append(tail)


def edge(slide, x1, y1, x2, y2, label=None, lx=None, ly=None, elbow=False,
         color=None, width=1.0, dashed=False):
    kind = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    conn = slide.shapes.add_connector(kind, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(color or C["line"])
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = 2
    _arrow_end(conn)
    if label:
        x = lx if lx is not None else (x1 + x2) / 2 - 0.50
        y = ly if ly is not None else (y1 + y2) / 2 - 0.10
        add_text(slide, x, y, 1.0, 0.18, label, 5.5, C["muted"])
    return conn


def status_node(slide, x, y, w, h, title, detail, domain_color, domain_fill, status,
                title_size=8.0, detail_size=5.8):
    label, status_color, status_fill = STATUS[status]
    obj = shape(slide, x, y, w, h, domain_fill, status_color, width=2.0,
                dashed=status != "ready")
    frame = obj.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Pt(5)
    frame.margin_top = Pt(6)
    frame.margin_bottom = Pt(3)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(1)
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = rgb(domain_color)
    if detail:
        p2 = frame.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = p2.space_after = Pt(0)
        r2 = p2.add_run()
        r2.text = detail
        r2.font.name = FONT
        r2.font.size = Pt(detail_size)
        r2.font.color.rgb = rgb(C["muted"])

    pill_w = 0.54 if status != "partial" else 0.67
    pill = shape(slide, x + w - pill_w - 0.05, y + 0.04, pill_w, 0.17,
                 status_fill, status_color, width=0.7)
    pill.text_frame.clear()
    pill.text_frame.margin_left = pill.text_frame.margin_right = Pt(1)
    pill.text_frame.margin_top = pill.text_frame.margin_bottom = Pt(0)
    pill.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p3 = pill.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    rr = p3.add_run()
    rr.text = label
    rr.font.name = FONT
    rr.font.size = Pt(4.5)
    rr.font.bold = True
    rr.font.color.rgb = rgb(status_color)
    return obj


def status_store(slide, x, y, w, h, title, detail, status):
    label, status_color, status_fill = STATUS[status]
    obj = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CAN, Inches(x), Inches(y), Inches(w), Inches(h))
    obj.fill.solid()
    obj.fill.fore_color.rgb = rgb(C["teal_fill"])
    obj.line.color.rgb = rgb(status_color)
    obj.line.width = Pt(2.0)
    if status != "ready":
        obj.line.dash_style = 2
    frame = obj.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(7.6)
    r.font.bold = True
    r.font.color.rgb = rgb(C["teal"])
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = detail
    r2.font.name = FONT
    r2.font.size = Pt(5.4)
    r2.font.color.rgb = rgb(C["muted"])
    add_text(slide, x + 0.05, y + h - 0.19, w - 0.10, 0.14, label, 4.5,
             status_color, True)
    return obj


def page_title(slide, number, title, subtitle):
    add_text(slide, 0.36, 0.15, 0.38, 0.25, f"{number:02d}", 7.5, C["muted"], True)
    add_text(slide, 0.78, 0.12, 6.2, 0.34, title, 14, C["ink"], True, PP_ALIGN.LEFT)
    add_text(slide, 6.65, 0.17, 6.28, 0.22, subtitle, 6.6, C["muted"], False, PP_ALIGN.RIGHT)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.36), Inches(0.52), Inches(12.93), Inches(0.52))
    line.line.color.rgb = rgb(C["grid"])
    line.line.width = Pt(0.7)


def footer(slide, note):
    add_text(slide, 0.38, 7.24, 9.4, 0.14, note, 5.4, C["muted"], False, PP_ALIGN.LEFT)
    add_text(slide, 10.15, 7.24, 2.77, 0.14, "Knowledge Base · GAP · 2026.07", 5.4,
             C["muted"], False, PP_ALIGN.RIGHT)


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(C["white"])
    return slide


def legend(slide, x=9.46, y=0.67):
    for idx, key in enumerate(("ready", "partial", "missing")):
        label, color, fill = STATUS[key]
        shape(slide, x + idx * 1.05, y, 0.18, 0.18, fill, color, width=1.0,
              dashed=key != "ready")
        add_text(slide, x + 0.21 + idx * 1.05, y, 0.76, 0.18, label, 5.6, color,
                 True, PP_ALIGN.LEFT)


def gap_overview(prs):
    s = blank_slide(prs)
    page_title(s, 1, "Knowledge Base 現況 vs 最終目標", "綠框＝程式庫已具備；橘框＝能力存在但尚未完整整合；紅框＝明確缺口")
    legend(s)

    boundary(s, 0.35, 1.45, 3.05, 3.70, "存取與既有 Web 系統", C["blue"])
    boundary(s, 3.75, 1.45, 5.25, 3.70, "Knowledge Base AI 應用系統", C["purple"])
    boundary(s, 2.10, 5.55, 8.75, 1.35, "Unified Knowledge Data Platform", C["teal"])

    # Main data paths. Dashed amber paths cross a lifecycle/integration gap.
    edge(s, 1.86, 2.42, 1.18, 2.75, "HTTPS", 1.34, 2.47, elbow=True)
    edge(s, 1.86, 2.42, 2.57, 2.75, "HTTPS", 2.03, 2.47, elbow=True)
    edge(s, 1.86, 2.42, 1.18, 3.83, "WSS", 1.20, 3.12, elbow=True)
    edge(s, 1.86, 2.42, 2.57, 3.83, "HTTPS", 2.41, 3.12, elbow=True)
    edge(s, 1.82, 3.12, 4.12, 2.10, "Search API", 2.85, 2.38, elbow=True, color=C["blue"])
    edge(s, 3.11, 3.12, 4.12, 2.10, "Admin API", 3.36, 2.65, elbow=True, color=C["purple"])
    edge(s, 1.82, 4.20, 5.95, 2.10, "WebSocket /ws", 3.31, 3.31, elbow=True, color=C["blue"])
    edge(s, 3.11, 4.20, 6.84, 2.73, "Ingest Task", 4.48, 3.65, elbow=True, color=C["green"])
    edge(s, 5.02, 2.45, 5.02, 2.72)
    edge(s, 6.84, 2.45, 6.84, 2.72)
    edge(s, 5.46, 3.52, 5.46, 3.82)
    edge(s, 6.40, 4.22, 8.08, 4.22, "Task", 6.95, 4.02)
    edge(s, 10.82, 2.45, 10.82, 2.78, color=C["amber"], dashed=True)
    edge(s, 10.82, 3.55, 10.82, 3.88, color=C["amber"], dashed=True)
    edge(s, 8.82, 4.22, 9.72, 4.22, "Tool / Session", 8.84, 4.01,
         color=C["amber"], dashed=True)
    edge(s, 4.88, 4.65, 3.32, 5.85, "Vector", 3.75, 5.15, elbow=True)
    edge(s, 5.58, 4.65, 5.17, 5.85, "Graph", 5.12, 5.15, elbow=True)
    edge(s, 6.28, 4.65, 7.01, 5.85, "Files", 6.48, 5.15, elbow=True)
    edge(s, 8.47, 4.65, 8.85, 5.85, "State", 8.38, 5.15, elbow=True,
         color=C["amber"], dashed=True)

    status_node(s, 3.20, 0.64, 1.82, 0.68, "使用者", "Search / Chat 路徑已有；身分治理未完整", C["blue"], C["blue_fill"], "partial", 8.6, 5.3)
    status_node(s, 5.23, 0.64, 1.82, 0.68, "管理者", "Admin 功能已有；RBAC 未形成", C["purple"], C["purple_fill"], "partial", 8.6, 5.3)
    status_node(s, 7.26, 0.64, 1.82, 0.68, "知識擁護者", "Upload / Watch 已有；角色未強制", C["green"], C["green_fill"], "partial", 8.6, 5.3)

    status_node(s, 0.84, 1.76, 2.04, 0.66, "Nginx TLS Gateway", "HTTPS :3030 / 靜態資源", C["blue"], C["blue_fill"], "ready")
    status_node(s, 0.54, 2.75, 1.28, 0.73, "Search UI", "查詢 / 引用", C["blue"], C["blue_fill"], "ready")
    status_node(s, 1.93, 2.75, 1.18, 0.73, "Admin UI", "管理 / 設定", C["purple"], C["purple_fill"], "ready")
    status_node(s, 0.54, 3.83, 1.28, 0.73, "Chat UI", "即時對話", C["blue"], C["blue_fill"], "ready")
    status_node(s, 1.93, 3.83, 1.18, 0.73, "Upload / Watch", "文件 / 索引", C["green"], C["green_fill"], "ready")

    status_node(s, 4.12, 1.76, 1.80, 0.69, "FastAPI / Uvicorn", "REST · Task · Admin", C["purple"], C["purple_fill"], "ready")
    status_node(s, 5.95, 1.76, 1.80, 0.69, "WebSocket Proxy", "/ws · event relay", C["purple"], C["purple_fill"], "ready")
    status_node(s, 4.12, 2.73, 1.80, 0.79, "SearchEngine", "Basic · Deep · Vector · Hybrid+", C["purple"], C["purple_fill"], "ready")
    status_node(s, 5.95, 2.73, 1.80, 0.79, "Document Pipeline", "Convert / Chunk 已有；OCR 條件式", C["purple"], C["purple_fill"], "partial")
    status_node(s, 4.52, 3.82, 1.88, 0.83, "Celery Workers", "Search / Ingest / Beat", C["purple"], C["purple_fill"], "ready")
    status_node(s, 7.33, 3.82, 1.49, 0.83, "Redis", "Broker · Cache · Locks", C["purple"], C["purple_fill"], "ready")

    status_node(s, 9.72, 1.78, 2.20, 0.67, "OpenClaw Gateway", "已串接；不在 KB Compose 生命週期", C["orange"], C["orange_fill"], "partial")
    status_node(s, 9.72, 2.78, 2.20, 0.77, "OpenClaw Runtime", "Identity / Workspace / Skills 已整合", C["orange"], C["orange_fill"], "partial")
    status_node(s, 9.72, 3.88, 2.20, 0.77, "Ollama LLM", "qwen3-coder-next · :11434；外部相依", C["orange"], C["orange_fill"], "partial")

    status_store(s, 2.45, 5.86, 1.72, 0.78, "Qdrant", "向量庫已具備；獨立容器", "partial")
    status_store(s, 4.29, 5.86, 1.72, 0.78, "Neo4j", "知識圖譜 / Compose volume", "ready")
    status_store(s, 6.13, 5.86, 1.72, 0.78, "File Store", "Raw / Processed / Assets", "ready")
    status_store(s, 7.97, 5.86, 1.72, 0.78, "Runtime State", "Task / locks 有；持久 Chat Memory 未完整", "partial")
    footer(s, "判定基準：目前 repository 能力與部署邊界；『部分具備』不代表不可用，而是尚未達到最終統一拓樸。")


def current_backbone(prs):
    s = blank_slide(prs)
    page_title(s, 2, "目前已具備的端到端主幹", "核心查詢、聊天、攝入與資料寫入皆已有對應程式與部署元件")
    legend(s, 9.46, 0.67)

    boundary(s, 0.48, 1.20, 2.08, 5.55, "入口與 UI", C["blue"])
    boundary(s, 2.88, 1.20, 3.02, 5.55, "API 與工作程序", C["purple"])
    boundary(s, 6.22, 1.20, 2.82, 5.55, "知識資料", C["teal"])
    boundary(s, 9.36, 1.20, 3.48, 5.55, "外部 AI Runtime", C["orange"])

    edge(s, 2.15, 2.03, 3.31, 2.03, "HTTPS")
    edge(s, 2.15, 3.05, 3.31, 3.05, "REST")
    edge(s, 2.15, 4.07, 3.31, 4.07, "WSS")
    edge(s, 4.31, 2.42, 4.31, 2.72, "enqueue", 4.43, 2.49)
    edge(s, 5.42, 3.14, 6.64, 2.30, "vector", 5.74, 2.46, elbow=True)
    edge(s, 5.42, 3.14, 6.64, 3.55, "graph", 5.75, 3.20, elbow=True)
    edge(s, 5.42, 4.35, 6.64, 4.80, "files", 5.76, 4.39, elbow=True)
    edge(s, 5.42, 4.35, 7.50, 5.82, "state", 6.25, 5.04, elbow=True)
    edge(s, 5.42, 3.14, 9.78, 2.29, "context / prompt", 7.34, 2.55, elbow=True, color=C["amber"], dashed=True)
    edge(s, 5.42, 4.07, 9.78, 3.54, "agent events", 7.33, 3.65, elbow=True, color=C["amber"], dashed=True)
    edge(s, 10.93, 3.95, 10.93, 4.63, "infer", 11.05, 4.17, color=C["amber"], dashed=True)

    status_node(s, 0.82, 1.62, 1.33, 0.80, "Nginx", "TLS / 靜態檔", C["blue"], C["blue_fill"], "ready")
    status_node(s, 0.82, 2.65, 1.33, 0.80, "Search / Admin", "Vue UI", C["blue"], C["blue_fill"], "ready")
    status_node(s, 0.82, 3.67, 1.33, 0.80, "Chat UI", "引用 / 串流", C["blue"], C["blue_fill"], "ready")
    status_node(s, 0.82, 4.69, 1.33, 0.80, "Upload / Watch", "上傳 / 定時掃描", C["green"], C["green_fill"], "ready")

    status_node(s, 3.31, 1.62, 2.11, 0.80, "FastAPI", "API / Task / WebSocket", C["purple"], C["purple_fill"], "ready")
    status_node(s, 3.31, 2.72, 2.11, 0.84, "SearchEngine", "Basic / Deep / Vector / Hybrid+", C["purple"], C["purple_fill"], "ready")
    status_node(s, 3.31, 3.67, 2.11, 0.80, "Celery", "Search / Ingest / Beat", C["purple"], C["purple_fill"], "ready")
    status_node(s, 3.31, 4.69, 2.11, 0.80, "Document Pipeline", "Convert / Chunk / conditional OCR", C["purple"], C["purple_fill"], "partial")

    status_store(s, 6.64, 1.78, 1.86, 0.86, "Qdrant", "向量檢索", "partial")
    status_store(s, 6.64, 3.09, 1.86, 0.86, "Neo4j", "圖譜 / 結構", "ready")
    status_store(s, 6.64, 4.40, 1.86, 0.86, "File Store", "原文 / 轉檔 / 資產", "ready")
    status_store(s, 6.64, 5.48, 1.86, 0.86, "Redis State", "任務 / 快取 / 鎖", "ready")

    status_node(s, 9.78, 1.85, 2.30, 0.88, "OpenClaw Gateway", "已代理 Agent / Session 事件", C["orange"], C["orange_fill"], "partial")
    status_node(s, 9.78, 3.10, 2.30, 0.85, "OpenClaw Runtime", "Identity / Workspace / Skills", C["orange"], C["orange_fill"], "partial")
    status_node(s, 9.78, 4.63, 2.30, 0.85, "Ollama", "qwen3-coder-next / :11434", C["orange"], C["orange_fill"], "partial")

    add_text(s, 9.75, 5.73, 2.44, 0.44,
             "整合程式已存在；差距在於部署、健康檢查、升級與故障復原尚未由 KB 統一管理。",
             5.9, C["amber"], True, PP_ALIGN.LEFT)
    footer(s, "已具備主幹：Browser → Nginx → FastAPI/Celery → Qdrant/Neo4j/File Store → OpenClaw/Ollama。")


def gap_roadmap(prs):
    s = blank_slide(prs)
    page_title(s, 3, "尚缺部分與補齊順序", "優先補共同機制；避免只針對單一機器、session 或部署路徑修補")

    # Four gap cards.
    cards = [
        (0.55, 1.24, "P0", "角色與存取治理", "現況", "UI 與 API 路徑存在，但三種角色尚未由登入、權限與稽核強制區隔。",
         "目標", "Identity Provider / RBAC / API scope / audit log", C["red"]),
        (3.72, 1.24, "P0", "Runtime State 完整化", "現況", "Redis 已承擔 task、cache、locks；Chat/Memory 分散於 OpenClaw session 與暫態狀態。",
         "目標", "統一 session metadata、持久 memory、TTL/清理與備份規則", C["red"]),
        (6.89, 1.24, "P1", "外部 AI 生命週期", "現況", "OpenClaw / Ollama 已串通，但仍是 KB Compose 外部相依。",
         "目標", "preflight、版本契約、health check、升級/回復與明確責任邊界", C["amber"]),
        (10.06, 1.24, "P1", "資料平台一致化", "現況", "Neo4j 與 File Store 由 KB 管理；Qdrant 以獨立容器運作，OCR 也依模型與格式條件啟用。",
         "目標", "統一部署/備份/監控；OCR capability check 與顯性失敗", C["amber"]),
    ]
    for x, y, priority, title, h1, b1, h2, b2, color in cards:
        shape(s, x, y, 2.72, 2.28, C["white"], color, width=1.5, dashed=True)
        shape(s, x + 0.13, y + 0.12, 0.42, 0.24, C["red_fill"] if priority == "P0" else C["amber_fill"], color, width=0.8)
        add_text(s, x + 0.13, y + 0.12, 0.42, 0.24, priority, 6.3, color, True)
        add_text(s, x + 0.62, y + 0.10, 1.92, 0.29, title, 9.2, C["ink"], True, PP_ALIGN.LEFT)
        add_text(s, x + 0.15, y + 0.49, 0.42, 0.22, h1, 6.1, C["muted"], True, PP_ALIGN.LEFT)
        add_text(s, x + 0.15, y + 0.71, 2.42, 0.57, b1, 6.2, C["muted"], False, PP_ALIGN.LEFT)
        add_text(s, x + 0.15, y + 1.33, 0.42, 0.22, h2, 6.1, color, True, PP_ALIGN.LEFT)
        add_text(s, x + 0.15, y + 1.55, 2.42, 0.55, b2, 6.2, C["ink"], False, PP_ALIGN.LEFT)

    add_text(s, 0.55, 3.86, 2.10, 0.28, "建議落地拓樸", 9.2, C["ink"], True, PP_ALIGN.LEFT)
    edge(s, 2.23, 5.11, 3.23, 5.11, "token / role", 2.26, 4.91, color=C["red"], dashed=True)
    edge(s, 5.25, 5.11, 6.24, 5.11, "state API", 5.28, 4.91, color=C["red"], dashed=True)
    edge(s, 8.26, 5.11, 9.26, 5.11, "contract", 8.30, 4.91, color=C["amber"], dashed=True)
    edge(s, 11.28, 5.11, 12.20, 5.11, "backup", 11.29, 4.91, color=C["amber"], dashed=True)

    status_node(s, 0.74, 4.68, 1.49, 0.86, "Identity / RBAC", "三角色與 API scope", C["red"], C["red_fill"], "missing")
    status_node(s, 3.23, 4.68, 2.02, 0.86, "FastAPI Gateway", "驗證 / 授權 / 稽核", C["purple"], C["purple_fill"], "partial")
    status_store(s, 6.24, 4.64, 2.02, 0.94, "Runtime State", "Session / Chat / Memory", "missing")
    status_node(s, 9.26, 4.68, 2.02, 0.86, "AI Runtime Manager", "OpenClaw / Ollama lifecycle", C["orange"], C["orange_fill"], "missing")
    status_store(s, 12.20, 4.64, 0.73, 0.94, "Data Ops", "統一備份", "missing")

    shape(s, 0.55, 6.07, 12.38, 0.63, C["gray_fill"], C["grid"], width=0.8)
    add_text(s, 0.78, 6.15, 11.92, 0.45,
             "驗證門檻：不同電腦 / session / runtime / 部署路徑皆需通過；服務缺失時必須顯性失敗，並可由健康檢查定位責任域。",
             7.0, C["ink"], True, PP_ALIGN.LEFT)
    footer(s, "本頁是目標差距，不代表全部都要容器化；重點是統一契約、可觀測、可備份、可回復。")


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "Knowledge Base Topology Gap Comparison"
    prs.core_properties.subject = "Current-state and target-state topology comparison"
    prs.core_properties.author = "Knowledge Base Architecture Team"
    gap_overview(prs)
    current_backbone(prs)
    gap_roadmap(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
