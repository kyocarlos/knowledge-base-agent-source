from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/home/da40_ai_gb10/knowledge-base/manual_ingest_customer_intro.pptx")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0.0, radius=False):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def set_text(frame, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Sans CJK TC"):
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)


def textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    set_text(shape.text_frame, text, size, color, bold=bold, align=align)
    return shape


def card(slide, left, top, width, height, title, body, accent, fill="FFFFFF"):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=fill,
        line="CBD5E1",
        radius=True,
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(7)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Noto Sans CJK TC"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = rgb(accent)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Noto Sans CJK TC"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = rgb("475569")
    return shape


def badge(slide, left, top, width, text, fill, color):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        Inches(0.30),
        fill=fill,
        line=fill,
        radius=True,
    )
    set_text(shape.text_frame, text, 8.5, color, align=PP_ALIGN.CENTER)
    return shape


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    for left, top, width, height, color in [
        (Inches(-0.4), Inches(-0.3), Inches(3.8), Inches(3.8), "DBEAFE"),
        (Inches(9.3), Inches(-0.3), Inches(4.3), Inches(3.8), "D1FAE5"),
        (Inches(8.8), Inches(4.8), Inches(4.4), Inches(2.7), "FEF3C7"),
    ]:
        deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        deco.fill.solid()
        deco.fill.fore_color.rgb = rgb(color)
        deco.fill.transparency = 0.80
        deco.line.fill.background()

    border = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.22),
        Inches(0.18),
        Inches(12.89),
        Inches(7.14),
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    border.fill.transparency = 1


def title_block(slide, title, subtitle):
    textbox(slide, Inches(0.52), Inches(0.32), Inches(4.0), Inches(0.20), "Knowledge Base / Manual Ingest", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.56), Inches(8.5), Inches(0.42), title, 24, "0F172A", bold=True)
    textbox(slide, Inches(0.52), Inches(0.98), Inches(9.5), Inches(0.34), subtitle, 12, "475569")


def add_cover(slide):
    add_background(slide)
    title_block(
        slide,
        "知識庫手動攝入方式",
        "使用者上傳文件後，系統會自動轉成可搜尋內容，並寫入知識圖譜與向量資料庫。",
    )

    badge(slide, Inches(10.52), Inches(0.48), Inches(2.10), "客戶說明用簡報", "EFF6FF", "1D4ED8")

    card(slide, Inches(0.72), Inches(1.72), Inches(3.75), Inches(1.40), "手動攝入是什麼", "手動攝入指的是由使用者主動上傳文件，系統在收到文件後立即執行轉換、切 chunk 與寫入資料庫。", "2563EB")
    card(slide, Inches(0.72), Inches(3.28), Inches(3.75), Inches(1.40), "適合的情境", "適合新報告上線、臨時文件更新，或需要立即讓文件加入搜尋時使用。", "16A34A")
    card(slide, Inches(0.72), Inches(4.84), Inches(3.75), Inches(1.40), "客戶怎麼理解", "可以把它理解成「上傳文件後，系統幫你整理成可查詢的資料」。", "D97706")

    card(slide, Inches(5.00), Inches(1.95), Inches(7.25), Inches(2.15), "完整流程一句話", "先上傳文件，再由系統自動完成格式轉換、內容切分、關聯建立與向量索引，最後就可以在搜尋頁直接查詢。", "047857")

    # simple flow
    blocks = [
        (0.82, "上傳文件", "Excel / PDF / Word / PPTX", "2563EB"),
        (3.05, "自動轉換", "變成可讀內容", "16A34A"),
        (5.28, "切成 chunks", "分段方便搜尋", "D97706"),
        (7.51, "寫入 Neo4j", "建立關聯脈絡", "1D4ED8"),
        (9.74, "寫入 Qdrant", "建立語意搜尋索引", "047857"),
    ]
    for x, t, d, a in blocks:
        card(slide, Inches(x), Inches(5.28), Inches(1.95), Inches(0.88), t, d, a, fill="FFFFFF")

    for x1, x2 in [(2.75, 3.05), (4.98, 5.28), (7.21, 7.51), (9.44, 9.74)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(5.72), Inches(x2), Inches(5.72))
        line.line.color.rgb = rgb("94A3B8")
        line.line.width = Pt(1.3)

    footer = slide.shapes.add_textbox(Inches(0.74), Inches(6.60), Inches(11.9), Inches(0.22))
    set_text(footer.text_frame, "本簡報將用一個實際案例說明：文件如何從上傳一路變成可搜尋資料。", 9.5, "64748B", align=PP_ALIGN.RIGHT)


def add_flow_slide(slide):
    add_background(slide)
    title_block(
        slide,
        "手動攝入流程",
        "以下流程可直接對客戶說明：上傳後，系統會自動處理，不需要人工逐步搬移資料。",
    )
    badge(slide, Inches(10.52), Inches(0.48), Inches(2.10), "操作步驟", "ECFDF5", "047857")

    steps = [
        ("1", "選擇文件", "使用者在上傳頁面選擇一份測試報告。", "2563EB"),
        ("2", "按下上傳並攝入", "系統開始接收文件並建立處理任務。", "16A34A"),
        ("3", "轉換與切分", "系統將文件轉成文字，再切成多個 chunk。", "D97706"),
        ("4", "寫入資料庫", "Neo4j 保存關聯，Qdrant 保存內容向量。", "1D4ED8"),
        ("5", "可供查詢", "完成後，這份文件就能在搜尋頁被查到。", "047857"),
    ]

    y = 1.72
    for idx, (num, title, body, accent) in enumerate(steps):
        badge(slide, Inches(0.78), Inches(y + idx * 0.86), Inches(0.42), num, "EFF6FF", accent)
        card(slide, Inches(1.30), Inches(y - 0.04 + idx * 0.86), Inches(5.85), Inches(0.68), title, body, accent)

    # right-side explanation block
    card(slide, Inches(7.55), Inches(1.82), Inches(4.95), Inches(1.45), "為什麼要手動攝入", "當客戶有最新報告、臨時資料或特定專案文件時，可以立即把文件加入知識庫，不必等待排程掃描。", "2563EB")
    card(slide, Inches(7.55), Inches(3.45), Inches(4.95), Inches(1.45), "系統會做哪些事", "系統會自動完成格式轉換、章節整理、關聯建立與向量索引，讓文件可被搜尋與引用。", "16A34A")
    card(slide, Inches(7.55), Inches(5.08), Inches(4.95), Inches(1.45), "客戶聽得懂的說法", "你只要上傳文件，後面的整理與入庫會自動完成。", "D97706")

    footer = slide.shapes.add_textbox(Inches(0.74), Inches(6.80), Inches(11.9), Inches(0.18))
    set_text(footer.text_frame, "手動攝入的核心，是讓新文件快速進入搜尋與問答流程。", 9.5, "64748B", align=PP_ALIGN.RIGHT)


def add_example_slide(slide):
    add_background(slide)
    title_block(
        slide,
        "實際範例說明",
        "以下以一份真實類型的測試報告示意手動攝入後，資料如何被整理進知識庫。",
    )
    badge(slide, Inches(10.52), Inches(0.48), Inches(2.10), "範例文件", "FEF3C7", "B45309")

    card(slide, Inches(0.72), Inches(1.80), Inches(4.20), Inches(1.25), "範例文件", "SIT-TR-SC-NR-Throughput-SCU2060-n79-EV-V13.8.xlsx", "2563EB")
    card(slide, Inches(0.72), Inches(3.20), Inches(4.20), Inches(1.25), "上傳方式", "在上傳頁選擇檔案後，按下「上傳並攝入」。", "16A34A")
    card(slide, Inches(0.72), Inches(4.60), Inches(4.20), Inches(1.25), "預期結果", "文件會先轉成文字，接著切成多個 chunk，最後可直接被搜尋。", "D97706")

    # mini pipeline with actual example sections
    card(slide, Inches(5.35), Inches(1.85), Inches(2.00), Inches(0.95), "2. Introduction", "成為第一個 chunk", "1D4ED8")
    card(slide, Inches(7.62), Inches(1.85), Inches(2.00), Inches(0.95), "3. Test Result Summary", "成為第二個 chunk", "047857")
    card(slide, Inches(9.89), Inches(1.85), Inches(2.00), Inches(0.95), "4. Performance Test", "成為第三個 chunk", "D97706")

    for x1, x2 in [(7.35, 7.62), (9.62, 9.89)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(2.30), Inches(x2), Inches(2.30))
        line.line.color.rgb = rgb("94A3B8")
        line.line.width = Pt(1.3)

    card(slide, Inches(5.35), Inches(3.30), Inches(6.54), Inches(1.55), "進入 Neo4j", "建立文件與章節之間的關係，方便系統知道這份報告屬於哪個專案，以及哪些章節彼此相關。", "2563EB")
    card(slide, Inches(5.35), Inches(5.00), Inches(6.54), Inches(1.55), "進入 Qdrant", "每個 chunk 都會被轉成向量，之後使用者詢問相近內容時，系統可以快速召回最相關段落。", "16A34A")

    footer = slide.shapes.add_textbox(Inches(0.74), Inches(6.82), Inches(11.9), Inches(0.18))
    set_text(footer.text_frame, "這個例子可以直接對客戶說明：手動上傳一份報告後，系統如何自動完成入庫。", 9.5, "64748B", align=PP_ALIGN.RIGHT)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

add_cover(prs.slides.add_slide(prs.slide_layouts[6]))
add_flow_slide(prs.slides.add_slide(prs.slide_layouts[6]))
add_example_slide(prs.slides.add_slide(prs.slide_layouts[6]))

prs.save(OUT)
print(f"saved {OUT}")
