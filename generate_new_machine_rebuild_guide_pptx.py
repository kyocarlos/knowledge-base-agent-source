from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("<project-root>/knowledge-base/new_machine_rebuild_guide.pptx")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0.0, radius=False):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def set_text(frame, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Sans CJK TC"):
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)


def textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, text, size, color, bold=bold, align=align)
    return box


def badge(slide, left, top, width, text, fill, color):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        Inches(0.30),
        fill=fill,
        line=fill,
        radius=True,
    )
    set_text(shape.text_frame, text, 8.5, color, align=PP_ALIGN.CENTER)
    return shape


def card(slide, left, top, width, height, title, body, accent, fill="FFFFFF"):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=fill,
        line="CBD5E1",
        radius=True,
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(7)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Noto Sans CJK TC"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = rgb(accent)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Noto Sans CJK TC"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = rgb("475569")
    return shape


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    for left, top, width, height, color in [
        (Inches(-0.35), Inches(-0.25), Inches(4.0), Inches(4.0), "DBEAFE"),
        (Inches(9.15), Inches(-0.35), Inches(4.5), Inches(3.9), "D1FAE5"),
        (Inches(8.6), Inches(4.8), Inches(4.8), Inches(2.8), "FEF3C7"),
    ]:
        deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        deco.fill.solid()
        deco.fill.fore_color.rgb = rgb(color)
        deco.fill.transparency = 0.80
        deco.line.fill.background()

    border = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.22),
        Inches(0.18),
        Inches(12.89),
        Inches(7.14),
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    border.fill.transparency = 1


def header(slide, title, subtitle, badge_text, badge_fill, badge_color):
    textbox(slide, Inches(0.52), Inches(0.32), Inches(3.6), Inches(0.20), "Knowledge Base / New Machine Rebuild", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.56), Inches(9.0), Inches(0.42), title, 24, "0F172A", bold=True)
    textbox(slide, Inches(0.52), Inches(0.98), Inches(9.8), Inches(0.32), subtitle, 12, "475569")
    badge(slide, Inches(10.72), Inches(0.50), Inches(1.80), badge_text, badge_fill, badge_color)


def add_cover(slide):
    add_background(slide)
    header(
        slide,
        "在另一台電腦重建 knowledge-base",
        "這份操作手冊說明如何從 code-only GitHub repo 開始，把系統在新電腦上重新架起來，然後再重新 ingest 資料。",
        "操作手冊",
        "EFF6FF",
        "1D4ED8",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(3.90), Inches(1.40), "這份手冊的目標", "你不需要把舊機器整包複製過去，只要 clone GitHub、補設定、起服務、再 ingest 文件即可。", "2563EB")
    card(slide, Inches(0.72), Inches(3.30), Inches(3.90), Inches(1.40), "最重要的原則", "GitHub 放可重建的程式碼，資料放獨立 bundle。這樣搬到新機器時才穩定。", "16A34A")
    card(slide, Inches(0.72), Inches(4.82), Inches(3.90), Inches(1.40), "你會得到什麼", "一套能在新電腦重新啟動的 knowledge-base 系統，再把需要的文件重新 ingest 進去。", "D97706")

    card(slide, Inches(5.05), Inches(1.95), Inches(7.05), Inches(2.10), "一句話版流程", "clone repository -> 安裝 Docker -> 建 symlink 或調整路徑 -> 複製 config -> 啟動服務 -> 還原資料 bundle（可選） -> 重新 ingest -> 驗證搜尋與 Neo4j/Qdrant。", "047857")

    flow = [
        (0.86, "clone repo", "GitHub code-only", "2563EB"),
        (2.60, "docker", "install engine", "16A34A"),
        (4.34, "path fix", "symlink or edit paths", "D97706"),
        (6.38, "config", "fill local secrets", "1D4ED8"),
        (8.42, "start stack", "Docker + restart script", "047857"),
        (10.34, "ingest", "rebuild your KB", "B45309"),
    ]
    for x, t, d, a in flow:
        card(slide, Inches(x), Inches(5.28), Inches(1.72), Inches(0.88), t, d, a)
    for x1, x2 in [(2.46, 2.60), (4.20, 4.34), (6.24, 6.38), (8.28, 8.42), (10.20, 10.34)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(5.72), Inches(x2), Inches(5.72))
        line.line.color.rgb = rgb("94A3B8")
        line.line.width = Pt(1.3)

    footer = slide.shapes.add_textbox(Inches(0.74), Inches(6.60), Inches(11.9), Inches(0.22))
    set_text(footer.text_frame, "建置完成後，系統會先可跑起來，再由你重新 ingest 新資料。", 9.5, "64748B", align=PP_ALIGN.RIGHT)


def add_prereq_slide(slide):
    add_background(slide)
    header(
        slide,
        "第一步：準備新電腦環境",
        "先安裝好系統需要的基礎工具，這一步完成後，再開始 clone 與設定。",
        "前置準備",
        "ECFDF5",
        "047857",
    )

    items = [
        ("1", "Git", "用來 clone GitHub repo。", "2563EB"),
        ("2", "Docker / Docker Compose", "用來啟動 Neo4j、Redis、Web API、Celery。", "16A34A"),
        ("3", "Python 3.12", "給後端腳本與工具使用。", "D97706"),
        ("4", "Node.js 18+", "給前端建置使用。", "1D4ED8"),
        ("5", "Ollama", "如果你使用本機模型，這是 LLM 服務。", "047857"),
    ]
    y = 1.72
    for idx, (num, title, body, accent) in enumerate(items):
        badge(slide, Inches(0.78), Inches(y + idx * 0.83), Inches(0.42), num, "EFF6FF", accent)
        card(slide, Inches(1.30), Inches(y - 0.04 + idx * 0.83), Inches(5.95), Inches(0.68), title, body, accent)

    card(slide, Inches(7.60), Inches(1.82), Inches(4.92), Inches(1.30), "建議先確認", "Docker 是否能正常啟動、`docker ps` 是否可執行、`python3.12 --version` 是否存在。", "2563EB")
    card(slide, Inches(7.60), Inches(3.35), Inches(4.92), Inches(1.30), "若要搬資料", "請先準備獨立資料 bundle，不要期待 GitHub repo 內會有 `data/`。", "16A34A")
    card(slide, Inches(7.60), Inches(4.88), Inches(4.92), Inches(1.30), "最常見卡點", "缺少 Docker、Node 或 Ollama；先補好環境，再繼續下一步。", "D97706")


def add_docker_slide(slide):
    add_background(slide)
    header(
        slide,
        "第二步：安裝 Docker 並啟動 Neo4j / Qdrant",
        "這個專案把 Neo4j 和 Qdrant 都放在 Docker 裡，先把 Docker 裝好，後面的服務才會起得來。",
        "docker db",
        "ECFDF5",
        "047857",
    )

    card(slide, Inches(0.72), Inches(1.72), Inches(5.95), Inches(1.20), "安裝 Docker", "Windows / macOS 建議用 Docker Desktop；Linux 可用 Docker Engine + Docker Compose plugin。", "2563EB")
    card(slide, Inches(0.72), Inches(3.06), Inches(5.95), Inches(1.20), "Linux 常用安裝指令", "sudo apt-get update && sudo apt-get install -y docker.io docker-compose-plugin && sudo systemctl enable --now docker", "16A34A")
    card(slide, Inches(0.72), Inches(4.40), Inches(5.95), Inches(1.20), "驗證 Docker", "確認 `docker version`、`docker compose version`、`docker ps` 都可以正常執行。", "D97706")

    card(slide, Inches(6.98), Inches(1.72), Inches(5.62), Inches(1.38), "Neo4j 在 Docker 裡怎麼跑", "由 `docker-compose.yml` 的 `neo4j` service 啟動，容器名是 `kb-neo4j`，主機對外是 `17474/17687`，容器內連線是 `bolt://neo4j:7687`。", "1D4ED8")
    card(slide, Inches(6.98), Inches(3.28), Inches(5.62), Inches(1.38), "Qdrant 在 Docker 裡怎麼跑", "由 `restart_kb.sh` 另外建立 `kb-qdrant` 容器，主機對外是 `6335/6336`；若容器不存在，腳本會自動 `docker run`。", "047857")
    card(slide, Inches(6.98), Inches(4.84), Inches(5.62), Inches(1.18), "快速啟動方式", "./restart_kb.sh 會一起拉起 Neo4j、Qdrant、Redis、Web API、Celery 與前端。", "B45309")


def add_clone_slide(slide):
    add_background(slide)
    header(
        slide,
        "第三步：Clone repo 並處理絕對路徑",
        "這個專案目前仍有一些硬編碼路徑，所以新電腦第一次建置時，先讓舊路徑能對上。",
        "clone / path",
        "EFF6FF",
        "1D4ED8",
    )

    card(slide, Inches(0.72), Inches(1.80), Inches(4.20), Inches(1.25), "Clone GitHub", "git clone git@github.com:kyocarlos/knowledge-base.git knowledge-base", "2563EB")
    card(slide, Inches(0.72), Inches(3.20), Inches(4.20), Inches(1.25), "切到專案目錄", "cd knowledge-base", "16A34A")
    card(slide, Inches(0.72), Inches(4.60), Inches(4.20), Inches(1.25), "建立相容 symlink", "sudo mkdir -p <project-root> && sudo ln -s \"$(pwd)\" <project-root>/knowledge-base", "D97706")

    card(slide, Inches(5.35), Inches(1.85), Inches(6.55), Inches(1.45), "為什麼要做這一步", "目前 `docker-compose.yml`、`restart_kb.sh`、`start.sh`、`config/config.yaml` 與部分 `src/` 還有原機器的絕對路徑。先建立 symlink，最快能讓現有腳本直接跑起來。", "1D4ED8")
    card(slide, Inches(5.35), Inches(3.45), Inches(6.55), Inches(1.45), "長期建議", "如果你之後要把它變成真正可移機版本，再把硬編碼路徑改成相對路徑或環境變數。", "047857")
    card(slide, Inches(5.35), Inches(5.05), Inches(6.55), Inches(1.15), "這一步做完後", "新電腦就能把這份專案當作原機器那個路徑來用。", "D97706")


def add_config_slide(slide):
    add_background(slide)
    header(
        slide,
        "第四步：複製設定檔與填入本機資訊",
        "GitHub 上只保留範本檔，你要在新電腦上建立自己的 `config/config.yaml`。",
        "config",
        "FEF3C7",
        "B45309",
    )

    card(slide, Inches(0.72), Inches(1.82), Inches(4.30), Inches(1.10), "複製範本", "cp config/config.yaml.example config/config.yaml", "2563EB")
    card(slide, Inches(0.72), Inches(3.05), Inches(4.30), Inches(1.10), "填入 Neo4j", "neo4j_uri: bolt://neo4j:7687", "16A34A")
    card(slide, Inches(0.72), Inches(4.28), Inches(4.30), Inches(1.10), "填入密碼", "neo4j_password: 你的密碼", "D97706")
    card(slide, Inches(0.72), Inches(5.51), Inches(4.30), Inches(1.10), "可調整的項目", "Ollama model、Qdrant URL、Redis、路徑與 ports。", "1D4ED8")

    card(slide, Inches(5.35), Inches(1.82), Inches(6.55), Inches(1.55), "設定原則", "把本機專屬的東西都留在 `config.yaml`，不要把它提交回 GitHub。GitHub 只留 `config.yaml.example`。", "2563EB")
    card(slide, Inches(5.35), Inches(3.55), Inches(6.55), Inches(1.55), "資料庫連線", "容器內 Neo4j 用 `bolt://neo4j:7687`；主機外部瀏覽器連線時才會用 `localhost:17687`。", "16A34A")
    card(slide, Inches(5.35), Inches(5.28), Inches(6.55), Inches(1.25), "如果你只想重建系統", "不用先恢復舊資料，先把服務起來最重要。", "047857")


def add_start_slide(slide):
    add_background(slide)
    header(
        slide,
        "第五步：啟動服務",
        "最簡單的方式是直接用專案提供的重啟腳本，讓前端、API、Neo4j、Redis、Celery 一起起來。",
        "start stack",
        "ECFDF5",
        "047857",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(5.00), Inches(1.05), "執行指令", "./restart_kb.sh", "2563EB")
    card(slide, Inches(0.72), Inches(3.02), Inches(5.00), Inches(1.05), "它會做什麼", "建置前端、啟動 Docker stack、檢查健康狀態、做 websocket smoke test。", "16A34A")
    card(slide, Inches(0.72), Inches(4.26), Inches(5.00), Inches(1.05), "如果失敗", "先檢查 symlink、Docker、Node、Ollama、以及本機 port 有沒有衝突。", "D97706")

    card(slide, Inches(5.98), Inches(1.90), Inches(6.00), Inches(1.30), "成功後你應該看到", "kb-web、kb-neo4j、kb-redis、kb-celery-search、kb-celery-ingest、kb-celery-beat 都處於 Up 狀態。", "1D4ED8")
    card(slide, Inches(5.98), Inches(3.45), Inches(6.00), Inches(1.30), "驗證頁面", "http://localhost:8000/health、http://localhost:3030/chat.html、Neo4j Browser。", "047857")
    card(slide, Inches(5.98), Inches(5.00), Inches(6.00), Inches(1.05), "補充", "如果你只想看 Neo4j 資料，請連到 `bolt://localhost:17687`。", "2563EB")


def add_data_slide(slide):
    add_background(slide)
    header(
        slide,
        "第六步：還原資料並重新 ingest",
        "如果你要沿用舊資料，就把資料包解壓；如果要乾淨重建，就直接重新 ingest。",
        "data / ingest",
        "FEF3C7",
        "B45309",
    )

    card(slide, Inches(0.72), Inches(1.80), Inches(4.25), Inches(1.20), "資料包來源", "獨立的 tar.gz bundle，不放在 GitHub repo 內。", "2563EB")
    card(slide, Inches(0.72), Inches(3.15), Inches(4.25), Inches(1.20), "還原方式", "tar -xzf backups/knowledge-base-data-YYYYMMDD_HHMMSS.tar.gz -C /path/to/knowledge-base", "16A34A")
    card(slide, Inches(0.72), Inches(4.50), Inches(4.25), Inches(1.20), "重新 ingest", "把新文件放進 data/raw/ 或透過上傳頁，系統會自動轉換並入庫。", "D97706")

    card(slide, Inches(5.35), Inches(1.85), Inches(6.55), Inches(1.40), "資料夾會長什麼樣子", "通常會恢復 `data/raw/`、`data/processed/`、`data/assets/`、`data/uploads/`。", "1D4ED8")
    card(slide, Inches(5.35), Inches(3.45), Inches(6.55), Inches(1.40), "如果不要舊資料", "可以直接跳過還原，從新文件開始 ingest，等於做一個乾淨的知識庫。", "047857")
    card(slide, Inches(5.35), Inches(5.05), Inches(6.55), Inches(1.10), "最終效果", "文件會進 Neo4j 與 Qdrant，之後可以在聊天頁面直接查詢。", "2563EB")


def add_verify_slide(slide):
    add_background(slide)
    header(
        slide,
        "第七步：驗證系統是否可用",
        "把最基本的連線、搜尋、和資料庫查詢都跑一遍，確認真的建成功。",
        "verify",
        "EFF6FF",
        "1D4ED8",
    )

    checks = [
        ("1", "Health check", "curl -fsS http://localhost:8000/health", "2563EB"),
        ("2", "Chat page", "http://localhost:3030/chat.html", "16A34A"),
        ("3", "Neo4j Browser", "bolt://localhost:17687", "D97706"),
        ("4", "Query Document", "MATCH (d:Document) RETURN count(d);", "1D4ED8"),
        ("5", "Search sample", "在聊天頁查一個你熟悉的報告名稱", "047857"),
    ]
    y = 1.72
    for idx, (num, title, body, accent) in enumerate(checks):
        badge(slide, Inches(0.78), Inches(y + idx * 0.82), Inches(0.42), num, "EFF6FF", accent)
        card(slide, Inches(1.30), Inches(y - 0.04 + idx * 0.82), Inches(5.80), Inches(0.68), title, body, accent)

    card(slide, Inches(7.55), Inches(1.82), Inches(4.95), Inches(1.30), "如果看到 No instance connected", "代表 Neo4j Browser 還沒連上資料庫，請手動輸入 `bolt://localhost:17687` 與帳密。", "D97706")
    card(slide, Inches(7.55), Inches(3.35), Inches(4.95), Inches(1.30), "如果能查到 Document", "表示 Neo4j 已正常承載報告資料。", "16A34A")
    card(slide, Inches(7.55), Inches(4.88), Inches(4.95), Inches(1.30), "如果聊天能回應並帶來源", "表示知識庫查詢鏈路已連通。", "2563EB")


def add_troubleshoot_slide(slide):
    add_background(slide)
    header(
        slide,
        "常見問題與排錯",
        "新機器最常遇到的是路徑、port、模型與瀏覽器連線方式。",
        "troubleshoot",
        "FEE2E2",
        "B91C1C",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(3.95), Inches(1.18), "看不到 Neo4j 資料", "Browser 內連到 `bolt://localhost:17687`，不是預設的 7687。", "B91C1C")
    card(slide, Inches(0.72), Inches(3.10), Inches(3.95), Inches(1.18), "restart_kb.sh 找不到路徑", "先建立 `<project-root>/knowledge-base` 的 symlink，或把硬編碼路徑改掉。", "D97706")
    card(slide, Inches(0.72), Inches(4.42), Inches(3.95), Inches(1.18), "Ollama 模型不存在", "先 `ollama pull gemma4:12b`，再重新啟動。", "16A34A")

    card(slide, Inches(4.95), Inches(1.92), Inches(7.55), Inches(1.35), "Docker port 衝突", "如果 17474 / 17687 / 3030 / 8000 / 6335 被別的服務佔用，先停掉衝突服務，再重新啟動。", "2563EB")
    card(slide, Inches(4.95), Inches(3.48), Inches(7.55), Inches(1.35), "資料沒進去", "先確認資料真的放進 `data/raw/` 或上傳頁，然後再看 ingest log 是否有完成。", "047857")
    card(slide, Inches(4.95), Inches(5.04), Inches(7.55), Inches(1.35), "建議的排錯順序", "路徑 -> Docker -> 模型 -> Neo4j / Qdrant -> 搜尋頁。", "1D4ED8")


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = [
    add_cover,
    add_prereq_slide,
    add_docker_slide,
    add_clone_slide,
    add_config_slide,
    add_start_slide,
    add_data_slide,
    add_verify_slide,
    add_troubleshoot_slide,
]

for builder in slides:
    builder(prs.slides.add_slide(prs.slide_layouts[6]))

prs.save(OUT)
print(f"saved {OUT}")
