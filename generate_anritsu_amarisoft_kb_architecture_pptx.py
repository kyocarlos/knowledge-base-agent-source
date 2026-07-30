from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("<project-root>/knowledge-base/anritsu_amarisoft_kb_architecture.pptx")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
FONT = "Noto Sans CJK TC"


def color(value):
    return RGBColor.from_string(value.replace("#", "").upper())


def box(slide, x, y, w, h, text, fill="FFFFFF", line="CBD5E1", size=12, bold=False, fg="0F172A"):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color(fill)
    s.line.color.rgb = color(line); s.line.width = Pt(1)
    tf = s.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(5)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color(fg)
    return s


def text(slide, x, y, w, h, value, size=12, fg="475569", bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = value; r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color(fg)
    return t


def line(slide, x1, y1, x2, y2, fg="94A3B8", width=2):
    l = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    l.line.color.rgb = color(fg); l.line.width = Pt(width)
    return l


def base(title, subtitle, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color("F8FAFC"); bg.line.fill.background()
    text(s, .55, .30, 10.8, .42, title, 24, "0F172A", True)
    text(s, .58, .79, 11.6, .30, subtitle, 11, "64748B")
    text(s, 12.45, .32, .35, .25, f"{n:02d}", 10, "2563EB", True, PP_ALIGN.RIGHT)
    return s


def slide1():
    s = base("Anritsu／Amarisoft 測試報告整合架構", "兩個測試環境獨立執行，報告經 HTTPS 傳送到 Knowledge Base，人工核准後正式入庫。", 1)
    box(s, .65, 1.45, 2.25, .52, "Anritsu 測試環境", "DBEAFE", "60A5FA", 14, True, "1E3A8A")
    box(s, .65, 2.20, 2.25, .80, "儀器／DUT\n測試控制 Agent", "FFFFFF", "60A5FA")
    box(s, .65, 3.35, 2.25, .80, "產生標準 Excel\n本機驗證與 Outbox", "FFFFFF", "60A5FA")
    line(s, 1.78, 3.00, 1.78, 3.35, "60A5FA")
    box(s, 10.42, 1.45, 2.25, .52, "Amarisoft 測試環境", "FEF3C7", "F59E0B", 14, True, "92400E")
    box(s, 10.42, 2.20, 2.25, .80, "儀器／DUT\n測試控制 Agent", "FFFFFF", "F59E0B")
    box(s, 10.42, 3.35, 2.25, .80, "產生標準 Excel\n本機驗證與 Outbox", "FFFFFF", "F59E0B")
    line(s, 11.55, 3.00, 11.55, 3.35, "F59E0B")
    box(s, 4.45, 1.45, 4.45, .52, "Knowledge Base 平台", "DCFCE7", "22C55E", 14, True, "166534")
    box(s, 3.70, 2.35, 2.35, .82, "HTTPS API\nToken／Hash／run_id 驗證", "FFFFFF", "22C55E")
    box(s, 6.05, 2.35, 2.35, .82, "Staging + Submission Registry\n待審狀態與稽核", "FFFFFF", "22C55E")
    box(s, 4.88, 3.75, 2.35, .82, "KB 待審管理頁\n核准／退回", "FFF7ED", "F97316", 13, True, "9A3412")
    box(s, 7.25, 3.75, 2.35, .82, "Redis／Celery\n核准後攝入", "FFFFFF", "22C55E")
    box(s, 4.20, 5.25, 2.35, .82, "Neo4j\n結構化測試數值", "ECFDF5", "10B981", 13, True, "065F46")
    box(s, 6.75, 5.25, 2.35, .82, "Qdrant\n報告內容與向量", "EFF6FF", "60A5FA", 13, True, "1E40AF")
    line(s, 2.90, 3.75, 3.70, 2.75, "64748B"); line(s, 10.42, 3.75, 8.40, 2.75, "64748B")
    line(s, 6.05, 3.17, 6.05, 3.75, "64748B"); line(s, 7.23, 4.57, 7.23, 5.25, "64748B"); line(s, 6.05, 4.57, 5.38, 5.25, "64748B")
    text(s, 3.25, 6.52, 6.8, .34, "核心原則：測試環境不直連資料庫；KB 統一驗證、審核、攝入與索引。", 14, "0F172A", True, PP_ALIGN.CENTER)


def slide2():
    s = base("Test Report 從產生到入庫的生命週期", "人工核准位於正式 Neo4j／Qdrant 寫入之前，確保可追蹤、可退回、可重送。", 2)
    states = [("received", "收到報告", "E0F2FE"), ("validating", "格式／欄位驗證", "E0F2FE"), ("pending_review", "待人工審核", "FEF3C7"), ("approved", "核准", "DCFCE7"), ("queued", "排入攝入佇列", "DCFCE7"), ("writing", "Neo4j + Qdrant", "DCFCE7"), ("completed", "可搜尋／可比較", "BBF7D0")]
    x = .45
    for i, (code, label, fill) in enumerate(states):
        box(s, x, 2.35, 1.66, .82, f"{code}\n{label}", fill, "94A3B8", 10 if len(code) > 10 else 11, True)
        if i < len(states)-1: line(s, x+1.66, 2.76, x+1.86, 2.76, "64748B", 2)
        x += 1.86
    box(s, 2.05, 4.15, 2.50, .88, "validation_failed\n缺欄位／格式錯誤", "FEE2E2", "F87171", 12, True, "991B1B")
    box(s, 5.40, 4.15, 2.50, .88, "rejected\n審核退回＋原因", "FEE2E2", "F87171", 12, True, "991B1B")
    box(s, 8.75, 4.15, 2.50, .88, "ingest_failed\n可重試／Reconciliation", "FEE2E2", "F87171", 12, True, "991B1B")
    line(s, 1.28, 3.17, 3.30, 4.15, "F87171", 1); line(s, 3.78, 3.17, 6.65, 4.15, "F87171", 1); line(s, 9.02, 3.17, 10.00, 4.15, "F87171", 1)
    text(s, .8, 5.72, 11.8, .62, "相同 environment + run_id + SHA-256 重送時回傳既有任務；相同 run_id 但內容不同則拒絕並回報衝突。", 15, "0F172A", True, PP_ALIGN.CENTER)


def slide3():
    s = base("標準報告契約與資料落點", "共同 Excel schema 讓兩套環境能使用同一套 parser，並支援精確指標比較與原文引用。", 3)
    box(s, .65, 1.55, 3.15, .55, "標準 Excel Report v1", "DBEAFE", "60A5FA", 14, True, "1E3A8A")
    box(s, .65, 2.30, 3.15, 2.58, "Manifest：run_id／environment／DUT／版本\nRadioConfig：RAT／band／頻寬／MIMO\nTestCases：case ID／條件／狀態\nMeasurements：數值／單位／門檻\nVerdicts：Pass／Fail／原因\nRawArtifacts：log／截圖／checksum", "FFFFFF", "60A5FA", 12)
    box(s, 4.70, 1.55, 3.15, .55, "Neo4j 結構化圖譜", "DCFCE7", "22C55E", 14, True, "166534")
    box(s, 4.70, 2.30, 3.15, 2.58, "TestEnvironment → TestRun\nTestRun → DUT／RadioConfig\nTestRun → TestCase\nTestCase → Measurement／Verdict\n保存數值、單位、門檻、判定\n用途：精確比較、趨勢、Fail 分析", "FFFFFF", "22C55E", 12)
    box(s, 8.75, 1.55, 3.15, .55, "Qdrant 向量索引", "E0F2FE", "60A5FA", 14, True, "1E40AF")
    box(s, 8.75, 2.30, 3.15, 2.58, "報告 Markdown 與 source chunks\nrun_id／environment\nproject／DUT／band\nprotocol／direction／verdict\n原始報告名稱與 chunk_index\n用途：語意搜尋、來源引用、Chat", "FFFFFF", "60A5FA", 12)
    line(s, 3.80, 3.60, 4.70, 3.60, "64748B"); line(s, 7.85, 3.60, 8.75, 3.60, "64748B")
    text(s, 1.0, 5.65, 11.3, .75, "共同欄位採資料驅動；Anritsu／Amarisoft 特有欄位放入可擴充 key/value 區，不在程式中硬編碼環境分支。", 16, "0F172A", True, PP_ALIGN.CENTER)


def slide4():
    s = base("導入順序與驗收重點", "先完成共同契約與單一環境試行，再接入第二套環境，降低跨環境 schema 漂移風險。", 4)
    steps = [("1", "凍結 Schema v1", "建立 Excel template、fixture 與欄位驗證器"), ("2", "建立上傳與待審", "Agent token、staging、submission registry、審核頁"), ("3", "接入正式攝入", "Canonical parser → Celery → Neo4j／Qdrant"), ("4", "Anritsu Pilot", "完成單環境端到端與失敗重送驗證"), ("5", "Amarisoft + 雙環境驗收", "比較查詢、趨勢、Fail 分析與權限測試")]
    y = 1.45
    for no, title, body in steps:
        box(s, .70, y, .58, .58, no, "2563EB", "2563EB", 18, True, "FFFFFF")
        box(s, 1.50, y-.02, 3.15, .62, title, "FFFFFF", "CBD5E1", 14, True, "0F172A")
        text(s, 5.00, y+.08, 7.35, .42, body, 13, "475569")
        if no != "5": line(s, .99, y+.58, .99, y+1.03, "94A3B8", 2)
        y += 1.05
    box(s, .70, 6.72, 11.70, .42, "驗收入口：使用 https://127.0.0.1:3030/chat.html 實測指定 run、Fail 案例、Anritsu／Amarisoft 比較與來源引用。", "ECFDF5", "10B981", 12, True, "065F46")


slide1(); slide2(); slide3(); slide4()
prs.save(OUT)
print(OUT)
