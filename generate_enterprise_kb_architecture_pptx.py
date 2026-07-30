from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "knowledge_base_enterprise_architecture.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Noto Sans CJK TC"
FONT_DISPLAY = "Noto Serif CJK TC"

COLORS = {
    "navy": "0B1F33",
    "navy_2": "102A43",
    "ink": "162B3D",
    "blue": "2F6FED",
    "blue_2": "5B8DEF",
    "blue_soft": "EAF1FF",
    "teal": "00A6A6",
    "teal_dark": "007F7F",
    "teal_soft": "E3F7F6",
    "amber": "F4B740",
    "amber_dark": "A96B00",
    "amber_soft": "FFF4D6",
    "red": "D94F4F",
    "red_soft": "FDECEC",
    "green": "2E9D68",
    "green_soft": "E6F6EE",
    "slate": "52677A",
    "muted": "71859A",
    "line": "D5DEE8",
    "panel": "F5F8FB",
    "white": "FFFFFF",
    "off_white": "F8FAFC",
    "dark_panel": "132E46",
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    fill,
    line=None,
    radius=True,
    transparency=0,
    line_width=1,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    value,
    size=12,
    color=None,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    font=None,
    margin=0,
    fit=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if fit else MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = Pt(margin)
    frame.margin_right = Pt(margin)
    frame.margin_top = Pt(margin)
    frame.margin_bottom = Pt(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = value
    run.font.name = font or FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color or COLORS["ink"])
    return box


def add_paragraphs(slide, x, y, w, h, paragraphs, fill=None, line=None, margin=9, radius=True):
    if fill:
        shape = add_rect(slide, x, y, w, h, fill, line=line, radius=radius)
        frame = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.margin_left = Pt(margin)
    frame.margin_right = Pt(margin)
    frame.margin_top = Pt(margin)
    frame.margin_bottom = Pt(margin)
    for index, item in enumerate(paragraphs):
        if index == 0:
            paragraph = frame.paragraphs[0]
        else:
            paragraph = frame.add_paragraph()
        if isinstance(item, str):
            item = {"text": item}
        paragraph.alignment = item.get("align", PP_ALIGN.LEFT)
        paragraph.space_before = Pt(item.get("space_before", 0))
        paragraph.space_after = Pt(item.get("space_after", 0))
        paragraph.line_spacing = item.get("line_spacing", 1.0)
        paragraph.level = item.get("level", 0)
        run = paragraph.add_run()
        run.text = item.get("text", "")
        run.font.name = item.get("font", FONT)
        run.font.size = Pt(item.get("size", 11))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = rgb(item.get("color", COLORS["ink"]))
    return shape


def add_pill(slide, x, y, w, text, fill, color, h=0.28, size=8.5, line=None):
    shape = add_rect(slide, x, y, w, h, fill, line=line or fill, radius=True)
    shape.text_frame.clear()
    shape.text_frame.word_wrap = False
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    shape.text_frame.margin_left = Pt(4)
    shape.text_frame.margin_right = Pt(4)
    shape.text_frame.margin_top = Pt(0)
    shape.text_frame.margin_bottom = Pt(0)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = rgb(color)
    return shape


def add_card(
    slide,
    x,
    y,
    w,
    h,
    title,
    body="",
    accent=None,
    fill=None,
    line=None,
    title_size=12,
    body_size=9.5,
    icon=None,
    title_color=None,
    body_color=None,
    radius=True,
):
    accent = accent or COLORS["blue"]
    fill = fill or COLORS["white"]
    line = line or COLORS["line"]
    add_rect(slide, x, y, w, h, fill, line=line, radius=radius)
    add_rect(slide, x, y, 0.07, h, accent, radius=False)

    text_x = x + 0.18
    text_w = w - 0.30
    if icon:
        add_rect(slide, x + 0.18, y + 0.16, 0.42, 0.42, accent, radius=True)
        add_text(
            slide,
            x + 0.18,
            y + 0.16,
            0.42,
            0.42,
            icon,
            size=8.5,
            color=COLORS["white"],
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_VERTICAL_ANCHOR.MIDDLE,
        )
        text_x = x + 0.72
        text_w = w - 0.86
    add_text(
        slide,
        text_x,
        y + 0.13,
        text_w,
        0.34,
        title,
        size=title_size,
        color=title_color or accent,
        bold=True,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
        fit=True,
    )
    if body:
        body_y = y + (0.62 if icon else 0.52)
        body_h = h - (0.70 if icon else 0.61)
        add_text(
            slide,
            x + 0.18,
            body_y,
            w - 0.34,
            body_h,
            body,
            size=body_size,
            color=body_color or COLORS["slate"],
            valign=MSO_VERTICAL_ANCHOR.TOP,
            fit=True,
        )


def add_flow_step(slide, x, y, w, h, number, title, body, accent, fill=None):
    add_rect(slide, x, y, w, h, fill or COLORS["white"], line=COLORS["line"], radius=True)
    add_rect(slide, x + 0.12, y + 0.13, 0.34, 0.34, accent, radius=True)
    add_text(
        slide,
        x + 0.12,
        y + 0.13,
        0.34,
        0.34,
        str(number),
        size=8.5,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_text(slide, x + 0.54, y + 0.12, w - 0.65, 0.38, title, 10.5, accent, True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
    add_text(slide, x + 0.15, y + 0.53, w - 0.30, h - 0.62, body, 8.6, COLORS["slate"], valign=MSO_VERTICAL_ANCHOR.TOP, fit=True)


def add_chevron(slide, x, y, w=0.24, h=0.34, fill=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill or COLORS["muted"])
    shape.line.fill.background()
    return shape


def add_connector(slide, x1, y1, x2, y2, color=None, width=1.6, elbow=False, dash=False):
    connector_type = MSO_CONNECTOR.ELBOW if elbow else MSO_CONNECTOR.STRAIGHT
    line = slide.shapes.add_connector(
        connector_type,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color or COLORS["muted"])
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = 2
    return line


def add_background(slide, dark=False):
    fill = COLORS["navy"] if dark else COLORS["off_white"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill, radius=False)
    if dark:
        for x, y, w, h, color, transparency in [
            (8.6, -0.6, 5.5, 5.5, COLORS["blue"], 82),
            (10.3, 4.6, 4.0, 4.0, COLORS["teal"], 84),
            (-1.2, 5.7, 4.0, 2.0, COLORS["blue_2"], 90),
        ]:
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(x),
                Inches(y),
                Inches(w),
                Inches(h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(color)
            shape.fill.transparency = transparency
            shape.line.fill.background()
    else:
        add_rect(slide, 0, 0, 0.18, SLIDE_H, COLORS["blue"], radius=False)
        add_rect(slide, 0.18, 0, SLIDE_W - 0.18, 0.08, COLORS["navy"], radius=False)


def add_header(slide, section, title, subtitle, page):
    add_text(slide, 0.52, 0.28, 5.8, 0.18, section.upper(), 8.8, COLORS["blue"], True)
    add_text(slide, 0.52, 0.52, 10.8, 0.46, title, 23, COLORS["navy"], True, font=FONT_DISPLAY, fit=True)
    add_text(slide, 0.53, 1.00, 11.3, 0.28, subtitle, 10.5, COLORS["slate"], fit=True)
    add_pill(slide, 11.93, 0.36, 0.78, f"{page:02d}", COLORS["navy"], COLORS["white"], h=0.34, size=9.5)


def add_footer(slide, takeaway, dark=False):
    y = 7.08
    line_color = "33506B" if dark else COLORS["line"]
    add_rect(slide, 0.52, y - 0.09, 12.18, 0.01, line_color, radius=False)
    add_text(slide, 0.52, y, 9.0, 0.19, f"報告重點｜{takeaway}", 8.2, "C4D3E0" if dark else COLORS["slate"], bold=True, fit=True)
    add_text(slide, 9.72, y, 2.98, 0.19, "Current-state architecture · 2026.07.17", 7.0, "8FA6BA" if dark else COLORS["muted"], align=PP_ALIGN.RIGHT, fit=True)


def add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, dark=True)

    add_pill(slide, 0.82, 0.58, 2.45, "ENTERPRISE ARCHITECTURE", COLORS["blue"], COLORS["white"], h=0.34, size=9)
    add_text(slide, 0.82, 1.42, 7.25, 1.16, "Knowledge Base\n企業知識庫系統架構", 31, COLORS["white"], True, font=FONT_DISPLAY, fit=True)
    add_text(
        slide,
        0.84,
        2.86,
        6.85,
        0.74,
        "從文件攝入、混合檢索到 OpenClaw AI 助理的完整現況架構\n主管報告版｜涵蓋應用、資料、部署、維運與治理",
        14,
        "C7D8E7",
        fit=True,
    )

    add_rect(slide, 0.82, 4.18, 6.75, 1.12, COLORS["dark_panel"], line="2B4A65", radius=True, transparency=5)
    add_text(slide, 1.05, 4.42, 1.35, 0.18, "ARCHITECTURE THESIS", 8.3, COLORS["teal"], True)
    add_text(
        slide,
        1.05,
        4.70,
        6.05,
        0.37,
        "以 FastAPI + Celery 解耦互動與批次工作，並以 Qdrant + Neo4j 支援語意與關聯雙引擎檢索。",
        12,
        COLORS["white"],
        True,
        fit=True,
    )

    add_text(slide, 8.38, 0.72, 3.8, 0.22, "PLATFORM LAYERS", 9, "AFC5D8", True)
    layers = [
        ("01", "Experience", "搜尋、聊天、上傳、管理", COLORS["blue"]),
        ("02", "API & Gateway", "Nginx、FastAPI、WebSocket", COLORS["blue_2"]),
        ("03", "Orchestration", "Redis、Celery、Beat", COLORS["teal"]),
        ("04", "Intelligence", "SearchEngine、Ollama、OpenClaw", COLORS["amber"]),
        ("05", "Data", "Qdrant、Neo4j、File Store", COLORS["green"]),
    ]
    y = 1.12
    for number, title, body, accent in layers:
        add_rect(slide, 8.35, y, 4.05, 0.82, COLORS["dark_panel"], line="2B4A65", radius=True)
        add_pill(slide, 8.55, y + 0.20, 0.48, number, accent, COLORS["white"], h=0.32, size=8)
        add_text(slide, 9.18, y + 0.12, 2.75, 0.28, title, 11.5, COLORS["white"], True, fit=True)
        add_text(slide, 9.18, y + 0.43, 2.84, 0.22, body, 8.5, "B8CAD9", fit=True)
        y += 0.94

    add_pill(slide, 0.84, 6.27, 1.34, "CURRENT STATE", COLORS["teal"], COLORS["white"], h=0.32, size=8.4)
    add_text(slide, 2.38, 6.27, 4.7, 0.32, "範圍：原始部署 + 可攜式 on-prem release", 9.2, "B8CAD9", valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_pill(slide, 10.78, 6.27, 1.62, "CONFIDENTIAL", COLORS["amber"], COLORS["navy"], h=0.32, size=8.4)
    add_footer(slide, "本簡報呈現已落地的系統能力，目標架構與改善項目另行標示。", dark=True)


def add_executive_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "01 · Executive Overview",
        "一張圖看懂：從企業文件到可追溯回答",
        "系統把多格式文件轉為可檢索知識，經由雙引擎召回與 AI 助理，交付具來源脈絡的答案與管理能力。",
        2,
    )

    add_text(slide, 0.55, 1.52, 1.85, 0.24, "服務對象", 10.5, COLORS["navy"], True)
    personas = [
        ("知識使用者", "搜尋、問答、來源追溯", "KU", COLORS["blue"]),
        ("資料維護者", "上傳、攝入、Chunk 修訂", "DM", COLORS["teal"]),
        ("平台管理者", "統計、日誌、排程、Skills", "PA", COLORS["amber"]),
    ]
    y = 1.88
    for title, body, icon, accent in personas:
        add_card(slide, 0.55, y, 2.05, 0.93, title, body, accent=accent, icon=icon, title_size=10.5, body_size=8.2)
        y += 1.08

    add_text(slide, 2.92, 1.52, 7.48, 0.24, "企業知識價值鏈", 10.5, COLORS["navy"], True)
    chain = [
        ("CAPTURE", "文件匯入", "PDF / Office / 圖片", COLORS["blue"]),
        ("UNDERSTAND", "解析與分塊", "Markdown / OCR / Assets", COLORS["blue_2"]),
        ("INDEX", "建立知識", "向量 + 圖譜", COLORS["teal"]),
        ("RETRIEVE", "混合檢索", "語意 + 關聯 + 規則", COLORS["green"]),
        ("ANSWER", "生成與引用", "Ollama / OpenClaw", COLORS["amber"]),
    ]
    x = 2.92
    card_w = 1.28
    for idx, (tag, title, body, accent) in enumerate(chain):
        add_rect(slide, x, 2.03, card_w, 2.54, COLORS["white"], line=COLORS["line"], radius=True)
        add_rect(slide, x, 2.03, card_w, 0.11, accent, radius=False)
        add_pill(slide, x + 0.12, 2.30, card_w - 0.24, tag, COLORS["panel"], accent, h=0.27, size=7.2, line=COLORS["line"])
        add_text(slide, x + 0.10, 2.80, card_w - 0.20, 0.52, title, 12, COLORS["navy"], True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
        add_text(slide, x + 0.10, 3.47, card_w - 0.20, 0.56, body, 8.5, COLORS["slate"], align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
        if idx < len(chain) - 1:
            add_chevron(slide, x + card_w + 0.08, 3.06, 0.22, 0.38, COLORS["muted"])
        x += 1.48

    add_text(slide, 10.62, 1.52, 2.16, 0.24, "管理輸出", 10.5, COLORS["navy"], True)
    outcomes = [
        ("可追溯回答", "來源文件、章節、引用占比", COLORS["blue"]),
        ("可操作資料", "文件版本、Chunk、資產", COLORS["teal"]),
        ("可觀測服務", "任務狀態、日誌、健康檢查", COLORS["amber"]),
    ]
    y = 1.88
    for title, body, accent in outcomes:
        add_card(slide, 10.62, y, 2.16, 0.93, title, body, accent=accent, title_size=10.5, body_size=8.2)
        y += 1.08

    add_rect(slide, 0.55, 5.25, 12.23, 1.25, COLORS["navy"], radius=True)
    add_text(slide, 0.78, 5.47, 1.32, 0.18, "CAPABILITY MAP", 8.2, COLORS["teal"], True)
    capabilities = [
        ("企業搜尋", "Basic / Deep / Vector / Hybrid"),
        ("AI 聊天", "KB Context + OpenClaw"),
        ("文件攝入", "Upload / Watch / Auto ingest"),
        ("管理中心", "Stats / Logs / Schedule"),
        ("Chunk 治理", "檢視 / 編輯 / 版本回復"),
        ("Skills 管理", "Workspace / skill files"),
    ]
    x = 2.20
    for title, body in capabilities:
        add_rect(slide, x, 5.43, 1.62, 0.72, COLORS["dark_panel"], line="31516D", radius=True)
        add_text(slide, x + 0.10, 5.51, 1.42, 0.22, title, 9.2, COLORS["white"], True, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, x + 0.08, 5.79, 1.46, 0.20, body, 7.2, "B9CAD9", align=PP_ALIGN.CENTER, fit=True)
        x += 1.72
    add_footer(slide, "平台價值不只在回答問題，而是把資料攝入、檢索、引用與維運串成一條可治理的知識供應鏈。")


def add_logical_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "02 · Logical Architecture",
        "Knowledge Base 完整邏輯架構",
        "五層架構將使用者體驗、同步 API、非同步任務、AI/檢索與持久化資料分離，降低互相耦合。",
        3,
    )

    rows = [
        {
            "label": "EXPERIENCE",
            "name": "體驗層",
            "accent": COLORS["blue"],
            "fill": COLORS["blue_soft"],
            "items": [
                ("Vue SPA", "搜尋 / 上傳 / 管理", "UI"),
                ("chat.html / ChatView", "KB 小幫手與引用", "CH"),
                ("Admin / Chunk / Skills", "治理與維運入口", "AD"),
            ],
        },
        {
            "label": "ACCESS",
            "name": "存取層",
            "accent": COLORS["blue_2"],
            "fill": "EEF4FF",
            "items": [
                ("Nginx TLS Gateway", "HTTPS 靜態頁與反向代理", "NG"),
                ("FastAPI / Uvicorn", "REST API、Pydantic、lifespan", "FA"),
                ("WebSocket Proxy", "Session、排隊、OpenClaw 轉送", "WS"),
            ],
        },
        {
            "label": "ORCHESTRATION",
            "name": "編排層",
            "accent": COLORS["teal"],
            "fill": COLORS["teal_soft"],
            "items": [
                ("Redis", "Broker / Result / Cache / Locks", "RD"),
                ("Celery Search", "搜尋與回答任務", "CS"),
                ("Celery Ingest", "轉檔與知識寫入", "CI"),
                ("Celery Beat", "Watch folder 定時掃描", "CB"),
            ],
        },
        {
            "label": "INTELLIGENCE",
            "name": "智能層",
            "accent": COLORS["amber_dark"],
            "fill": COLORS["amber_soft"],
            "items": [
                ("SearchEngine", "Basic / Deep / Vector / Hybrid", "SE"),
                ("Document Pipeline", "MarkItDown / Chunk / Assets", "DP"),
                ("Embedding", "BAAI bge-base-zh · 768D", "EM"),
                ("Ollama + OpenClaw", "推論、工具與串流回答", "AI"),
            ],
        },
        {
            "label": "DATA",
            "name": "資料層",
            "accent": COLORS["green"],
            "fill": COLORS["green_soft"],
            "items": [
                ("Qdrant", "knowledge_base 向量集合", "QD"),
                ("Neo4j", "文件 / 實體 / 報告圖譜", "N4"),
                ("File Store", "raw / watch / processed / assets", "FS"),
                ("OpenClaw Runtime", "Identity / Workspace / Skills", "OC"),
            ],
        },
    ]
    y = 1.45
    for row_index, row in enumerate(rows):
        add_rect(slide, 0.50, y, 12.25, 0.96, row["fill"], line=COLORS["line"], radius=True)
        add_rect(slide, 0.50, y, 1.34, 0.96, row["accent"], radius=True)
        add_text(slide, 0.63, y + 0.18, 1.08, 0.20, row["label"], 7.2, COLORS["white"], True, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, 0.63, y + 0.49, 1.08, 0.22, row["name"], 10.2, COLORS["white"], True, align=PP_ALIGN.CENTER, fit=True)

        item_count = len(row["items"])
        available = 10.55
        gap = 0.14
        item_w = (available - (item_count - 1) * gap) / item_count
        x = 2.02
        for title, body, icon in row["items"]:
            add_rect(slide, x, y + 0.13, item_w, 0.70, COLORS["white"], line=COLORS["line"], radius=True)
            add_rect(slide, x + 0.10, y + 0.25, 0.38, 0.38, row["accent"], radius=True)
            add_text(slide, x + 0.10, y + 0.25, 0.38, 0.38, icon, 7.6, COLORS["white"], True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
            add_text(slide, x + 0.56, y + 0.18, item_w - 0.65, 0.24, title, 9.5, COLORS["navy"], True, fit=True)
            add_text(slide, x + 0.56, y + 0.48, item_w - 0.65, 0.18, body, 7.4, COLORS["slate"], fit=True)
            x += item_w + gap
        if row_index < len(rows) - 1:
            add_connector(slide, 12.37, y + 0.96, 12.37, y + 1.08, row["accent"], width=2.0)
            add_chevron(slide, 12.22, y + 0.98, 0.30, 0.19, row["accent"])
        y += 1.06

    add_pill(slide, 0.53, 6.70, 1.52, "PRIMARY FLOW", COLORS["navy"], COLORS["white"], h=0.26, size=7.5)
    add_text(slide, 2.20, 6.70, 9.78, 0.22, "HTTPS / WSS → FastAPI → Redis / Celery → Search & AI → Qdrant / Neo4j / File Store", 8.4, COLORS["slate"], True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
    add_footer(slide, "同步入口保持輕量，耗時搜尋與攝入交由背景 worker；資料庫與 LLM 則透過明確介面被各流程共用。")


def add_query_chat_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "03 · Runtime Flow",
        "查詢與 AI 聊天：兩條互補的執行路徑",
        "搜尋 API 以 Celery 任務交付完整 KB 結果；聊天路徑則經 WebSocket 串接 OpenClaw，並在前端補入 KB 脈絡與引用。",
        4,
    )

    tracks = [
        {
            "y": 1.48,
            "h": 2.24,
            "label": "KB SEARCH",
            "title": "非同步知識庫搜尋",
            "accent": COLORS["blue"],
            "fill": COLORS["blue_soft"],
            "steps": [
                ("提出問題", "SearchView / API client"),
                ("POST /search", "FastAPI 驗證並建立 task"),
                ("Redis + Celery", "search queue / task state"),
                ("雙引擎召回", "Qdrant + Neo4j / rules"),
                ("融合與生成", "重排、去重、Ollama"),
                ("輪詢結果", "/tasks/{id} + citations"),
            ],
        },
        {
            "y": 4.02,
            "h": 2.24,
            "label": "AI CHAT",
            "title": "OpenClaw 串流聊天",
            "accent": COLORS["teal"],
            "fill": COLORS["teal_soft"],
            "steps": [
                ("載入設定", "chat-config / browser session"),
                ("準備 KB Context", "依題型檢索與來源整理"),
                ("WSS /ws", "瀏覽器連線 FastAPI"),
                ("併發與排隊", "session lock / global slot"),
                ("OpenClaw + LLM", "gateway tools / Ollama"),
                ("串流呈現", "agent events + fallback citation"),
            ],
        },
    ]

    for track in tracks:
        add_rect(slide, 0.50, track["y"], 12.25, track["h"], track["fill"], line=COLORS["line"], radius=True)
        add_rect(slide, 0.50, track["y"], 1.18, track["h"], track["accent"], radius=True)
        add_text(slide, 0.63, track["y"] + 0.38, 0.92, 0.22, track["label"], 7.5, COLORS["white"], True, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, 0.63, track["y"] + 0.80, 0.92, 0.64, track["title"], 12, COLORS["white"], True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
        add_text(slide, 0.63, track["y"] + 1.62, 0.92, 0.18, "AS-IS", 7.2, COLORS["white"], True, align=PP_ALIGN.CENTER)

        x = 1.88
        step_w = 1.62
        for index, (title, body) in enumerate(track["steps"], start=1):
            add_flow_step(slide, x, track["y"] + 0.35, step_w, 1.55, index, title, body, track["accent"])
            if index < len(track["steps"]):
                add_chevron(slide, x + step_w + 0.07, track["y"] + 0.96, 0.20, 0.30, track["accent"])
            x += 1.82

    add_pill(slide, 0.52, 6.58, 1.26, "RESILIENCE", COLORS["amber"], COLORS["navy"], h=0.27, size=7.5)
    add_text(slide, 1.93, 6.57, 10.74, 0.25, "任務重試與 timeout、idempotency key、WebSocket 自動重連、queue/lock TTL、來源不足時的 deterministic fallback。", 8.5, COLORS["slate"], fit=True)
    add_footer(slide, "搜尋路徑確保 KB 結果完整；聊天路徑優先互動體感，並以 KB context 與來源補強回答可信度。")


def add_ingestion_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "04 · Knowledge Supply Chain",
        "文件攝入管線：從原始檔到可查詢知識",
        "同一條攝入骨幹支援 API 上傳、watch folder 與人工匯入，再依文件型態選擇向量、圖譜或混合寫入。",
        5,
    )

    columns = [
        (0.48, 1.70, "01", "輸入通道", COLORS["blue"]),
        (2.62, 1.70, "02", "任務接收", COLORS["blue_2"]),
        (4.77, 1.70, "03", "轉換與增強", COLORS["teal"]),
        (7.34, 1.70, "04", "模式路由", COLORS["amber_dark"]),
        (10.20, 1.70, "05", "知識寫入", COLORS["green"]),
    ]
    for x, y, num, title, accent in columns:
        add_pill(slide, x, y - 0.24, 0.48, num, accent, COLORS["white"], h=0.28, size=7.5)
        add_text(slide, x + 0.58, y - 0.23, 1.55, 0.25, title, 10.2, COLORS["navy"], True, valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)

    inputs = [
        ("API Upload", "/api/upload + task id", "UP"),
        ("Watch Folder", "Celery Beat 定時掃描", "WF"),
        ("Raw / Manual", "批次或維運匯入", "RM"),
    ]
    y = 2.05
    for title, body, icon in inputs:
        add_card(slide, 0.48, y, 1.82, 0.88, title, body, accent=COLORS["blue"], icon=icon, title_size=9.4, body_size=7.6)
        y += 1.01

    add_chevron(slide, 2.36, 3.10, 0.22, 0.42, COLORS["muted"])
    add_card(slide, 2.62, 2.05, 1.85, 1.27, "Redis 任務狀態", "file hash index\n狀態 TTL / queue position", accent=COLORS["blue_2"], icon="RD", title_size=10, body_size=8.1)
    add_card(slide, 2.62, 3.54, 1.85, 1.27, "Celery Ingest", "單檔背景處理\nprocessing lock / retry", accent=COLORS["blue_2"], icon="CI", title_size=10, body_size=8.1)

    add_chevron(slide, 4.53, 3.10, 0.22, 0.42, COLORS["muted"])
    transforms = [
        ("MarkItDown", "Office / PDF / Text → Markdown"),
        ("PDF / Excel 增強", "頁面快照、圖片、表格與 OCR"),
        ("Chunk + Metadata", "標題感知分塊、來源與 image refs"),
    ]
    y = 2.05
    for title, body in transforms:
        add_card(slide, 4.77, y, 2.23, 0.88, title, body, accent=COLORS["teal"], title_size=9.8, body_size=7.8)
        y += 1.01

    add_chevron(slide, 7.08, 3.10, 0.22, 0.42, COLORS["muted"])
    routes = [
        ("REPORT", "報告圖譜 + 向量", COLORS["amber_dark"]),
        ("VECTOR-ONLY", "Lab / Project / Automation", COLORS["amber_dark"]),
        ("SEMANTIC", "4G/5G / WiFi + LLM 實體萃取", COLORS["amber_dark"]),
    ]
    y = 2.05
    for title, body, accent in routes:
        add_card(slide, 7.34, y, 2.40, 0.88, title, body, accent=accent, fill=COLORS["amber_soft"], title_size=9.5, body_size=7.8)
        y += 1.01

    add_chevron(slide, 9.83, 3.10, 0.24, 0.42, COLORS["muted"])
    stores = [
        ("Qdrant", "768D embeddings + chunk payload", "QD", COLORS["green"]),
        ("Neo4j", "Document / Entity / Report graph", "N4", COLORS["green"]),
        ("File Store", "processed / assets / source metadata", "FS", COLORS["green"]),
    ]
    y = 2.05
    for title, body, icon, accent in stores:
        add_card(slide, 10.20, y, 2.58, 0.88, title, body, accent=accent, icon=icon, title_size=9.8, body_size=7.7)
        y += 1.01

    add_rect(slide, 0.48, 5.43, 12.30, 1.12, COLORS["navy"], radius=True)
    add_text(slide, 0.72, 5.62, 1.22, 0.19, "CONTROL PLANE", 8, COLORS["teal"], True)
    controls = [
        ("SHA-256 去重", "避免相同內容重複攝入"),
        ("Cleanup before write", "同文件重攝入先清舊資料"),
        ("Index refresh", "完成後更新 index 與統計"),
        ("Failure visibility", "狀態、錯誤與任務歷史可查"),
    ]
    x = 2.05
    for title, body in controls:
        add_rect(slide, x, 5.63, 2.42, 0.64, COLORS["dark_panel"], line="31516D", radius=True)
        add_text(slide, x + 0.12, 5.70, 2.18, 0.19, title, 8.7, COLORS["white"], True, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, x + 0.10, 5.96, 2.22, 0.17, body, 7.2, "B9CAD9", align=PP_ALIGN.CENTER, fit=True)
        x += 2.55
    add_footer(slide, "攝入不是單一步驟，而是帶有去重、轉檔、資產保留、模式路由、雙庫寫入與可觀測狀態的資料產品管線。")


def add_retrieval_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "05 · Retrieval & AI",
        "混合檢索與答案生成架構",
        "SearchEngine 依題型與模式調度向量、圖譜與規則式來源，再經重排、去重與引用整合後交付答案。",
        6,
    )

    add_rect(slide, 0.52, 1.45, 12.22, 0.62, COLORS["navy"], radius=True)
    add_text(slide, 0.72, 1.64, 1.16, 0.20, "QUERY ROUTER", 8.2, COLORS["teal"], True)
    modes = ["AUTO", "BASIC", "DEEP", "VECTOR", "HYBRID", "HYBRID+", "REPORT GRAPH"]
    x = 2.10
    widths = [0.88, 0.92, 0.88, 1.02, 1.05, 1.13, 1.50]
    for label, width in zip(modes, widths):
        add_pill(slide, x, 1.60, width, label, COLORS["dark_panel"], COLORS["white"], h=0.30, size=7.3, line="31516D")
        x += width + 0.16
    add_text(slide, 10.64, 1.62, 1.75, 0.20, "規則判斷 + 題型提示", 7.4, "B9CAD9", align=PP_ALIGN.RIGHT, fit=True)

    add_text(slide, 0.58, 2.29, 2.90, 0.22, "KNOWLEDGE SOURCES", 9.5, COLORS["navy"], True)
    sources = [
        ("Qdrant Vector", "BAAI/bge-base-zh-v1.5\nCosine similarity + metadata filter", "QD", COLORS["blue"]),
        ("Neo4j Graph", "Document / Entity / Report\nSection / TestItem / Case / Metric", "N4", COLORS["teal"]),
        ("Processed Content", "Markdown、來源 metadata\n規則式摘要與 cleaned context", "FS", COLORS["amber_dark"]),
    ]
    y = 2.65
    for title, body, icon, accent in sources:
        add_card(slide, 0.55, y, 3.04, 1.03, title, body, accent=accent, icon=icon, title_size=10.5, body_size=8.1)
        y += 1.14

    add_text(slide, 4.08, 2.29, 4.35, 0.22, "RETRIEVAL FUSION", 9.5, COLORS["navy"], True)
    fusion = [
        ("1", "平行召回", "向量 raw sources + 圖譜 raw sources"),
        ("2", "業務重排", "題型、專案、case、數值與報告偏好"),
        ("3", "去重與分類", "doc/chunk 去重、storage category、source type"),
        ("4", "引用封裝", "citation distribution、原始文件與片段摘要"),
    ]
    y = 2.65
    for num, title, body in fusion:
        add_rect(slide, 4.05, y, 4.42, 0.76, COLORS["white"], line=COLORS["line"], radius=True)
        add_pill(slide, 4.20, y + 0.19, 0.38, num, COLORS["navy"], COLORS["white"], h=0.32, size=8)
        add_text(slide, 4.74, y + 0.10, 1.34, 0.26, title, 9.8, COLORS["navy"], True, fit=True)
        add_text(slide, 4.74, y + 0.40, 3.47, 0.20, body, 7.8, COLORS["slate"], fit=True)
        y += 0.90

    add_connector(slide, 3.60, 4.10, 4.02, 4.10, COLORS["muted"], width=2.0)
    add_chevron(slide, 3.79, 3.94, 0.25, 0.33, COLORS["muted"])
    add_connector(slide, 8.48, 4.10, 8.90, 4.10, COLORS["muted"], width=2.0)
    add_chevron(slide, 8.69, 3.94, 0.25, 0.33, COLORS["muted"])

    add_text(slide, 9.02, 2.29, 3.66, 0.22, "ANSWER DELIVERY", 9.5, COLORS["navy"], True)
    outputs = [
        ("Deterministic Answer", "報告 / WiFi / comparison 題型可直接呈現 KB 結果", COLORS["blue"]),
        ("Ollama Synthesis", "以檢索 context 生成摘要或混合回答", COLORS["teal"]),
        ("OpenClaw Assistant", "WebSocket 串流、工具能力與 KB 引用補強", COLORS["amber_dark"]),
    ]
    y = 2.65
    for title, body, accent in outputs:
        add_card(slide, 9.02, y, 3.70, 1.03, title, body, accent=accent, title_size=10.5, body_size=8.1)
        y += 1.14

    add_rect(slide, 0.55, 6.22, 12.17, 0.56, COLORS["red_soft"], line="F3CACA", radius=True)
    add_pill(slide, 0.72, 6.36, 1.04, "GUARDRAIL", COLORS["red"], COLORS["white"], h=0.27, size=7.4)
    add_text(slide, 1.95, 6.33, 10.48, 0.26, "找到來源時不得誤稱『查無資料』；片段不足需明確說明限制，並保留來源、引用與 deterministic excerpt fallback。", 8.5, COLORS["ink"], True, fit=True)
    add_footer(slide, "核心差異化在於多路召回後的業務重排與引用封裝，而不是單純把問題交給 LLM。")


def add_data_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "06 · Data Architecture",
        "資料、狀態與知識資產的責任邊界",
        "不同儲存層各自承擔文件生命週期、語意召回、關聯查詢、短期狀態與 AI runtime，不以單一資料庫包辦所有需求。",
        7,
    )

    stores = [
        {
            "title": "FILE STORE",
            "subtitle": "原始資料與可回溯資產",
            "accent": COLORS["blue"],
            "icon": "FS",
            "items": ["raw / watch / uploads", "processed Markdown", "PDF/Excel 圖片與頁面快照", "source.json / chunk versions"],
        },
        {
            "title": "QDRANT",
            "subtitle": "語意內容索引",
            "accent": COLORS["blue_2"],
            "icon": "QD",
            "items": ["collection: knowledge_base", "768D BGE embeddings", "chunk content + doc metadata", "Cosine search / doc filter"],
        },
        {
            "title": "NEO4J",
            "subtitle": "結構化關聯知識",
            "accent": COLORS["teal"],
            "icon": "N4",
            "items": ["Document / Entity", "Report / Section", "TestItem / TestCase / Metric", "跨文件關聯與統計"],
        },
        {
            "title": "REDIS",
            "subtitle": "短期任務與協調狀態",
            "accent": COLORS["amber_dark"],
            "icon": "RD",
            "items": ["Celery broker / result backend", "搜尋與 ingest task state", "cache / queue / locks", "TTL 與 active slot"],
        },
        {
            "title": "OPENCLAW",
            "subtitle": "AI 助理 runtime",
            "accent": COLORS["green"],
            "icon": "OC",
            "items": ["identity / device auth", "workspace / memory", "skills files", "session / gateway config"],
        },
    ]
    x = 0.48
    for store in stores:
        add_rect(slide, x, 1.54, 2.40, 4.20, COLORS["white"], line=COLORS["line"], radius=True)
        add_rect(slide, x, 1.54, 2.40, 0.82, store["accent"], radius=True)
        add_rect(slide, x + 0.16, 1.72, 0.42, 0.42, COLORS["white"], radius=True, transparency=8)
        add_text(slide, x + 0.16, 1.72, 0.42, 0.42, store["icon"], 8, store["accent"], True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
        add_text(slide, x + 0.70, 1.64, 1.52, 0.22, store["title"], 10, COLORS["white"], True, fit=True)
        add_text(slide, x + 0.70, 1.93, 1.52, 0.18, store["subtitle"], 7.3, COLORS["white"], fit=True)
        y = 2.62
        for item in store["items"]:
            add_rect(slide, x + 0.18, y + 0.03, 0.09, 0.09, store["accent"], radius=True)
            add_text(slide, x + 0.36, y - 0.03, 1.84, 0.34, item, 8.3, COLORS["slate"], valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
            y += 0.62
        x += 2.48

    add_rect(slide, 0.48, 5.98, 12.32, 0.80, COLORS["navy"], radius=True)
    governance = [
        ("IDENTITY", "doc_name 與 source metadata 串起文件、chunk、圖譜與原始檔"),
        ("DURABILITY", "release 以 volumes 保存 Redis / Neo4j / Qdrant，app/data 獨立掛載"),
        ("ISOLATION", "目前沒有內建 tenant namespace；跨環境共庫需避免名稱與 point id 衝突"),
    ]
    x = 0.72
    widths = [3.52, 3.72, 3.84]
    for (tag, body), width in zip(governance, widths):
        add_pill(slide, x, 6.18, 0.92, tag, COLORS["dark_panel"], COLORS["teal"], h=0.27, size=7.1, line="31516D")
        add_text(slide, x + 1.08, 6.10, width - 1.08, 0.38, body, 7.8, COLORS["white"], fit=True)
        x += width + 0.34
    add_footer(slide, "Qdrant、Neo4j 與檔案系統是互補資產；Redis 是可淘汰狀態，OpenClaw runtime 則屬於獨立的 AI 執行邊界。")


def add_deployment_topology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "07 · Deployment Topology",
        "現行站台與可攜式 on-prem 發行拓撲",
        "兩種部署共用同一應用程式，但 Qdrant、資料卷、路徑與端口的封裝方式不同；拓撲必須分開理解。",
        8,
    )

    panels = [
        {
            "x": 0.48,
            "title": "CURRENT SITE · 原始站台",
            "badge": "VERIFIED ENTRY",
            "badge_body": "https://127.0.0.1:3030/chat.html",
            "accent": COLORS["blue"],
            "external": ["Host Qdrant :6335", "Host Ollama :11434", "OpenClaw Gateway"],
            "internal": [
                ("Nginx", "3030 → TLS 443"),
                ("FastAPI Web", "Uvicorn :8000 ×4"),
                ("Redis", "broker / result / cache"),
                ("Neo4j", "container graph DB"),
                ("Search Worker", "Celery queue=search"),
                ("Ingest Worker", "Celery queue=ingest"),
                ("Celery Beat", "watch scan scheduler"),
            ],
            "storage": "bind mounts: config / data / uploads / frontend runtime",
        },
        {
            "x": 6.80,
            "title": "ON-PREM RELEASE · 可攜式發行包",
            "badge": "CONFIGURABLE",
            "badge_body": "HTTPS port / project / passwords / worker counts",
            "accent": COLORS["teal"],
            "external": ["Host Ollama", "Host OpenClaw Gateway", "Installer / Upgrade"],
            "internal": [
                ("Nginx", "configurable HTTPS"),
                ("FastAPI Web", "release workers"),
                ("Redis", "persistent volume"),
                ("Neo4j", "persistent volumes"),
                ("Qdrant", "bundled volume"),
                ("Search Worker", "configurable concurrency"),
                ("Ingest Worker", "isolated queue"),
                ("Celery Beat", "scheduled ingest"),
            ],
            "storage": "app/data + runtime/openclaw + named volumes + release metadata",
        },
    ]

    for panel in panels:
        x = panel["x"]
        add_rect(slide, x, 1.48, 6.00, 5.23, COLORS["white"], line=COLORS["line"], radius=True)
        add_rect(slide, x, 1.48, 6.00, 0.52, panel["accent"], radius=True)
        add_text(slide, x + 0.22, 1.63, 5.56, 0.22, panel["title"], 10.6, COLORS["white"], True, fit=True)

        add_text(slide, x + 0.22, 2.16, 1.36, 0.20, "HOST / EXTERNAL", 7.7, panel["accent"], True)
        ex_x = x + 0.22
        for item in panel["external"]:
            add_pill(slide, ex_x, 2.43, 1.71, item, COLORS["panel"], COLORS["ink"], h=0.36, size=7.1, line=COLORS["line"])
            ex_x += 1.84

        add_connector(slide, x + 3.00, 2.84, x + 3.00, 3.15, panel["accent"], width=2.0)
        down = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(x + 2.86), Inches(2.94), Inches(0.28), Inches(0.28))
        down.fill.solid()
        down.fill.fore_color.rgb = rgb(panel["accent"])
        down.line.fill.background()

        add_text(slide, x + 0.22, 3.14, 1.62, 0.20, "DOCKER COMPOSE BOUNDARY", 7.7, panel["accent"], True)
        add_rect(slide, x + 0.20, 3.43, 5.60, 2.08, COLORS["panel"], line=panel["accent"], radius=True, line_width=1.4)
        cols = 4
        cell_w = 1.24
        gap = 0.10
        start_x = x + 0.36
        start_y = 3.61
        for idx, (title, body) in enumerate(panel["internal"]):
            row = idx // cols
            col = idx % cols
            cx = start_x + col * (cell_w + gap)
            cy = start_y + row * 0.87
            add_rect(slide, cx, cy, cell_w, 0.72, COLORS["white"], line=COLORS["line"], radius=True)
            add_text(slide, cx + 0.08, cy + 0.08, cell_w - 0.16, 0.20, title, 8.3, COLORS["navy"], True, align=PP_ALIGN.CENTER, fit=True)
            add_text(slide, cx + 0.06, cy + 0.36, cell_w - 0.12, 0.18, body, 6.8, COLORS["slate"], align=PP_ALIGN.CENTER, fit=True)

        add_pill(slide, x + 0.22, 5.76, 1.20, panel["badge"], panel["accent"], COLORS["white"], h=0.30, size=7.0)
        add_text(slide, x + 1.55, 5.72, 4.12, 0.38, panel["badge_body"], 7.6, COLORS["slate"], valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
        add_rect(slide, x + 0.22, 6.18, 5.55, 0.32, COLORS["navy"], radius=True)
        add_text(slide, x + 0.34, 6.23, 5.30, 0.18, panel["storage"], 7.1, COLORS["white"], align=PP_ALIGN.CENTER, fit=True)

    add_footer(slide, "原始站台的 Qdrant/Ollama 位於 host；on-prem release 將 Qdrant 納入 Compose，但 Ollama 與 OpenClaw gateway 仍是外部 runtime。")


def add_operations_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "08 · Reliability & Governance",
        "可靠度、安全與可維運性：現況控制與管理關注",
        "系統已具備任務隔離、健康檢查與資料一致性控制；企業化下一階段應聚焦認證、秘密管理、模組化與租戶隔離。",
        9,
    )

    panels = [
        {
            "x": 0.50,
            "title": "已落地控制 · EXISTING CONTROLS",
            "accent": COLORS["green"],
            "fill": COLORS["green_soft"],
            "cards": [
                ("傳輸與入口", "TLS 1.2/1.3、Nginx proxy、HTML no-cache、內外端口分離"),
                ("任務可靠度", "Search/Ingest queue 分流、acks late、worker lost reject、retry/timeout"),
                ("一致性控制", "SHA-256 去重、processing lock、idempotency key、cleanup-before-write"),
                ("維運可視性", "health/stats、task state、logs/SSE、admin API、Beat schedule"),
            ],
        },
        {
            "x": 6.78,
            "title": "管理關注 · NEXT CONTROLS",
            "accent": COLORS["red"],
            "fill": COLORS["red_soft"],
            "cards": [
                ("身份與授權", "補齊 Web/API 認證、RBAC 與管理操作稽核；收斂 CORS allow-all"),
                ("秘密與設定", "資料庫密碼、token 與 endpoint 移入 secret store / deployment config"),
                ("應用模組化", "大型 FastAPI __init__.py 拆分 router/service/repository 與契約測試"),
                ("資料隔離與 DR", "tenant namespace、備份/還原演練、RPO/RTO、跨環境資料防碰撞"),
            ],
        },
    ]
    for panel in panels:
        add_rect(slide, panel["x"], 1.48, 6.02, 3.94, panel["fill"], line=COLORS["line"], radius=True)
        add_rect(slide, panel["x"], 1.48, 6.02, 0.50, panel["accent"], radius=True)
        add_text(slide, panel["x"] + 0.22, 1.62, 5.56, 0.22, panel["title"], 10.5, COLORS["white"], True, fit=True)
        y = 2.17
        for idx, (title, body) in enumerate(panel["cards"], start=1):
            add_rect(slide, panel["x"] + 0.22, y, 5.58, 0.67, COLORS["white"], line=COLORS["line"], radius=True)
            add_pill(slide, panel["x"] + 0.35, y + 0.17, 0.36, str(idx), panel["accent"], COLORS["white"], h=0.31, size=7.6)
            add_text(slide, panel["x"] + 0.86, y + 0.08, 1.20, 0.22, title, 9, COLORS["navy"], True, fit=True)
            add_text(slide, panel["x"] + 2.06, y + 0.08, 3.50, 0.43, body, 7.6, COLORS["slate"], valign=MSO_VERTICAL_ANCHOR.MIDDLE, fit=True)
            y += 0.78

    add_text(slide, 0.52, 5.72, 2.20, 0.22, "OPERATING FEEDBACK LOOP", 8.8, COLORS["navy"], True)
    loop = [
        ("Health & Stats", "服務/worker/DB"),
        ("Logs & SSE", "錯誤與即時事件"),
        ("Browser E2E", "真實聊天流程"),
        ("Backup / Preflight", "安裝與還原準備"),
        ("Review & Improve", "根因修正與回歸"),
    ]
    x = 0.52
    for idx, (title, body) in enumerate(loop):
        add_rect(slide, x, 6.07, 2.16, 0.63, COLORS["navy"], radius=True)
        add_text(slide, x + 0.10, 6.14, 1.96, 0.19, title, 8.6, COLORS["white"], True, align=PP_ALIGN.CENTER, fit=True)
        add_text(slide, x + 0.08, 6.39, 2.00, 0.16, body, 7.0, "B9CAD9", align=PP_ALIGN.CENTER, fit=True)
        if idx < len(loop) - 1:
            add_chevron(slide, x + 2.21, 6.23, 0.24, 0.31, COLORS["teal"])
        x += 2.48
    add_footer(slide, "企業級不是只看功能可用；需把身份、秘密、租戶、備援與 SLO 變成可持續驗證的控制面。")


def add_management_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, dark=True)
    add_text(slide, 0.72, 0.44, 4.8, 0.20, "09 · MANAGEMENT TAKEAWAY", 9.2, COLORS["teal"], True)
    add_text(slide, 0.72, 0.80, 9.5, 0.60, "主管結論：架構已具備企業知識平台骨幹", 27, COLORS["white"], True, font=FONT_DISPLAY, fit=True)
    add_text(slide, 0.74, 1.42, 10.8, 0.32, "現行系統可支援 on-prem 部署、雙引擎檢索、背景攝入與 OpenClaw 助理；下一階段應從『可運作』走向『可治理、可擴展、可稽核』。", 11.5, "C7D8E7", fit=True)
    add_pill(slide, 11.76, 0.50, 0.78, "10", COLORS["blue"], COLORS["white"], h=0.34, size=9.5)

    columns = [
        ("READY NOW", "現有優勢", COLORS["blue"], ["Hybrid RAG + GraphRAG", "FastAPI / Celery 非同步骨幹", "可攜式 on-prem release", "來源、Chunk 與管理 API"]),
        ("DEPENDENCIES", "關鍵依賴", COLORS["amber"], ["Ollama 模型效能與可用性", "OpenClaw gateway / identity", "資料品質與攝入成功率", "Neo4j / Qdrant 持久性"]),
        ("DECISIONS", "需要主管支持", COLORS["teal"], ["明確服務 SLO 與容量目標", "指定資料治理與平台 owner", "核定安全與秘密管理優先級", "決定是否需要多租戶隔離"]),
    ]
    x = 0.72
    for tag, title, accent, bullets in columns:
        add_rect(slide, x, 2.04, 3.86, 2.44, COLORS["dark_panel"], line="31516D", radius=True)
        add_pill(slide, x + 0.18, 2.22, 1.16, tag, accent, COLORS["navy"] if accent == COLORS["amber"] else COLORS["white"], h=0.30, size=7.4)
        add_text(slide, x + 0.18, 2.68, 3.48, 0.26, title, 12.5, COLORS["white"], True, fit=True)
        y = 3.12
        for bullet in bullets:
            add_rect(slide, x + 0.20, y + 0.06, 0.08, 0.08, accent, radius=True)
            add_text(slide, x + 0.38, y, 3.22, 0.25, bullet, 8.5, "C7D8E7", fit=True)
            y += 0.31
        x += 4.08

    add_text(slide, 0.72, 4.82, 3.6, 0.20, "90-DAY ARCHITECTURE PRIORITIES", 8.8, COLORS["teal"], True)
    roadmap = [
        ("P0", "安全基線", "Secrets / CORS / API auth", COLORS["red"]),
        ("P1", "模組邊界", "FastAPI routers + services", COLORS["blue"]),
        ("P1", "營運 SLO", "metrics / alert / backup drill", COLORS["teal"]),
        ("P2", "規模治理", "tenant namespace / DR", COLORS["amber"]),
    ]
    x = 0.72
    for idx, (priority, title, body, accent) in enumerate(roadmap):
        add_rect(slide, x, 5.18, 2.82, 0.90, COLORS["dark_panel"], line="31516D", radius=True)
        add_pill(slide, x + 0.14, 5.35, 0.50, priority, accent, COLORS["white"] if accent != COLORS["amber"] else COLORS["navy"], h=0.28, size=7.2)
        add_text(slide, x + 0.78, 5.25, 1.82, 0.23, title, 9.6, COLORS["white"], True, fit=True)
        add_text(slide, x + 0.78, 5.57, 1.82, 0.18, body, 7.5, "B9CAD9", fit=True)
        if idx < len(roadmap) - 1:
            add_chevron(slide, x + 2.92, 5.47, 0.26, 0.32, accent)
        x += 3.14

    add_rect(slide, 0.72, 6.36, 11.84, 0.50, COLORS["blue"], radius=True)
    add_text(slide, 0.94, 6.47, 11.38, 0.23, "建議定位：將 Knowledge Base 視為可持續營運的企業 AI 平台，而非單一聊天介面或一次性 RAG 專案。", 10.2, COLORS["white"], True, align=PP_ALIGN.CENTER, fit=True)
    add_footer(slide, "既有骨幹已足以支撐下一階段投資；管理重點是把技術能力轉成有 owner、有 SLO、有控制面的平台產品。", dark=True)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    props = prs.core_properties
    props.title = "Knowledge Base 企業級系統架構"
    props.subject = "Knowledge Base current-state enterprise architecture"
    props.author = "Knowledge Base Architecture Team"
    props.keywords = "Knowledge Base, FastAPI, Celery, Qdrant, Neo4j, Ollama, OpenClaw, Architecture"
    props.comments = "Generated from the current repository architecture on 2026-07-17."

    add_cover(prs)
    add_executive_overview(prs)
    add_logical_architecture(prs)
    add_query_chat_flow(prs)
    add_ingestion_flow(prs)
    add_retrieval_architecture(prs)
    add_data_architecture(prs)
    add_deployment_topology(prs)
    add_operations_security(prs)
    add_management_summary(prs)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    output = build_presentation()
    print(f"saved {output}")
