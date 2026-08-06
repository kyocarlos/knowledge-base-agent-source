from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "sub2api_simple_architecture.pptx"
FONT = "Noto Sans CJK TC"
W, H = 13.333, 7.5
C = {
    "navy": "0B1F33", "ink": "193247", "muted": "607589", "line": "C8D5DF",
    "blue": "2F6FED", "blue_fill": "EAF1FF", "teal": "009E9A", "teal_fill": "E4F7F5",
    "purple": "7653A6", "purple_fill": "F3EEFA", "amber": "D68B16", "amber_fill": "FFF3D6",
    "gray_fill": "F4F7FA", "white": "FFFFFF",
}


def rgb(value):
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill, line=None, rounded=True, width=1.2):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()
    return shape


def text(slide, x, y, w, h, value, size=12, color=None, bold=False,
         align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = valign
    frame.margin_left = frame.margin_right = Pt(5)
    frame.margin_top = frame.margin_bottom = Pt(2)
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color or C["ink"])
    return box


def arrow(slide, x1, y1, x2, y2, color=None):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = rgb(color or C["line"])
    connector.line.width = Pt(2.0)
    end = connector._element.spPr.get_or_add_ln()
    from pptx.oxml.xmlchemy import OxmlElement
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    end.append(tail)
    return connector


def node(slide, x, y, w, h, title, body, fill, accent):
    rect(slide, x, y, w, h, fill, accent, True, 1.8)
    text(slide, x + 0.10, y + 0.16, w - 0.20, 0.38, title, 16, accent, True)
    text(slide, x + 0.14, y + 0.70, w - 0.28, h - 0.82, body, 10, C["muted"])


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C["white"])

    text(slide, 0.55, 0.30, 8.0, 0.45, "Sub2API 簡易架構圖", 24, C["ink"], True, PP_ALIGN.LEFT)
    text(slide, 0.58, 0.82, 8.5, 0.28, "OpenAI-compatible API Gateway，統一管理模型路由與用量", 10, C["muted"], False, PP_ALIGN.LEFT)

    # Main request path
    node(slide, 0.72, 2.12, 2.45, 1.45, "OpenClaw / 客戶端", "外部應用程式\n呼叫 /v1 API", C["blue_fill"], C["blue"])
    node(slide, 5.02, 1.85, 3.28, 2.02, "Sub2API", "API Gateway\nAPI Key → 分組 → 帳號\n/v1/models\n/v1/chat/completions", C["teal_fill"], C["teal"])
    node(slide, 10.15, 2.12, 2.45, 1.45, "Ollama LLM", "本機模型\n例如 gemma4:12b\nqwen3-coder-next", C["purple_fill"], C["purple"])
    arrow(slide, 3.17, 2.84, 5.02, 2.84, C["blue"])
    arrow(slide, 8.30, 2.84, 10.15, 2.84, C["teal"])
    text(slide, 3.55, 2.52, 1.15, 0.25, "HTTPS / API", 8, C["muted"], True)
    text(slide, 8.55, 2.52, 1.15, 0.25, "upstream", 8, C["muted"], True)

    # Supporting services
    rect(slide, 3.48, 4.55, 6.40, 1.55, C["gray_fill"], C["line"], True, 1.2)
    text(slide, 3.72, 4.73, 2.2, 0.25, "Sub2API 內部服務", 10, C["ink"], True, PP_ALIGN.LEFT)
    node(slide, 4.02, 5.12, 2.25, 0.68, "PostgreSQL", "群組、帳號、API key、用量", C["white"], C["amber"])
    node(slide, 7.05, 5.12, 2.25, 0.68, "Redis", "快取、工作狀態、限流", C["white"], C["blue"])
    arrow(slide, 6.66, 3.87, 5.15, 5.12, C["line"])
    arrow(slide, 6.66, 3.87, 8.18, 5.12, C["line"])

    # Isolation note
    rect(slide, 0.72, 5.18, 2.45, 0.95, C["amber_fill"], C["amber"], True, 1.2)
    text(slide, 0.88, 5.32, 2.12, 0.55, "Knowledge Base\n獨立系統，不共用 Sub2API DB", 9, C["amber"], True)
    text(slide, 0.72, 6.72, 11.8, 0.22, "重點：客戶端只連 Sub2API API；不直接連 PostgreSQL、Redis 或 Ollama。", 10, C["navy"], True, PP_ALIGN.CENTER)
    text(slide, 10.20, 6.72, 2.40, 0.22, "Port: 18080", 8, C["muted"], False, PP_ALIGN.RIGHT)

    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
