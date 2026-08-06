from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "sub2api_development_progress.pptx"
W, H = 13.333, 7.5
FONT = "Noto Sans CJK TC"

C = {
    "navy": "0B1F33", "navy2": "153A56", "ink": "193247", "muted": "607589",
    "blue": "2F6FED", "blue_soft": "EAF1FF", "teal": "009E9A", "teal_soft": "E4F7F5",
    "green": "2E9D68", "green_soft": "E7F6ED", "amber": "D68B16", "amber_soft": "FFF3D6",
    "red": "C94C4C", "red_soft": "FCEBEC", "line": "D7E0E8", "panel": "F5F8FB",
    "white": "FFFFFF", "purple": "7653A6", "purple_soft": "F3EEFA",
}


def rgb(value):
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill, line=None, rounded=True, width=1.0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(width)
    else:
        shape.line.fill.background()
    return shape


def text(slide, x, y, w, h, value, size=11, color=None, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP, fit=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if fit else MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = frame.margin_right = Pt(3)
    frame.margin_top = frame.margin_bottom = Pt(1)
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    p.line_spacing = 1.0
    r = p.add_run()
    r.text = value
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color or C["ink"])
    return box


def line(slide, x1, y1, x2, y2, color=None, width=1.3):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = rgb(color or C["line"])
    conn.line.width = Pt(width)
    return conn


def title(slide, no, heading, sub=""):
    text(slide, 0.42, 0.16, 0.40, 0.22, f"{no:02d}", 8, C["muted"], True)
    text(slide, 0.86, 0.11, 7.0, 0.34, heading, 16, C["ink"], True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    if sub:
        text(slide, 7.7, 0.16, 5.15, 0.22, sub, 7, C["muted"], False, PP_ALIGN.RIGHT)
    line(slide, 0.42, 0.55, 12.9, 0.55, C["line"], 0.8)


def footer(slide, page):
    text(slide, 0.42, 7.22, 8.0, 0.16, "Sub2API 開發進度 · 2026-08-03 · 敏感設定已遮罩", 5.8, C["muted"])
    text(slide, 11.9, 7.22, 0.95, 0.16, f"{page:02d} / 08", 5.8, C["muted"], False, PP_ALIGN.RIGHT)


def pill(slide, x, y, w, label, fill, color):
    obj = rect(slide, x, y, w, 0.28, fill, fill, True)
    text(slide, x, y + 0.01, w, 0.23, label, 7, color, True, PP_ALIGN.CENTER, MSO_VERTICAL_ANCHOR.MIDDLE)
    return obj


def card(slide, x, y, w, h, heading, body, accent, fill=None, body_size=9.5):
    rect(slide, x, y, w, h, fill or C["white"], C["line"], True)
    rect(slide, x, y, 0.07, h, accent, accent, False)
    text(slide, x + 0.18, y + 0.14, w - 0.30, 0.28, heading, 11, accent, True, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    text(slide, x + 0.18, y + 0.53, w - 0.30, h - 0.64, body, body_size, C["muted"])


def node(slide, x, y, w, h, heading, body, fill, accent):
    rect(slide, x, y, w, h, fill, accent, True, 1.4)
    text(slide, x + 0.08, y + 0.10, w - 0.16, 0.28, heading, 10, accent, True, PP_ALIGN.CENTER, MSO_VERTICAL_ANCHOR.MIDDLE)
    text(slide, x + 0.10, y + 0.47, w - 0.20, h - 0.55, body, 7.8, C["muted"], False, PP_ALIGN.CENTER)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C["white"])
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1 Cover
    s = new_slide(prs)
    rect(s, 0, 0, W, H, C["navy"], C["navy"], False)
    rect(s, 0, 5.92, W, 1.58, C["navy2"], C["navy2"], False)
    text(s, 0.72, 0.85, 8.2, 0.45, "Sub2API", 30, C["white"], True)
    text(s, 0.76, 1.55, 10.7, 0.72, "開發進度與下一階段規劃", 25, "DDEBFF", True)
    text(s, 0.78, 2.55, 8.6, 0.55, "Ollama OpenAI-compatible Gateway · 群組路由 · API Key · 用量計費", 12, "B8D4E8")
    pill(s, 0.80, 3.62, 1.30, "已完成部署", C["green"], C["white"])
    pill(s, 2.24, 3.62, 1.36, "正在驗證", C["amber"], C["white"])
    pill(s, 3.74, 3.62, 1.62, "待強化治理", C["purple"], C["white"])
    node(s, 8.55, 4.20, 1.45, 0.95, "OpenClaw", "外部客戶端", C["blue_soft"], C["blue"])
    node(s, 10.35, 4.20, 1.45, 0.95, "Sub2API", "API Gateway", C["teal_soft"], C["teal"])
    line(s, 10.0, 4.67, 10.35, 4.67, "9EC4D8", 2)
    text(s, 0.78, 6.42, 5.0, 0.25, "內部開發簡報 · 2026-08-03", 8, "A9C1D4")

    # 2 executive status
    s = new_slide(prs); title(s, 2, "目前狀態總覽", "結論先行"); footer(s, 2)
    card(s, 0.58, 0.93, 2.85, 1.42, "部署", "Sub2API 0.1.168 已以獨立 Compose project 運行。\nWeb health：HTTP 200", C["green"], C["green_soft"])
    card(s, 3.65, 0.93, 2.85, 1.42, "Ollama", "主機 Ollama 已接入 OpenAI-compatible upstream。\n實際回傳 22 個模型。", C["blue"], C["blue_soft"])
    card(s, 6.72, 0.93, 2.85, 1.42, "路由", "Ollama Local、Anderson_H、WifiSit01_DA40 已建立或驗證。", C["teal"], C["teal_soft"])
    card(s, 9.79, 0.93, 2.85, 1.42, "焦點問題", "Test_Local 仍使用預設 13 個 OpenAI 模型清單，尚未建立 API key。", C["amber"], C["amber_soft"])
    text(s, 0.62, 2.78, 4.3, 0.28, "已完成能力", 12, C["ink"], True)
    card(s, 0.58, 3.18, 3.85, 2.65, "隔離部署", "獨立 network、PostgreSQL、Redis 與資料目錄。\n不加入 Knowledge Base 的 Compose network。\n不修改 Neo4j、Qdrant、KB Redis。", C["green"], C["white"])
    card(s, 4.68, 3.18, 3.85, 2.65, "API Gateway", "提供 OpenAI-compatible：\n/v1/models\n/v1/chat/completions\n支援群組、API key、用量紀錄與 upstream 路由。", C["blue"], C["white"])
    card(s, 8.78, 3.18, 3.85, 2.65, "管理與計費", "已具備群組、帳號、模型白名單、token 用量與成本估算。\n正式對外前仍需完成安全治理。", C["purple"], C["white"])

    # 3 topology
    s = new_slide(prs); title(s, 3, "部署拓撲與資料隔離", "Knowledge Base 與 Sub2API 並行"); footer(s, 3)
    rect(s, 0.58, 0.92, 5.75, 5.65, C["panel"], C["blue"], True, 1.4)
    text(s, 0.82, 1.10, 2.5, 0.28, "Sub2API Compose project", 12, C["blue"], True)
    node(s, 1.02, 1.75, 1.58, 0.98, "Web", "0.0.0.0:18080\ncontainer :8080", C["blue_soft"], C["blue"])
    node(s, 3.22, 1.75, 1.58, 0.98, "PostgreSQL", "內部 5432\n群組/帳號/用量", C["teal_soft"], C["teal"])
    node(s, 2.12, 3.45, 1.58, 0.98, "Redis", "內部 6379\ncache / queue", C["purple_soft"], C["purple"])
    line(s, 2.60, 2.73, 2.60, 3.45, C["line"], 1.5); line(s, 3.22, 2.24, 3.22, 2.24, C["line"])
    text(s, 0.92, 5.25, 5.0, 0.65, "資料目錄：sub2api-deploy/data/\n網路：sub2api_sub2api-network\n容器：sub2api、sub2api-postgres、sub2api-redis", 8.5, C["muted"])
    rect(s, 6.98, 0.92, 5.75, 5.65, C["panel"], C["teal"], True, 1.4)
    text(s, 7.22, 1.10, 3.3, 0.28, "Knowledge Base project", 12, C["teal"], True)
    node(s, 7.45, 1.75, 1.55, 0.98, "FastAPI", "Search / Upload\nChat proxy", C["blue_soft"], C["blue"])
    node(s, 9.65, 1.75, 1.55, 0.98, "Celery", "Search / Ingest\nbackground task", C["amber_soft"], C["amber"])
    node(s, 8.55, 3.45, 1.55, 0.98, "Neo4j", "knowledge graph", C["teal_soft"], C["teal"])
    node(s, 10.75, 3.45, 1.55, 0.98, "Qdrant", "vector search", C["purple_soft"], C["purple"])
    line(s, 8.22, 2.73, 8.95, 3.45, C["line"], 1.4); line(s, 10.42, 2.73, 9.35, 3.45, C["line"], 1.4); line(s, 10.42, 2.73, 11.15, 3.45, C["line"], 1.4)
    text(s, 7.30, 5.25, 5.0, 0.65, "Sub2API 未加入 KB network。\n兩套 PostgreSQL / Redis 完全隔離。\n目前 live 狀態：knowledge-base running(8)、sub2api running(3)。", 8.5, C["muted"])

    # 4 timeline
    s = new_slide(prs); title(s, 4, "已完成里程碑", "2026-07-30 至 2026-08-03"); footer(s, 4)
    line(s, 1.00, 2.18, 12.2, 2.18, C["line"], 2)
    events = [
        (1.0, "07/30", "隔離安裝", "Sub2API 0.1.168\n獨立 Compose / DB / Redis", C["green"]),
        (3.3, "07/30", "接入 Ollama", "host.docker.internal\n同步 22 個模型", C["blue"]),
        (5.6, "07/31", "Anderson_H", "獨立群組與 API key\n指定模型 chat 200", C["teal"]),
        (7.9, "07/31", "WifiSit01", "OpenAI OAuth 綁定\ngpt-5.4 chat 200", C["purple"]),
        (10.2, "08/03", "Test_Local", "發現仍是預設模型清單\n進入修正與驗證", C["amber"]),
    ]
    for x, date, head, body, color in events:
        rect(s, x - 0.07, 2.08, 0.14, 0.14, color, color, True)
        text(s, x - 0.42, 1.38, 0.84, 0.28, date, 8, color, True, PP_ALIGN.CENTER)
        card(s, x - 0.72, 2.62, 1.72, 1.78, head, body, color, C["white"], 8.2)
    card(s, 0.78, 5.10, 3.80, 0.92, "驗證已通過", "健康檢查、容器健康、Ollama chat、群組路由與用量資料均已完成局部驗證。", C["green"], C["green_soft"], 8.5)
    card(s, 4.78, 5.10, 3.80, 0.92, "目前限制", "部分群組的 /v1/models 清單不等於 upstream 實際模型清單。", C["amber"], C["amber_soft"], 8.5)
    card(s, 8.78, 5.10, 3.80, 0.92, "下一步", "完成 Test_Local 模型清單、API key、外部呼叫與安全收斂。", C["blue"], C["blue_soft"], 8.5)

    # 5 Ollama flow
    s = new_slide(prs); title(s, 5, "Ollama 整合資料流", "已完成 upstream connectivity"); footer(s, 5)
    node(s, 0.72, 2.00, 2.05, 1.18, "OpenClaw / Client", "OpenAI-compatible\n/v1/chat/completions", C["blue_soft"], C["blue"])
    node(s, 3.38, 2.00, 2.05, 1.18, "Sub2API", "API key → group\naccount selection", C["teal_soft"], C["teal"])
    node(s, 6.04, 2.00, 2.05, 1.18, "Ollama Local", "account id 1\nactive / schedulable", C["purple_soft"], C["purple"])
    node(s, 8.70, 2.00, 2.05, 1.18, "Host Ollama", "127.0.0.1:11434\n/v1/models", C["green_soft"], C["green"])
    node(s, 11.36, 2.00, 1.25, 1.18, "LLM", "22 models", C["amber_soft"], C["amber"])
    for x in (2.77, 5.43, 8.09, 10.75):
        line(s, x, 2.59, x + 0.60, 2.59, C["line"], 2)
    text(s, 0.82, 1.30, 10.9, 0.30, "實際 upstream 已回傳 22 個模型；問題發生在 group-level /v1/models 白名單，而非 Ollama 連線。", 10, C["ink"], True)
    card(s, 1.02, 4.18, 3.35, 1.28, "已驗證", "Ollama 直接 /v1/chat/completions\n使用 gemma4:12b 回 HTTP 200。", C["green"], C["green_soft"])
    card(s, 4.98, 4.18, 3.35, 1.28, "已驗證", "Anderson_H group id 6\n使用 Ollama 模型路由成功。", C["blue"], C["blue_soft"])
    card(s, 8.94, 4.18, 3.35, 1.28, "需注意", "Ollama cloud 模型可能依賴外部授權；embedding 模型不適合一般 chat。", C["amber"], C["amber_soft"])

    # 6 Test_Local root cause
    s = new_slide(prs); title(s, 6, "Test_Local 問題定位", "目前最明確的功能缺口"); footer(s, 6)
    rect(s, 0.62, 0.95, 5.5, 5.45, C["red_soft"], C["red"], True, 1.4)
    text(s, 0.92, 1.22, 4.85, 0.34, "現況：群組清單不是 upstream 清單", 14, C["red"], True)
    text(s, 0.94, 1.86, 4.72, 2.10, "Test_Local group id 8\n\nplatform = openai\nstatus = active\nmodels_list_config.enabled = true\n目前清單 = 預設 13 個 OpenAI 模型\nAPI key = 尚未建立", 11, C["ink"], False)
    pill(s, 0.95, 4.92, 2.15, "不是 Ollama 無模型", C["red"], C["white"])
    rect(s, 6.78, 0.95, 5.9, 5.45, C["panel"], C["blue"], True, 1.4)
    text(s, 7.08, 1.22, 5.1, 0.34, "建議修正流程", 14, C["blue"], True)
    steps = [
        ("01", "同步 upstream", "在 Ollama Local account 執行 sync upstream models"),
        ("02", "更新群組白名單", "在 Test_Local 選入 Ollama 實際模型並儲存"),
        ("03", "建立專用 API key", "API key 綁定 group id 8，避免無法路由"),
        ("04", "端到端驗證", "測 /v1/models 與 gemma4:12b chat completion"),
    ]
    for i, (num, head, body) in enumerate(steps):
        y = 1.88 + i * 0.88
        pill(s, 7.10, y, 0.52, num, C["blue"], C["white"])
        text(s, 7.82, y - 0.01, 2.1, 0.25, head, 10.5, C["ink"], True)
        text(s, 9.55, y - 0.01, 2.65, 0.34, body, 8, C["muted"])
    text(s, 7.10, 5.62, 5.0, 0.42, "預期結果：Test_Local /v1/models 顯示實際允許的 Ollama 模型，而非預設 OpenAI 清單。", 8.5, C["blue"], True)

    # 7 security
    s = new_slide(prs); title(s, 7, "正式對外前的治理缺口", "功能可用 ≠ 已完成生產化"); footer(s, 7)
    card(s, 0.64, 0.98, 3.82, 1.45, "網路暴露", "目前 Web 綁定 0.0.0.0:18080，使用 HTTP。\n正式對外應放在 VPN 或 HTTPS reverse proxy 後。", C["red"], C["red_soft"])
    card(s, 4.76, 0.98, 3.82, 1.45, "Secrets", ".env 與 data/config.yaml 含敏感設定。\n需限制權限、輪換 secret、避免備份外洩。", C["red"], C["red_soft"])
    card(s, 8.88, 0.98, 3.82, 1.45, "Upstream allowlist", "目前為支援 host Ollama 而關閉 URL allowlist。\n需重新設計 egress/SSRF 防護。", C["amber"], C["amber_soft"])
    card(s, 0.64, 3.00, 3.82, 1.45, "版本固定", "目前 image 使用 latest。\n正式環境應固定 release tag 或 digest。", C["amber"], C["amber_soft"])
    card(s, 4.76, 3.00, 3.82, 1.45, "群組隔離", "WifiSit01 OAuth account 目前曾同時綁多個群組。\n需明確定義隔離政策。", C["purple"], C["purple_soft"])
    card(s, 8.88, 3.00, 3.82, 1.45, "模型能力", "模型清單、模型可用性與 upstream 實際能力需分開驗證。", C["blue"], C["blue_soft"])
    text(s, 0.72, 5.20, 11.4, 0.30, "建議順序：VPN/HTTPS → secret rotation → allowlist → image pinning → group/account policy → automated smoke tests", 11, C["ink"], True, PP_ALIGN.CENTER)
    rect(s, 1.40, 5.78, 10.50, 0.48, C["navy"], C["navy"], True)
    text(s, 1.55, 5.88, 10.20, 0.24, "禁止把 Sub2API PostgreSQL、Redis 或 upstream secret 直接提供給外部 Agent。", 9, C["white"], True, PP_ALIGN.CENTER, MSO_VERTICAL_ANCHOR.MIDDLE)

    # 8 roadmap
    s = new_slide(prs); title(s, 8, "下一階段開發與驗收", "從可用走向可交付"); footer(s, 8)
    text(s, 0.68, 0.92, 3.2, 0.28, "P0：先完成可重複驗證", 12, C["red"], True)
    card(s, 0.64, 1.32, 3.82, 2.30, "Test_Local 完整打通", "同步 account 1 的 Ollama 模型\n更新 group 8 白名單\n建立 Test_Local API key\n驗證 /v1/models + chat completion", C["red"], C["red_soft"])
    card(s, 4.76, 1.32, 3.82, 2.30, "自動化 smoke test", "每次部署驗證：\nhealth、models、指定模型 chat、錯誤群組、無可用 account、rate limit。", C["blue"], C["blue_soft"])
    card(s, 8.88, 1.32, 3.82, 2.30, "外部 OpenClaw", "使用專用 API key\n明確指定 model\n不依賴空的 /v1/models 自動探索\n記錄 request / group / model", C["teal"], C["teal_soft"])
    text(s, 0.68, 4.15, 3.2, 0.28, "P1：完成生產治理", 12, C["amber"], True)
    card(s, 0.64, 4.54, 3.82, 1.50, "安全", "VPN、HTTPS、allowlist、secret rotation、固定 image digest。", C["amber"], C["amber_soft"])
    card(s, 4.76, 4.54, 3.82, 1.50, "可觀測性", "群組、account、model、latency、token、cost與錯誤碼 dashboard。", C["purple"], C["purple_soft"])
    card(s, 8.88, 4.54, 3.82, 1.50, "文件與交付", "建立外部 OpenClaw 設定文件、API 使用範例與回復 SOP。", C["green"], C["green_soft"])
    rect(s, 0.64, 6.42, 12.06, 0.42, C["navy"], C["navy"], True)
    text(s, 0.82, 6.50, 11.70, 0.22, "完成標準：Test_Local 可由獨立 API key 穩定呼叫指定 Ollama 模型，且部署重啟後設定、路由、用量與安全策略仍一致。", 8.8, C["white"], True, PP_ALIGN.CENTER, MSO_VERTICAL_ANCHOR.MIDDLE)

    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
