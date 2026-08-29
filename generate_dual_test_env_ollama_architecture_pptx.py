from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/home/da40_ai_gb10/knowledge-base/dual_test_env_ollama_architecture.pptx")
prs = None


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0.0, radius=False):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def set_text(frame, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Sans CJK TC"):
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return frame


def textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, text, size, color, bold=bold, align=align)
    return box


def badge(slide, left, top, width, text, fill, color):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        Inches(0.30),
        fill=fill,
        line=fill,
        radius=True,
    )
    set_text(shape.text_frame, text, 8.5, color, align=PP_ALIGN.CENTER)
    return shape


def card(slide, left, top, width, height, title, body, accent, fill="FFFFFF"):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=fill,
        line="CBD5E1",
        radius=True,
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(7)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Noto Sans CJK TC"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = rgb(accent)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Noto Sans CJK TC"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = rgb("475569")
    return shape


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    for left, top, width, height, color in [
        (Inches(-0.4), Inches(-0.3), Inches(3.8), Inches(3.8), "DBEAFE"),
        (Inches(9.3), Inches(-0.3), Inches(4.3), Inches(3.8), "D1FAE5"),
        (Inches(8.8), Inches(4.8), Inches(4.4), Inches(2.7), "FEF3C7"),
    ]:
        deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        deco.fill.solid()
        deco.fill.fore_color.rgb = rgb(color)
        deco.fill.transparency = 0.80
        deco.line.fill.background()

    border = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.22),
        Inches(0.18),
        Inches(12.89),
        Inches(7.14),
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    border.fill.transparency = 1


def header(slide, title, subtitle, badge_text, badge_fill, badge_color):
    textbox(slide, Inches(0.52), Inches(0.32), Inches(4.2), Inches(0.20), "OpenClaw / Shared Ollama Architecture", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.56), Inches(9.6), Inches(0.42), title, 24, "0F172A", bold=True)
    textbox(slide, Inches(0.52), Inches(0.98), Inches(10.5), Inches(0.34), subtitle, 12, "475569")
    badge(slide, Inches(10.72), Inches(0.50), Inches(1.80), badge_text, badge_fill, badge_color)


def connector(slide, x1, y1, x2, y2, color="8EA2BE"):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(2)
    return line


def add_cover(slide):
    add_background(slide)
    header(
        slide,
        "雙測試環境共用 DGX GB10 Ollama 的 OpenClaw 架構",
        "Anritsu MT8000 與 Amarisoft 各自有獨立的 OpenClaw AI Agent，用來切換與控制儀器；兩邊共用同一台 DGX GB10 上的 Ollama 推論服務。",
        "架構簡報",
        "EFF6FF",
        "1D4ED8",
    )

    card(slide, Inches(0.72), Inches(1.72), Inches(4.20), Inches(1.28), "核心概念", "兩個測試環境各自獨立運作，但 LLM 推論層集中到同一台 DGX GB10，方便統一維運與模型管理。", "2563EB")
    card(slide, Inches(0.72), Inches(3.20), Inches(4.20), Inches(1.28), "Anritsu MT8000", "OpenClaw Agent A 使用 qwen3.5:35b，負責 Anritsu 測試流程與儀器切換。", "16A34A")
    card(slide, Inches(0.72), Inches(4.68), Inches(4.20), Inches(1.28), "Amarisoft", "OpenClaw Agent B 使用 gemma4:12b，負責 Amarisoft 測試流程與儀器切換。", "D97706")

    card(
        slide,
        Inches(5.20),
        Inches(1.88),
        Inches(6.90),
        Inches(1.62),
        "投影片會說明什麼",
        "1. 整體架構與資料流向\n2. 兩個環境各自的責任邊界\n3. 共用 Ollama 的模型分派方式\n4. 這樣設計的好處與注意事項",
        "047857",
    )

    badge(slide, Inches(5.22), Inches(4.08), Inches(1.52), "共用推論層", "ECFDF5", "047857")
    badge(slide, Inches(6.88), Inches(4.08), Inches(1.62), "環境隔離", "FEF3C7", "B45309")
    badge(slide, Inches(8.62), Inches(4.08), Inches(1.64), "模型分工", "DBEAFE", "1D4ED8")

    card(slide, Inches(5.20), Inches(4.55), Inches(6.90), Inches(1.20), "一句話總結", "兩個 OpenClaw Agent 分別控制不同的測試儀器環境，但都把腦袋接到同一台 DGX GB10 上的 Ollama，依環境使用不同模型。", "1D4ED8")


def add_overview_slide(slide):
    add_background(slide)
    header(
        slide,
        "雙環境共用同一台 DGX GB10 的 LLM",
        "Anritsu 與 Amarisoft 各自保留獨立的 OpenClaw 控制層，但兩邊都把推論請求送到同一台 DGX GB10 上的 Ollama。",
        "總覽",
        "ECFDF5",
        "047857",
    )

    badge(slide, Inches(0.86), Inches(1.72), Inches(1.76), "Anritsu 環境", "DBEAFE", "1D4ED8")
    badge(slide, Inches(10.70), Inches(1.72), Inches(1.86), "Amarisoft 環境", "FEF3C7", "B45309")
    badge(slide, Inches(4.95), Inches(1.60), Inches(2.85), "共用 DGX GB10 LLM", "ECFDF5", "047857")

    card(
        slide,
        Inches(0.72),
        Inches(2.10),
        Inches(3.20),
        Inches(2.05),
        "OpenClaw A",
        "獨立控制 Anritsu MT8000\n模型：qwen3.5:35b\n職責：流程判斷、儀器切換、回覆生成",
        "2563EB",
    )

    card(
        slide,
        Inches(9.40),
        Inches(2.10),
        Inches(3.20),
        Inches(2.05),
        "OpenClaw B",
        "獨立控制 Amarisoft\n模型：gemma4:12b\n職責：流程判斷、儀器切換、回覆生成",
        "D97706",
    )

    card(
        slide,
        Inches(4.48),
        Inches(2.02),
        Inches(4.26),
        Inches(2.38),
        "DGX GB10 / Ollama",
        "同一台推論主機提供兩邊共用的 LLM 服務。\n\nEndpoint: http://61.216.9.52:11434\nOpenAI compatible: /v1\n\n左、右兩個 OpenClaw 都連到這一層。",
        "047857",
    )

    connector(slide, Inches(3.92), Inches(3.02), Inches(4.48), Inches(3.02), "8EA2BE")
    connector(slide, Inches(8.74), Inches(3.02), Inches(9.40), Inches(3.02), "8EA2BE")

    connector(slide, Inches(3.92), Inches(3.98), Inches(4.48), Inches(3.38), "8EA2BE")
    connector(slide, Inches(8.74), Inches(3.98), Inches(8.74), Inches(3.38), "8EA2BE")

    badge(slide, Inches(0.90), Inches(4.48), Inches(1.82), "獨立控制層", "EFF6FF", "1D4ED8")
    badge(slide, Inches(2.86), Inches(4.48), Inches(1.78), "各自工具", "EEF2FF", "4338CA")
    badge(slide, Inches(5.00), Inches(4.48), Inches(1.88), "共用 LLM", "ECFDF5", "047857")
    badge(slide, Inches(10.02), Inches(4.48), Inches(1.94), "互不干擾", "FEF3C7", "B45309")

    card(
        slide,
        Inches(0.82),
        Inches(5.58),
        Inches(11.85),
        Inches(1.00),
        "設計原則",
        "Anritsu 與 Amarisoft 各自保留自己的 OpenClaw 與操作邏輯，避免互相干擾；LLM 推論則集中在同一台 DGX GB10，讓模型升級、監控與資源管理更簡單。",
        "2563EB",
    )

def add_anritsu_slide(slide):
    add_background(slide)
    header(
        slide,
        "Anritsu MT8000 環境",
        "這一組 OpenClaw Agent 只負責 Anritsu 測試流程，並使用 DGX GB10 上的 qwen3.5:35b。",
        "Anritsu",
        "DBEAFE",
        "1D4ED8",
    )

    card(slide, Inches(0.72), Inches(1.72), Inches(3.65), Inches(1.05), "使用者 / 測試工程師", "輸入測試目標、選擇情境、要求切換儀器。", "2563EB")
    card(slide, Inches(0.72), Inches(3.00), Inches(3.65), Inches(1.05), "OpenClaw Agent A", "解析指令、選擇步驟、調度 Anritsu 相關工具。", "16A34A")
    card(slide, Inches(0.72), Inches(4.28), Inches(3.65), Inches(1.05), "Anritsu MT8000", "執行測試、切換設定、回傳結果。", "D97706")

    connector(slide, Inches(4.38), Inches(2.24), Inches(5.22), Inches(2.24), "8EA2BE")
    connector(slide, Inches(4.38), Inches(3.52), Inches(5.22), Inches(3.52), "8EA2BE")

    card(slide, Inches(5.25), Inches(1.72), Inches(3.70), Inches(1.20), "模型配置", "model: qwen3.5:35b\n用途：Anritsu 測試判斷、流程編排、回覆生成。", "047857")
    card(slide, Inches(5.25), Inches(3.10), Inches(3.70), Inches(1.20), "LLM 角色", "負責理解操作意圖，例如切換測項、改變測試步驟或查詢目前狀態。", "1D4ED8")
    card(slide, Inches(5.25), Inches(4.48), Inches(3.70), Inches(1.20), "控制邊界", "只處理 Anritsu 這個環境的儀器，不直接碰 Amarisoft。", "B45309")

    card(
        slide,
        Inches(9.28),
        Inches(1.96),
        Inches(3.15),
        Inches(3.02),
        "Anritsu 流程建議",
        "1. 使用者發送需求\n2. OpenClaw Agent A 判斷是否要切換儀器\n3. 呼叫 Anritsu 工具或腳本\n4. 將結果送回使用者\n\n優點：\n- 風險隔離\n- 可獨立調整 prompt\n- 出錯時只影響單一環境",
        "2563EB",
    )

    badge(slide, Inches(0.84), Inches(5.72), Inches(1.72), "專用 Agent", "EFF6FF", "1D4ED8")
    badge(slide, Inches(2.68), Inches(5.72), Inches(1.78), "qwen3.5:35b", "ECFDF5", "047857")
    badge(slide, Inches(4.58), Inches(5.72), Inches(1.86), "Anritsu 專屬", "FEF3C7", "B45309")
    badge(slide, Inches(6.56), Inches(5.72), Inches(1.82), "儀器切換", "EEF2FF", "4338CA")


def add_amarisoft_slide(slide):
    add_background(slide)
    header(
        slide,
        "Amarisoft 環境",
        "另一組 OpenClaw Agent 專門負責 Amarisoft 測試流程，並使用 DGX GB10 上的 gemma4:12b。",
        "Amarisoft",
        "FEF3C7",
        "B45309",
    )

    card(slide, Inches(0.72), Inches(1.72), Inches(3.65), Inches(1.05), "使用者 / 測試工程師", "選擇 Amarisoft 測試情境、下達切換或執行命令。", "2563EB")
    card(slide, Inches(0.72), Inches(3.00), Inches(3.65), Inches(1.05), "OpenClaw Agent B", "解析指令、管理流程、調度 Amarisoft 工具。", "16A34A")
    card(slide, Inches(0.72), Inches(4.28), Inches(3.65), Inches(1.05), "Amarisoft", "執行控制、切換配置、輸出測試結果。", "D97706")

    connector(slide, Inches(4.38), Inches(2.24), Inches(5.22), Inches(2.24), "8EA2BE")
    connector(slide, Inches(4.38), Inches(3.52), Inches(5.22), Inches(3.52), "8EA2BE")

    card(slide, Inches(5.25), Inches(1.72), Inches(3.70), Inches(1.20), "模型配置", "model: gemma4:12b\n用途：Amarisoft 測試編排、操作回覆、狀態判斷。", "047857")
    card(slide, Inches(5.25), Inches(3.10), Inches(3.70), Inches(1.20), "LLM 角色", "負責理解 Amarisoft 相關操作意圖，並產生對應的控制流程。", "1D4ED8")
    card(slide, Inches(5.25), Inches(4.48), Inches(3.70), Inches(1.20), "控制邊界", "只處理 Amarisoft 的儀器與腳本，與 Anritsu 完全分離。", "B45309")

    card(
        slide,
        Inches(9.28),
        Inches(1.96),
        Inches(3.15),
        Inches(3.02),
        "Amarisoft 流程建議",
        "1. 使用者提出需求\n2. OpenClaw Agent B 判斷步驟與切換邏輯\n3. 呼叫 Amarisoft 工具或腳本\n4. 回傳執行結果\n\n優點：\n- 規則獨立\n- Prompt 可分開調整\n- 可用不同模型對應不同任務",
        "16A34A",
    )

    badge(slide, Inches(0.84), Inches(5.72), Inches(1.72), "專用 Agent", "EFF6FF", "1D4ED8")
    badge(slide, Inches(2.68), Inches(5.72), Inches(1.78), "gemma4:12b", "ECFDF5", "047857")
    badge(slide, Inches(4.58), Inches(5.72), Inches(1.86), "Amarisoft 專屬", "FEF3C7", "B45309")
    badge(slide, Inches(6.56), Inches(5.72), Inches(1.82), "流程隔離", "EEF2FF", "4338CA")


def add_ops_slide(slide):
    add_background(slide)
    header(
        slide,
        "部署與維運重點",
        "這樣的架構可以把測試控制與 LLM 推論分離，讓兩個環境共用算力，但仍保留各自的操作邏輯。",
        "維運",
        "ECFDF5",
        "047857",
    )

    card(slide, Inches(0.72), Inches(1.70), Inches(3.80), Inches(1.15), "OpenClaw 分工", "Anritsu 與 Amarisoft 各自有自己的 agent、工具與流程，不共享控制狀態。", "2563EB")
    card(slide, Inches(0.72), Inches(3.02), Inches(3.80), Inches(1.15), "Ollama 共用", "兩邊都指向同一台 DGX GB10，但使用不同模型名稱。", "16A34A")
    card(slide, Inches(0.72), Inches(4.34), Inches(3.80), Inches(1.15), "模型管理", "可集中更新、監控與資源分配，避免兩套 AI 各自漂移。", "D97706")

    card(slide, Inches(4.88), Inches(1.78), Inches(3.95), Inches(1.25), "建議的穩定做法", "把 OpenClaw 的模型名稱、base URL、timeout 與權限邊界都明確寫在設定檔。", "047857")
    card(slide, Inches(4.88), Inches(3.20), Inches(3.95), Inches(1.25), "風險控管", "若共用 LLM 造成延遲，可再加排隊、快取或模型分流，但控制層仍要保持獨立。", "1D4ED8")
    card(slide, Inches(4.88), Inches(4.62), Inches(3.95), Inches(1.25), "對外說法", "每個測試環境有自己的 AI 操作員，但後面的腦袋可以共用同一台推論伺服器。", "B45309")

    card(slide, Inches(9.10), Inches(1.92), Inches(3.15), Inches(3.10), "成功標準", "1. Anritsu 與 Amarisoft 可以獨立切換儀器\n2. 兩邊都能呼叫共用 Ollama\n3. 模型選用互不干擾\n4. 之後維護只要管一台 DGX GB10", "2563EB")

    badge(slide, Inches(9.38), Inches(5.36), Inches(2.50), "可擴充架構", "DBEAFE", "1D4ED8")


def build():
    global prs
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs.slides.add_slide(prs.slide_layouts[6]))
    add_overview_slide(prs.slides.add_slide(prs.slide_layouts[6]))
    add_anritsu_slide(prs.slides.add_slide(prs.slide_layouts[6]))
    add_amarisoft_slide(prs.slides.add_slide(prs.slide_layouts[6]))
    add_ops_slide(prs.slides.add_slide(prs.slide_layouts[6]))

    prs.save(OUT)
    print(f"Saved to {OUT}")


if __name__ == "__main__":
    build()
