from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "docs" / "km-modernization" / "KM_MODERNIZATION_WP0-WP13_ROADMAP.pptx"
BLUE = RGBColor(20, 55, 92)
TEAL = RGBColor(0, 120, 125)
INK = RGBColor(35, 42, 50)
MUTED = RGBColor(95, 105, 115)
LIGHT = RGBColor(237, 244, 247)
GREEN = RGBColor(33, 122, 78)
AMBER = RGBColor(168, 103, 0)


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(5)
        if align:
            p.alignment = align
    return box


def base(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = RGBColor(255, 255, 255)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(.25))
    band.fill.solid(); band.fill.fore_color.rgb = TEAL; band.line.fill.background()
    add_text(slide, title, .6, .55, 12, .6, 27, BLUE, True)
    if subtitle:
        add_text(slide, subtitle, .62, 1.15, 12, .35, 11, MUTED)
    return slide


def bullets(slide, items, x=.8, y=1.65, w=11.8, size=18, color=INK):
    add_text(slide, "\n".join("• " + item for item in items), x, y, w, 5.5, size, color)


def card(slide, title, body, x, y, w, h, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT; shape.line.color.rgb = color
    add_text(slide, title, x+.18, y+.15, w-.35, .35, 16, color, True)
    add_text(slide, body, x+.18, y+.62, w-.35, h-.75, 12, INK)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BLUE
    add_text(s, "Knowledge Base Modernization", .75, 1.25, 12, 1, 34, RGBColor(255,255,255), True)
    add_text(s, "WP0–WP13 implementation roadmap\nWP0/WP1 current delivery status and execution gates", .8, 2.45, 11, 1.2, 22, RGBColor(220,240,242))
    add_text(s, "2026-08-10  |  branch: agent/wp1-job-config-reliability", .8, 6.45, 11, .4, 12, RGBColor(210,220,230))

    s = base(prs, "簡報目的", "先建立可靠底座，再逐步導入正式資料契約、治理與 Agentic RAG")
    bullets(s, ["本簡報不是一次性重構清單，而是依 Gate 分階段交付。", "每個 WP 獨立分支／PR；前一個 Gate 未通過，不修改下一階段 production path。", "既有 Portal、chat、search、report upload/review、ingest 與 A2A dry-run 保持相容。", "任何未驗證的外部依賴，必須標示為 blocker，不以 HTTP 200 或 UI demo 取代驗收。"])

    s = base(prs, "目標架構與資料流")
    for i, (t, b) in enumerate([("Portal / OpenClaw", "使用者查詢、報告審核、受控工具"), ("AI KM API", "FastAPI contract、Auth、Audit、Job"), ("Job Layer", "Redis + Celery\nqueue / retry / state"), ("Knowledge Services", "Package validation\nVector / Graph / Time-series"), ("External Systems", "CSIT、Anritsu、Amarisoft\nA2A bridge 預設關閉")]):
        card(s, t, b, .45 + i*2.58, 2.2, 2.25, 1.65, TEAL if i in (1,3) else BLUE)
        if i < 4:
            add_text(s, "→", 2.72 + i*2.58, 2.75, .35, .35, 24, TEAL, True, PP_ALIGN.CENTER)
    add_text(s, "資料原則：Parser → Knowledge Package → Validation → Publish Ledger → Projection / Query", .9, 5.2, 11.5, .55, 19, BLUE, True, PP_ALIGN.CENTER)

    s = base(prs, "今天完成的 WP0", "FastAPI contract baseline，commit 2c46c834")
    bullets(s, ["新增 app/ FastAPI shell、/api/v1 health/version、統一 response/error envelope。", "X-Trace-ID middleware 與穩定 exception boundary；不回傳 stack、path、secret。", "保留 legacy /health、/search、/ws、report、ingest 與 A2A dry-run。", "測試證據：76 passed、compileall、Compose config、Vue build、release package、Portal smoke。", "未部署 production；PR #2 base=agent/km-modify-codex-plan，head=agent/wp0-fastapi-contract。"])

    s = base(prs, "今天完成的 WP1", "Job/config reliability baseline，commit 7cfa1d6e")
    bullets(s, ["JobConfig typed env：concurrency、TTL、timeout、retry、queue、Beat。", "JobStatus：queued / running / succeeded / failed / retrying / cancelled。", "保留 search/ingest，新增 default/document/indexing queue contract。", "X-Trace-ID → Celery headers；transient／non-retryable retry policy。", "Compose：named volume、可配置 host paths、restart policy、Beat opt-in、Neo4j healthcheck。", "驗證：83 passed；實際隔離 Docker 啟動、worker restart、Redis persistence、API idempotency 成功。"])

    s = base(prs, "WP1 實際驗證證據與發現")
    card(s, "Worker recovery", "隔離 Compose project\n重啟 celery_ingest_worker\nRedis key 保留\nworker 回復運作", .7, 1.7, 3.7, 2.1, GREEN)
    card(s, "Idempotency", "相同 Excel + identity\n第二次回傳相同 task_id\nduplicate=true\n避免重複攝入", 4.8, 1.7, 3.7, 2.1, GREEN)
    card(s, "根因修正", "發現 duplicate path 缺少 import\n發現 Neo4j healthcheck 硬編碼密碼與錯誤 PATH\n均已修正並推送", 8.9, 1.7, 3.7, 2.1, AMBER)
    add_text(s, "剩餘：CI/PR gate、長時間故障注入、正式部署環境 restore/recovery 演練", 1, 5.2, 11.3, .55, 19, AMBER, True, PP_ALIGN.CENTER)

    s = base(prs, "歷史規劃草稿（不作 v2.6 基準）", "本頁僅保留歷史內容；目前 Phase/WP 請以 v2.6 source index 為準")
    bullets(s, ["本頁原為舊版 WP2–WP4 草稿，不能用來判定目前 Phase 或 WP 進度。", "目前唯一有效規劃為 01_AI_KM_Phase規劃_v2.6.xlsx，正式對照見 docs/km-modernization/07-v2.6-source-index.md。", "v2.6 使用 Phase 1～5、WP0～WP13，WP10 拆為 WP10A／WP10B；不使用正式商業 Phase 0。", "原始規格、Evidence、PR、CI 與驗收證據必須分開核對，不因設計文件存在而增加程式完成率。"])

    s = base(prs, "Phase 1：WP5–WP9", "正式報告 workflow、時序資料、Graph、治理與 Portal")
    bullets(s, ["WP5 Report Publish Ledger：draft → validated → review → approved → publishing → published。", "WP6 TimescaleDB：iPerf／PHY／RF／KPI 明細與 summary；bulk ingest、trend、compare、ACL。", "WP7 Neo4j ontology：固定節點、canonical ID、MERGE、source/lineage、受控 query template。", "WP8 RBAC/Citation/Audit：deny-by-default、citation completeness、query/tool/upload/review/publish audit。", "WP9 Portal/OpenClaw MVP：upload → review → approve → publish → authorized search；A2A 仍隔離、disabled、dry-run。"])

    s = base(prs, "Phase 2–4：WP10A–WP12", "品質、受控 Agentic RAG 與 AI analysis")
    bullets(s, ["WP10A Document Intelligence：文字、圖片、Excel、表格專用 pipeline；OCR 按需，來源追蹤到 page/cell。", "WP10B Entity normalization：alias、master data、ontology mapping、Qdrant/Neo4j rebuild。", "WP11 Controlled Tools：CSIT、Vector、Graph、Time-Series、Automation；禁止自由 SQL/Cypher/SSH。", "WP12 Analysis：Root Cause、Benchmark、Similar Case、Recommendation；區分 evidence / inference / confidence。"])

    s = base(prs, "Phase 5：WP13 與 A2A 邊界")
    bullets(s, ["WP13 Multi-Agent／Predictive／Proactive workflow 只有在單 Agent 品質、資料量與 KPI 達標後啟動。", "KM 是中心 agent；Anritsu／Amarisoft 是受控階層子 agent。", "A2A bridge 預設 disabled，dry-run 不碰儀器；real execute 需明確授權、確認、timeout、audit 與停止條件。", "Automation query 與 execute 分離；高權限 execute 必須明確確認，不能由模型自行擴權。"])

    s = base(prs, "跨階段安全與相容性規則")
    bullets(s, ["不得硬編碼 token、密碼、host path；.env 不入 Git，log 不含 secret 或完整敏感 payload。", "新 API 以 additive contract 進入；legacy routes response shape 保持相容。", "資料庫各有責任：Qdrant=semantic、Neo4j=ontology/evidence、TimescaleDB=timeseries、CSIT=source of record。", "每次 migration 必須有 forward、rollback、backup/restore 證據。", "A2A、Portal、search、report、ingest 的既有功能不能因新 WP 被順便重構。"])

    s = base(prs, "Gate 與測試矩陣")
    bullets(s, ["G0：API contract、trace、exception、legacy smoke。", "G1：queue routing、retry/non-retry、timeout、worker restart、idempotency、config。", "G2：Knowledge Package validation 與 routing；invalid 不 mutation。", "G3–G5：Qdrant/CSIT/publish ledger/Timescale/Neo4j 與 Portal E2E。", "G6–G8：document quality、agent routing、citation、analysis evaluation。", "任何 Gate 未通過，不得以 UI demo、HTTP 200 或單一 happy path 宣稱完成。"])

    s = base(prs, "Rollback 與正式導入策略")
    bullets(s, ["程式碼：每個 WP 獨立 branch/PR；穩定 tag 作回復點，不直接 reset shared branch。", "資料：Neo4j dump、Qdrant snapshot、PostgreSQL dump、Redis/SQLite registry、uploads/assets、config/env、image manifest。", "部署：獨立 Compose project、container name、port、volume 做 UAT；未驗證版本不得寫正式資料。", "失敗時：停止新版本 → 保留 log/trace → 回復 image/config/code → 還原資料快照 → 執行 health、query、ingest、citation smoke。"])

    s = base(prs, "目前 GitHub 狀態與下一步")
    bullets(s, ["WP0 branch：agent/wp0-fastapi-contract，commit 2c46c834；PR #2 尚需 CI gate。", "WP1 branch：agent/wp1-job-config-reliability，最新 commit 7cfa1d6e，已推送。", "WP0 workflow 推送曾被 PAT 拒絕，原因是缺少 workflow scope；需具 workflow 權限的 credential 才能讓 Actions 執行。", "下一步：補 CI → 建立 WP1 PR → review/merge gate → WP2 Knowledge Package，不跨階段搶做資料庫重構。"])

    s = base(prs, "交付結論")
    add_text(s, "先可靠，再正式化資料；先治理，再放大 Agent 能力。", 1, 2.0, 11.3, .8, 30, BLUE, True, PP_ALIGN.CENTER)
    add_text(s, "WP0/WP1 已建立可驗證底座，但每個後續 WP 都必須以獨立 Gate、測試證據與 rollback 條件推進。", 1.2, 3.35, 10.9, 1, 20, INK, False, PP_ALIGN.CENTER)
    add_text(s, "產出檔案：docs/km-modernization/KM_MODERNIZATION_WP0-WP13_ROADMAP.pptx", 1, 6.1, 11.3, .4, 12, MUTED, False, PP_ALIGN.CENTER)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
