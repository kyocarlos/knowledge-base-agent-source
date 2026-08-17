#!/usr/bin/env python3
"""Reconcile the 17-slide W33 supervisor deck from its canonical JSON data."""

from __future__ import annotations

import argparse
import json
import re
import shutil
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PERCENT_PATTERN = re.compile(r"\d+(?:\.\d+)?\s*%")


@dataclass(frozen=True)
class ProgressMetrics:
    wp0: int | float
    wp1: int | float
    phase1: int | float
    program: int | float
    cutoff: str


def rounded_progress(total: int | float, count: int) -> int | float:
    value = round(total / count, 1)
    return int(value) if value.is_integer() else value


def load_metrics(path: Path) -> ProgressMetrics:
    data = json.loads(path.read_text(encoding="utf-8"))
    packages = data.get("work_packages")
    if not isinstance(packages, list) or not packages:
        raise SystemExit("W33 JSON must contain a non-empty work_packages list")

    progress_by_id: dict[str, int | float] = {}
    phase1_progress: list[int | float] = []
    all_progress: list[int | float] = []
    for package in packages:
        package_id = package.get("id")
        progress = package.get("progress")
        if not isinstance(package_id, str) or not isinstance(progress, (int, float)):
            raise SystemExit("each work package must contain string id and numeric progress")
        progress_by_id[package_id] = progress
        all_progress.append(progress)
        if package.get("phase") == 1:
            phase1_progress.append(progress)

    for required_id in ("WP0", "WP1"):
        if required_id not in progress_by_id:
            raise SystemExit(f"W33 JSON is missing {required_id}")
    if not phase1_progress:
        raise SystemExit("W33 JSON has no Phase 1 work packages")

    computed_phase1 = rounded_progress(sum(phase1_progress), len(phase1_progress))
    computed_program = rounded_progress(sum(all_progress), len(all_progress))
    declared_phase1 = data.get("phase_progress", {}).get("1")
    declared_program = data.get("program_progress")
    if declared_phase1 != computed_phase1:
        raise SystemExit(
            f"phase_progress mismatch: declared={declared_phase1} computed={computed_phase1}"
        )
    if declared_program != computed_program:
        raise SystemExit(
            f"program_progress mismatch: declared={declared_program} computed={computed_program}"
        )

    period = data.get("period", {})
    cutoff = period.get("cutoff")
    timezone = period.get("timezone")
    if not isinstance(cutoff, str) or not isinstance(timezone, str):
        raise SystemExit("W33 JSON period must contain cutoff and timezone")

    return ProgressMetrics(
        wp0=progress_by_id["WP0"],
        wp1=progress_by_id["WP1"],
        phase1=declared_phase1,
        program=declared_program,
        cutoff=f"{cutoff} {timezone}",
    )


def percent(value: int | float) -> str:
    return f"{value:g}%"


def replace_percentages(value: str, replacements: list[str]) -> str:
    index = 0

    def replace(_match: re.Match[str]) -> str:
        nonlocal index
        if index >= len(replacements):
            return _match.group(0)
        replacement = replacements[index]
        index += 1
        return replacement

    return PERCENT_PATTERN.sub(replace, value)


def reconcile_text(value: str, metrics: ProgressMetrics) -> str:
    if "不能宣稱" in value:
        return value
    has_wp0 = re.search(r"\bWP0\b", value, re.IGNORECASE) is not None
    has_wp1 = re.search(r"\bWP1\b", value, re.IGNORECASE) is not None
    if has_wp0 and has_wp1:
        return replace_percentages(value, [percent(metrics.wp0), percent(metrics.wp1)])
    if has_wp0:
        return replace_percentages(value, [percent(metrics.wp0)])
    if has_wp1:
        return replace_percentages(value, [percent(metrics.wp1)])
    if re.search(r"Phase\s*1.*進度", value, re.IGNORECASE):
        return replace_percentages(value, [percent(metrics.phase1)])
    if re.search(r"(?:全計畫|整體計畫|全案).*進度", value):
        return replace_percentages(value, [percent(metrics.program)])
    return value


def reconcile_paragraph(paragraph, metrics: ProgressMetrics) -> int:
    before = paragraph.text
    after = reconcile_text(before, metrics)
    if after == before:
        return 0
    if paragraph.runs:
        paragraph.runs[0].text = after
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = after
    return 1


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)


def reconcile_shape(shape, metrics: ProgressMetrics) -> int:
    changes = 0
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            changes += reconcile_paragraph(paragraph, metrics)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    changes += reconcile_paragraph(paragraph, metrics)
    return changes


def shape_text(shape) -> list[str]:
    values: list[str] = []
    if getattr(shape, "has_text_frame", False):
        values.extend(p.text for p in shape.text_frame.paragraphs if p.text.strip())
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                values.extend(p.text for p in cell.text_frame.paragraphs if p.text.strip())
    return values


def assert_displayed_metrics(rendered_text: str, metrics: ProgressMetrics) -> None:
    required = (percent(metrics.wp0), percent(metrics.wp1), percent(metrics.phase1))
    for value in required:
        if value not in rendered_text:
            raise SystemExit(f"required JSON-backed metric missing from deck: {value}")

    for line in rendered_text.splitlines():
        if "不能宣稱" in line:
            continue
        has_wp0 = re.search(r"\bWP0\b", line, re.IGNORECASE) is not None
        has_wp1 = re.search(r"\bWP1\b", line, re.IGNORECASE) is not None
        values = PERCENT_PATTERN.findall(line)
        if has_wp0 and has_wp1 and len(values) >= 2 and values[:2] != [percent(metrics.wp0), percent(metrics.wp1)]:
            raise SystemExit(f"combined WP metric line differs from W33 JSON: {line}")
        if has_wp0 and not has_wp1 and values and values[0] != percent(metrics.wp0):
            raise SystemExit(f"WP0 metric line differs from W33 JSON: {line}")
        if has_wp1 and not has_wp0 and values and values[0] != percent(metrics.wp1):
            raise SystemExit(f"WP1 metric line differs from W33 JSON: {line}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--data", required=True, type=Path)
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--text-output", required=True, type=Path)
    parser.add_argument("--qa-output", required=True, type=Path)
    args = parser.parse_args()

    metrics = load_metrics(args.data)
    presentation = Presentation(args.input)
    expected_slides = 17
    if len(presentation.slides) != expected_slides:
        raise SystemExit(f"expected {expected_slides} slides, found {len(presentation.slides)}")

    changed_paragraphs = 0
    all_text: list[str] = []
    out_of_bounds: list[dict[str, int]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in walk_shapes(slide.shapes):
            changed_paragraphs += reconcile_shape(shape, metrics)
            slide_lines.extend(shape_text(shape))
            if (
                min(shape.left, shape.top) < 0
                or shape.left + shape.width > presentation.slide_width
                or shape.top + shape.height > presentation.slide_height
            ):
                out_of_bounds.append({"slide": slide_number, "shape_id": shape.shape_id})
        all_text.append(f"## Slide {slide_number}\n" + "\n".join(slide_lines))

    rendered_text = "\n\n".join(all_text)
    assert_displayed_metrics(rendered_text, metrics)
    if out_of_bounds:
        raise SystemExit(f"shapes outside slide bounds: {out_of_bounds}")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.text_output.parent.mkdir(parents=True, exist_ok=True)
    args.qa_output.parent.mkdir(parents=True, exist_ok=True)
    if changed_paragraphs:
        presentation.save(args.output)
    else:
        shutil.copyfile(args.input, args.output)
    args.text_output.write_text(rendered_text + "\n", encoding="utf-8")
    args.qa_output.write_text(
        json.dumps(
            {
                "slides": expected_slides,
                "reconciliation": "passed",
                "metrics": {
                    "WP0": metrics.wp0,
                    "WP1": metrics.wp1,
                    "Phase 1": metrics.phase1,
                    "program": metrics.program,
                },
                "cutoff": metrics.cutoff,
                "out_of_bounds_shapes": 0,
                "metric_source": str(args.data),
                "program_metric_in_deck": "not displayed; canonical value is recorded in JSON and Markdown",
                "cutoff_in_deck": "not displayed; canonical cutoff is recorded in Evidence, JSON and Markdown",
                "render_check": "pending",
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
