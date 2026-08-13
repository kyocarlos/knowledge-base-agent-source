#!/usr/bin/env python3
"""Build an itemized WP0/WP1 change ledger inside the W33 Phase 1 deck."""

from __future__ import annotations

import argparse
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


COLORS = {
    "navy": "17324D", "teal": "087F8C", "green": "2E7D5B", "amber": "B16D00",
    "red": "A83B3B", "ink": "26333D", "muted": "667783", "light": "F4F7F8",
    "pale": "E7F1F2", "white": "FFFFFF", "line": "CDD9DE", "alt": "F9FBFB",
}
ROWS_PER_SLIDE = 8


WP0 = [
    ("新增", "正式 FastAPI 套件邊界", "建立 app/api、app/core、app/schemas，將平台 Contract 與 legacy 程式分層。", "app/__init__.py 等"),
    ("新增", "正式應用程式入口", "以 create_app() 工廠建立 production FastAPI app。", "app/main.py"),
    ("新增", "服務 Metadata", "設定 title、description、version，供 OpenAPI 與部署識別。", "app/main.py"),
    ("新增", "Typed AppSettings", "服務名、版本、環境、commit 改由 immutable settings 管理。", "app/core/config.py"),
    ("新增", "環境值安全驗證", "限制設定字元與長度；不合法值在啟動時顯性失敗。", "app/core/config.py"),
    ("新增", "Versioned Router", "建立 /api/v1 平台 API prefix，不改名既有 API。", "app/api/v1/router.py"),
    ("新增", "Health Endpoint", "新增 /api/v1/health，回傳 status/live/ready 與 trace。", "router.py"),
    ("新增", "Liveness Endpoint", "新增 /api/v1/health/live，供程序存活探測。", "router.py"),
    ("新增", "Readiness Endpoint", "新增 /api/v1/health/ready，保留後續依賴探測擴充點。", "router.py"),
    ("新增", "Version Endpoint", "新增 /api/v1/version，回傳 service/version/environment/commit。", "router.py"),
    ("新增", "成功回應 Envelope", "所有 v1 成功回應統一 data/error/trace_id 結構。", "schemas/common.py"),
    ("新增", "錯誤回應 Envelope", "錯誤統一 code/message，避免各 endpoint 自訂格式。", "schemas/common.py"),
    ("新增", "Health/Version Schema", "HealthData、VersionData 固定欄位與型別。", "schemas/common.py"),
    ("新增", "Schema 嚴格模式", "Pydantic extra=forbid，拒絕未定義欄位。", "schemas/common.py"),
    ("新增", "Trace ID 輸入驗證", "只接受 1～128 字元的安全 X-Trace-ID。", "app/core/trace.py"),
    ("新增", "Trace ID 自動產生", "缺失或非法 header 時產生 UUID correlation ID。", "app/core/trace.py"),
    ("新增", "Request Trace Context", "trace_id 同步寫入 request.state 與 ContextVar。", "app/core/trace.py"),
    ("新增", "Trace Response Header", "所有新舊路由回應加入 X-Trace-ID。", "app/core/trace.py"),
    ("新增", "HTTP Trace Logging", "記錄 trace、method、path、status，支援請求追查。", "app/core/trace.py"),
    ("新增", "Logging Context Filter", "現有 logging handler 自動帶入目前 trace_id。", "app/core/logging.py"),
    ("新增", "422 錯誤映射", "v1 validation error 固定為 validation_error，不暴露原始 payload。", "app/core/exceptions.py"),
    ("新增", "HTTP 錯誤映射", "v1 HTTPException 固定 http_<status> code 與 trace。", "app/core/exceptions.py"),
    ("新增", "500 秘密安全處理", "記錄 exception type，但 client 只收到 Internal server error。", "app/core/exceptions.py"),
    ("新增", "Security Context 邊界", "提供匿名安全 context；WP0 不自行發明 Identity Provider。", "app/core/security.py"),
    ("保留", "Legacy Lifespan", "沿用既有 startup/shutdown context，避免資源初始化回歸。", "app/main.py"),
    ("保留", "Legacy Routes", "掛回 Portal/chat/search/report/ingest/A2A 等既有 routes。", "app/main.py"),
    ("保留", "Legacy Middleware", "沿用既有 CORS 與 middleware 設定。", "app/main.py"),
    ("修改", "Framework Route 去重", "排除舊 openapi/docs/redoc，避免新舊 framework route 衝突。", "app/main.py"),
    ("修改", "Docker 正式入口", "Image 複製 app/，CMD 由 src.web_api:app 改為 app.main:app。", "Dockerfile"),
    ("修改", "Compose／啟停入口", "Compose command 與 start/stop process pattern 改用 app.main:app。", "docker-compose.yml、start.sh"),
    ("修改", "Release Package 入口", "Release bundle 納入 app/；runtime Uvicorn 與 frontend build 路徑校正。", "release/build_release.sh"),
    ("新增", "WP0 測試與 CI 基線", "新增 contract/architecture/legacy tests、pytest.ini、dev requirements、README 指令。", "3 test files；2c46c834"),
]


WP1 = [
    ("新增", "Canonical JobStatus", "固定 queued/running/succeeded/failed/retrying/cancelled 六種狀態。", "app/core/job_config.py"),
    ("新增", "Immutable JobConfig", "背景任務參數收斂為 frozen dataclass，不再散落硬編碼。", "job_config.py"),
    ("新增", "Concurrency 設定", "KB_MAX_CONCURRENT_PROCESSING 控制 worker concurrency。", "job_config.py、Compose"),
    ("新增", "Lock／Result TTL", "processing lock 與 result expiration 改為環境可配置。", "job_config.py"),
    ("新增", "Soft／Hard Timeout", "任務 soft limit 與 hard limit 改為環境可配置。", "job_config.py"),
    ("新增", "Retry Policy 設定", "最大重試次數與 countdown 改為環境可配置。", "job_config.py"),
    ("新增", "Queue 名稱設定", "default/document/indexing queues 統一由 JobConfig 提供。", "job_config.py"),
    ("新增", "Beat Enable 設定", "提供 KB_CELERY_BEAT_ENABLED 布林設定及驗證。", "job_config.py"),
    ("新增", "設定型別驗證", "整數最小值與布林字串不合法時立即失敗。", "job_config.py"),
    ("新增", "Celery Trace Headers", "celery_headers() 只傳遞 trace_id，不改 task payload。", "job_config.py"),
    ("修改", "Celery Default Queue", "task_default_queue 由 JobConfig 控制。", "src/web_api/tasks.py"),
    ("修改", "Queue 宣告", "顯式宣告 default/document/indexing/search/ingest queues。", "tasks.py"),
    ("修改", "Search Routing", "search_task 與 watch_folder_scan 固定送至 search queue。", "tasks.py"),
    ("修正", "Ingest Routing", "ingest_task 從錯誤的 search queue 改送 ingest queue。", "tasks.py"),
    ("修改", "Result Expiration", "Celery result_expires 採 JobConfig result TTL。", "tasks.py"),
    ("修改", "Worker Concurrency", "Celery worker_concurrency 採 JobConfig。", "tasks.py"),
    ("新增", "Lost Worker 保護", "啟用 task_acks_late 與 task_reject_on_worker_lost。", "tasks.py"),
    ("新增", "Canonical State Mapping", "converting/completed/queued 等 legacy 狀態映射為 job_status。", "tasks.py"),
    ("保留", "Legacy Status 相容", "保留原 status/progress/status_text/step，只新增 job_status。", "tasks.py"),
    ("新增", "Search Trace Propagation", "HTTP X-Trace-ID 經 Celery header 送入 search task。", "src/web_api/__init__.py"),
    ("新增", "Upload Trace Propagation", "Upload/ingest 提交時將 trace_id 傳入 ingest task。", "src/web_api/__init__.py"),
    ("新增", "Review Trace Propagation", "報告 approve 後排入 ingest 時傳遞 trace_id。", "report_routes.py"),
    ("新增", "Task State Trace", "攝入 state 保存 trace_id 與 celery_task_id。", "tasks.py"),
    ("新增", "Failure Trace Log", "搜尋失敗 log 加入跨程序 trace_id。", "tasks.py"),
    ("修改", "Ingest Max Retries", "ingest_file_task max_retries 由 JobConfig 控制。", "tasks.py"),
    ("修改", "Search Retry Countdown", "search self.retry countdown 不再固定 5 秒。", "tasks.py"),
    ("新增", "Error Classification", "輸入錯誤不重試；Connection/Timeout 視為暫時性基礎設施錯誤。", "job_config.py、retry tests"),
    ("保留", "Idempotent Ingest", "沿用 document hash、Idempotency-Key、registry duplicate/conflict 機制。", "ingest_conflict_protection.py"),
    ("修改", "可配置部署根目錄", "config/data/upload mounts 改用 KB_CONFIG_ROOT/KB_DATA_ROOT/KB_UPLOAD_ROOT。", "docker-compose.yml"),
    ("新增", "Service Restart Policy", "Redis/PostgreSQL/Neo4j/web/search/ingest 設 unless-stopped。", "docker-compose.yml"),
    ("保留", "專用 Search Worker", "search worker 固定 -Q search 並納入健康檢查。", "docker-compose.yml"),
    ("保留", "專用 Ingest Worker", "ingest worker 固定 -Q ingest，concurrency=1。", "docker-compose.yml"),
    ("保留", "Celery Beat", "Beat 獨立 service，納入 scheduler Gate。", "docker-compose.yml"),
    ("新增", "WP1 Unit Tests", "新增 JobConfig、queue/status contract、retry/non-retry tests。", "3 WP1 test files"),
    ("新增", "WP1 CI Gate", "backend/frontend/repository hygiene workflow，可手動執行。", "wp1-job-reliability.yml"),
    ("新增", "Pre-WP01 Checkpoint", "備份 source/config/data、dirty patch 與精確 application image。", "pre_wp01_backup.py"),
    ("新增", "Neo4j Logical Export", "使用 APOC stream 產生可驗證 Cypher export。", "pre_wp01_backup.py"),
    ("新增", "Qdrant Snapshot", "逐 collection 建立 snapshot 並下載保存。", "pre_wp01_backup.py"),
    ("新增", "PostgreSQL Dump", "report registry 使用 custom-format pg_dump。", "pre_wp01_backup.py"),
    ("新增", "Redis／SQLite Backup", "保存 Redis archive，SQLite 使用 online backup API。", "pre_wp01_backup.py"),
    ("新增", "Checkpoint Integrity", "產生 checkpoint.json、SHA256SUMS、inspect 與 rollback.env。", "pre_wp01_backup.py"),
    ("新增", "Rollback Dry-run", "未加 --execute 時只驗證 manifest/SHA/image，不操作服務。", "rollback_pre_wp01.py"),
    ("新增", "Production Confirmation", "正式回退需 --execute 與 PRE_WP01_ROLLBACK 雙重確認。", "rollback_pre_wp01.py"),
    ("新增", "資料安全邊界", "Application rollback 不刪 volume，也不隱式執行 Level 2 data restore。", "rollback_pre_wp01.py"),
    ("新增", "Shadow Rollback Drill", "隔離 project/network/volume；驗證 200→503→rollback 200。", "drill_pre_wp01_rollback.py"),
    ("新增", "Candidate Gate", "隔離依賴驗證 legacy/v1/Agent auth、web、workers、Beat 與清理。", "drill_wp01_candidate.py"),
    ("修正", "Candidate Search Probe", "Gate 補入 /search POST 與 X-Trace-ID，攔截 Request 注入回歸。", "drill_wp01_candidate.py"),
    ("修改", "安全生命週期工具", "restart_kb 拆分只讀 status、app-only restart、checkpoint deploy 與自動 rollback。", "restart_kb.sh；8 tests"),
]


def rgb(name: str) -> RGBColor:
    return RGBColor.from_string(COLORS[name])


def textbox(slide, text, x, y, w, h, size=11, color="ink", bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text; p.alignment = align
    p.font.name = "Microsoft JhengHei"; p.font.size = Pt(size)
    p.font.bold = bold; p.font.color.rgb = rgb(color)
    return box


def rect(slide, x, y, w, h, fill, line="line"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line); shape.line.width = Pt(0.7)
    return shape


def add_ledger_slide(prs, wp, page, page_count, rows, page_number):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    textbox(slide, f"{wp} 修改內容（{page}/{page_count}）", 0.6, 0.4, 12.0, 0.48, 24, "navy", True)
    subtitle = "每一列代表一項實際新增、修改、修正或保留相容性的功能／機制"
    textbox(slide, subtitle, 0.62, 0.96, 11.8, 0.26, 11, "muted")

    columns = [(0.65, 0.58, "編號"), (1.23, 0.67, "動作"), (1.90, 2.25, "功能／機制"),
               (4.15, 5.45, "實際改動"), (9.60, 3.05, "主要檔案／證據")]
    header_y, header_h = 1.36, 0.42
    for x, w, label in columns:
        rect(slide, x, header_y, w, header_h, "navy", "white")
        textbox(slide, label, x + 0.03, header_y + 0.02, w - 0.06, header_h - 0.04,
                10.5, "white", True, PP_ALIGN.CENTER)

    row_y, row_h = 1.78, 0.57
    action_color = {"新增": "teal", "修改": "amber", "修正": "red", "保留": "green"}
    for idx, (action, feature, change, evidence) in enumerate(rows):
        y = row_y + idx * row_h
        fill = "white" if idx % 2 == 0 else "alt"
        for x, w, _ in columns:
            rect(slide, x, y, w, row_h, fill)
        item_no = (page - 1) * ROWS_PER_SLIDE + idx + 1
        textbox(slide, f"{wp}-{item_no:02d}", 0.67, y + 0.02, 0.54, row_h - 0.04, 9.2, "navy", True, PP_ALIGN.CENTER)
        textbox(slide, action, 1.26, y + 0.02, 0.61, row_h - 0.04, 9.5,
                action_color[action], True, PP_ALIGN.CENTER)
        textbox(slide, feature, 1.96, y + 0.02, 2.13, row_h - 0.04, 10.3, "ink", True)
        textbox(slide, change, 4.21, y + 0.02, 5.33, row_h - 0.04, 9.8, "ink")
        textbox(slide, evidence, 9.66, y + 0.02, 2.93, row_h - 0.04, 9.2, "muted")

    start = (page - 1) * ROWS_PER_SLIDE + 1
    end = start + len(rows) - 1
    color = "teal" if wp == "WP0" else "green"
    rect(slide, 0.65, 6.46, 12.0, 0.38, "pale")
    textbox(slide, f"本頁 {wp}-{start:02d}～{wp}-{end:02d}｜{wp} 合計 {len(WP0) if wp == 'WP0' else len(WP1)} 條",
            0.85, 6.53, 11.6, 0.18, 9.5, color, True, PP_ALIGN.CENTER)
    textbox(slide, str(page_number), 12.2, 7.02, 0.5, 0.22, 10, "muted", False, PP_ALIGN.RIGHT)
    return slide


def move_last_slides_after(prs, count, after_index):
    ids = prs.slides._sldIdLst
    added = list(ids)[-count:]
    for item in added: ids.remove(item)
    for offset, item in enumerate(added, 1): ids.insert(after_index + offset, item)


def build(source: Path, output: Path):
    prs = Presentation(source)
    if len(prs.slides) != 7:
        raise ValueError(f"來源必須是原始 7 頁簡報，實際 {len(prs.slides)} 頁")
    total_detail = math.ceil(len(WP0) / ROWS_PER_SLIDE) + math.ceil(len(WP1) / ROWS_PER_SLIDE)
    total_slides = 7 + total_detail
    added = []
    page_number = 6
    for wp, items in (("WP0", WP0), ("WP1", WP1)):
        count = math.ceil(len(items) / ROWS_PER_SLIDE)
        for page in range(1, count + 1):
            rows = items[(page - 1) * ROWS_PER_SLIDE:page * ROWS_PER_SLIDE]
            added.append(add_ledger_slide(prs, wp, page, count, rows, page_number))
            page_number += 1
    move_last_slides_after(prs, len(added), 4)
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text") and shape.text.strip() == "1 / 7":
            p = shape.text_frame.paragraphs[0]
            p.text = f"1 / {total_slides}"; p.font.name = "Aptos"; p.font.size = Pt(10)
            p.font.color.rgb = RGBColor.from_string("AFC1CD"); p.alignment = PP_ALIGN.RIGHT
    prs.save(output)
    return total_slides


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()
    print(f"generated {build(args.source, args.output)} slides: {args.output}")


if __name__ == "__main__":
    main()
