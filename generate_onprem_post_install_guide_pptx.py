from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("<project-root>/knowledge-base/onprem_post_install_connection_guide.pptx")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0.0, radius=False):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def set_text(frame, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Sans CJK TC"):
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)


def textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, text, size, color, bold=bold, align=align)
    return box


def badge(slide, left, top, width, text, fill, color):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        Inches(0.30),
        fill=fill,
        line=fill,
        radius=True,
    )
    set_text(shape.text_frame, text, 8.5, color, align=PP_ALIGN.CENTER)
    return shape


def card(slide, left, top, width, height, title, body, accent, fill="FFFFFF"):
    shape = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=fill,
        line="CBD5E1",
        radius=True,
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(7)

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Noto Sans CJK TC"
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = rgb(accent)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(3)
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Noto Sans CJK TC"
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = rgb("475569")
    return shape


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    for left, top, width, height, color in [
        (Inches(-0.4), Inches(-0.3), Inches(3.8), Inches(3.8), "DBEAFE"),
        (Inches(9.3), Inches(-0.3), Inches(4.3), Inches(3.8), "D1FAE5"),
        (Inches(8.8), Inches(4.8), Inches(4.4), Inches(2.7), "FEF3C7"),
    ]:
        deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        deco.fill.solid()
        deco.fill.fore_color.rgb = rgb(color)
        deco.fill.transparency = 0.80
        deco.line.fill.background()

    border = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.22),
        Inches(0.18),
        Inches(12.89),
        Inches(7.14),
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    border.fill.transparency = 1


def header(slide, title, subtitle, badge_text, badge_fill, badge_color):
    textbox(slide, Inches(0.52), Inches(0.32), Inches(4.0), Inches(0.20), "Knowledge Base / On-Prem Guide", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.56), Inches(9.3), Inches(0.42), title, 24, "0F172A", bold=True)
    textbox(slide, Inches(0.52), Inches(0.98), Inches(10.0), Inches(0.34), subtitle, 12, "475569")
    badge(slide, Inches(10.72), Inches(0.50), Inches(1.80), badge_text, badge_fill, badge_color)


def add_cover(slide):
    add_background(slide)
    header(
        slide,
        "安裝完成後，還要做哪些設定才會連線成功？",
        "這份簡報專門給不熟系統的人看：安裝包裝好後，哪些已經自動完成，哪些還要再確認一次。",
        "安裝後設定",
        "EFF6FF",
        "1D4ED8",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(3.90), Inches(1.38), "你會看到什麼", "安裝包解壓後，系統會先裝好大部分內容，但仍需要確認 OpenClaw 連線、資料來源與驗證步驟。", "2563EB")
    card(slide, Inches(0.72), Inches(3.30), Inches(3.90), Inches(1.38), "這份簡報的目的", "把「安裝完成後要做什麼」整理成簡單步驟，讓客戶可以照著做，不需要懂程式碼。", "16A34A")
    card(slide, Inches(0.72), Inches(4.82), Inches(3.90), Inches(1.38), "最重要的一句話", "KB 安裝包會把系統架起來，但若要真正連線成功，還要確認 OpenClaw gateway、資料檔案與驗證流程。", "D97706")

    card(slide, Inches(5.05), Inches(1.95), Inches(7.05), Inches(2.18), "本次簡報的主題", "1. 安裝包已自動完成哪些事\n2. 安裝後還需要確認哪些設定\n3. 如果畫面還是顯示未連線，要怎麼手動修正\n4. 如果要把資料搬進去，raw 檔案要放哪裡\n5. 怎麼判斷系統真的好了", "047857")

    flow = [
        (0.86, "安裝包", "先把系統裝起來", "2563EB"),
        (3.00, "OpenClaw", "確認 gateway 連線", "16A34A"),
        (5.14, "host nginx", "可選，做外部網址", "D97706"),
        (7.28, "raw 資料", "把文件放進去", "1D4ED8"),
        (9.42, "驗證", "聊天與搜尋測試", "047857"),
    ]
    for x, t, d, a in flow:
        card(slide, Inches(x), Inches(5.30), Inches(1.92), Inches(0.88), t, d, a)
    for x1, x2 in [(2.78, 3.00), (4.92, 5.14), (7.06, 7.28), (9.20, 9.42)]:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(5.74), Inches(x2), Inches(5.74))
        line.line.color.rgb = rgb("94A3B8")
        line.line.width = Pt(1.3)


def add_autofill_slide(slide):
    add_background(slide)
    header(
        slide,
        "第一件事：哪些已經自動完成？",
        "先知道安裝包幫你做了什麼，這樣比較好判斷後面還缺什麼。",
        "已自動",
        "ECFDF5",
        "047857",
    )

    left_items = [
        ("1", "KB 服務本體", "Docker Compose、Web、Neo4j、Redis、Qdrant 都會起來。", "2563EB"),
        ("2", "OpenClaw 身分", "installer 會沿用主機上的 identity，避免金鑰空白。", "16A34A"),
        ("3", "聊天預設值", "避免舊的 queue / lock 讓聊天一開始就卡住。", "D97706"),
        ("4", "資料夾結構", "會先建立 raw / processed / assets / uploads。", "1D4ED8"),
    ]
    y = 1.72
    for idx, (num, title, body, accent) in enumerate(left_items):
        badge(slide, Inches(0.78), Inches(y + idx * 0.90), Inches(0.42), num, "EFF6FF", accent)
        card(slide, Inches(1.30), Inches(y - 0.04 + idx * 0.90), Inches(5.90), Inches(0.73), title, body, accent)

    card(slide, Inches(7.55), Inches(1.82), Inches(4.95), Inches(1.36), "這代表什麼", "大部分系統都已經準備好，不需要你手動安裝程式碼或自己 build。", "2563EB")
    card(slide, Inches(7.55), Inches(3.38), Inches(4.95), Inches(1.36), "但還要確認", "OpenClaw gateway 是否真的在聽正確的 port，以及你是否要對外開啟 nginx 站台。", "16A34A")
    card(slide, Inches(7.55), Inches(4.94), Inches(4.95), Inches(1.36), "客戶最常誤會的地方", "以為安裝包裝好就不用再確認。其實最後還是要做一次連線驗證。", "D97706")


def add_connect_slide(slide):
    add_background(slide)
    header(
        slide,
        "第二件事：讓 KB 連到 OpenClaw",
        "如果聊天頁顯示「未連線」，通常不是系統壞掉，而是 gateway 位址還沒對上。",
        "連線設定",
        "EFF6FF",
        "1D4ED8",
    )

    card(slide, Inches(0.72), Inches(1.72), Inches(4.75), Inches(1.10), "同一台主機時", "OpenClaw gateway 應該使用本機 IP，port 用 18790。", "2563EB")
    card(slide, Inches(0.72), Inches(3.00), Inches(4.75), Inches(1.10), "正確的樣子", "ws://<主機IP>:18790/ws", "16A34A")
    card(slide, Inches(0.72), Inches(4.28), Inches(4.75), Inches(1.10), "如果還沒連上", "把 `.env` 與 `install-state.env` 裡的 gateway 設定改成主機 IP + 18790。", "D97706")

    card(slide, Inches(5.80), Inches(1.72), Inches(6.00), Inches(1.28), "你可以直接記這句", "KB 安裝包預設會幫你接上 OpenClaw，但如果畫面還是未連線，就把 gateway 改成本機 IP 與 18790。", "1D4ED8")
    card(slide, Inches(5.80), Inches(3.22), Inches(6.00), Inches(1.28), "安裝後要看的檔案", ".env、install-state.env、runtime/openclaw/identity/device.json。", "047857")
    card(slide, Inches(5.80), Inches(4.72), Inches(6.00), Inches(1.28), "判斷成功的訊號", "chat.html 左上變成「已連線」，而且可以送出問題。", "D97706")


def add_nginx_slide(slide):
    add_background(slide)
    header(
        slide,
        "第三件事：OpenClaw 主機 nginx 要不要開？",
        "這一步是選配，不是每個客戶都要。只有想讓外部入口是 https://<host>:18789 時才需要。",
        "選配",
        "FEF3C7",
        "B45309",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(4.30), Inches(1.12), "如果不要對外網址", "可以直接用 KB 的 18443 網站，不必改主機 nginx。", "2563EB")
    card(slide, Inches(0.72), Inches(3.06), Inches(4.30), Inches(1.12), "如果要對外網址", "安裝時加上 --configure-openclaw-nginx，系統才會幫你寫 OpenClaw 的 nginx 站台。", "16A34A")
    card(slide, Inches(0.72), Inches(4.34), Inches(4.30), Inches(1.12), "注意", "主機上要先有 nginx，且要有權限寫 /etc/nginx。", "D97706")

    card(slide, Inches(5.35), Inches(1.82), Inches(6.55), Inches(1.45), "這個功能的用途", "讓 OpenClaw 的外部入口變成 https://<主機IP>:18789，方便客戶直接開網址使用。", "1D4ED8")
    card(slide, Inches(5.35), Inches(3.45), Inches(6.55), Inches(1.45), "如果不啟用會怎樣", "系統還是能跑，只是沒有那個對外入口；KB 內部頁面仍可正常使用。", "047857")
    card(slide, Inches(5.35), Inches(5.08), Inches(6.55), Inches(1.12), "一句話記法", "這是加分功能，不是系統必需功能。", "B45309")


def add_data_slide(slide):
    add_background(slide)
    header(
        slide,
        "第四件事：把文件放進 raw，再做 ingest",
        "如果客戶要用自己的文件回答問題，就要把文件放進 raw 目錄或透過上傳頁處理。",
        "資料",
        "ECFDF5",
        "047857",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(4.25), Inches(1.15), "原始文件放哪裡", "<onprem-root>/knowledge-base-onprem/app/data/raw", "2563EB")
    card(slide, Inches(0.72), Inches(3.08), Inches(4.25), Inches(1.15), "處理後內容", "系統會自動產生 processed、assets、uploads。", "16A34A")
    card(slide, Inches(0.72), Inches(4.38), Inches(4.25), Inches(1.15), "上傳後會做什麼", "轉文字、切段、寫 Neo4j、寫 Qdrant。", "D97706")

    card(slide, Inches(5.35), Inches(1.82), Inches(6.55), Inches(1.45), "如果不放資料會怎樣", "系統可以先起來，但聊天時可能只能回答通用內容，無法引用你的報告。", "1D4ED8")
    card(slide, Inches(5.35), Inches(3.45), Inches(6.55), Inches(1.45), "如果要測攝入", "把原始系統 data/raw 的檔案複製到新機器的 raw，然後再做 ingest。", "047857")
    card(slide, Inches(5.35), Inches(5.08), Inches(6.55), Inches(1.12), "最終目的", "讓客戶能用自己的文件去問小幫手。", "2563EB")


def add_verify_slide(slide):
    add_background(slide)
    header(
        slide,
        "第五件事：怎樣算真的成功連線？",
        "安裝完成後，要用最簡單的方法檢查三件事：服務有沒有起來、OpenClaw 有沒有連上、聊天有沒有回覆。",
        "驗證",
        "EFF6FF",
        "1D4ED8",
    )

    checks = [
        ("1", "看服務狀態", "docker compose ps 看到 web、nginx、redis、neo4j、qdrant 都是 Up。", "2563EB"),
        ("2", "開聊天頁", "https://<主機IP>:18443/chat.html", "16A34A"),
        ("3", "確認已連線", "畫面左上應該顯示「已連線」，不是「未連線」。", "D97706"),
        ("4", "送一個簡單問題", "例如「你在嘛？」或「今天天氣如何？」。", "1D4ED8"),
        ("5", "看回覆", "如果能回覆，代表整條鏈路已通。", "047857"),
    ]
    y = 1.72
    for idx, (num, title, body, accent) in enumerate(checks):
        badge(slide, Inches(0.78), Inches(y + idx * 0.82), Inches(0.42), num, "EFF6FF", accent)
        card(slide, Inches(1.30), Inches(y - 0.04 + idx * 0.82), Inches(5.95), Inches(0.68), title, body, accent)

    card(slide, Inches(7.55), Inches(1.82), Inches(4.95), Inches(1.28), "如果還是未連線", "通常是 gateway 位址不對，或 OpenClaw gateway 還沒有真正啟動。", "D97706")
    card(slide, Inches(7.55), Inches(3.30), Inches(4.95), Inches(1.28), "如果有連線但沒回覆", "再看是否有資料、或是否還卡在舊的 queue / lock。", "16A34A")
    card(slide, Inches(7.55), Inches(4.78), Inches(4.95), Inches(1.28), "最簡單的成功標準", "能開頁面、能看到已連線、能送出問題、能收到回答。", "2563EB")


def add_troubleshoot_slide(slide):
    add_background(slide)
    header(
        slide,
        "常見問題：如果還是連不上怎麼辦？",
        "這一頁給現場操作的人看，遇到問題時先從最容易的地方檢查。",
        "排錯",
        "FEE2E2",
        "B91C1C",
    )

    card(slide, Inches(0.72), Inches(1.78), Inches(3.95), Inches(1.15), "看到未連線", "先檢查 gateway 是否是本機 IP + 18790。", "B91C1C")
    card(slide, Inches(0.72), Inches(3.08), Inches(3.95), Inches(1.15), "聊天沒回覆", "看 WebSocket log，有沒有連到 OpenClaw gateway。", "D97706")
    card(slide, Inches(0.72), Inches(4.38), Inches(3.95), Inches(1.15), "要對外網址", "啟用 host nginx 選項，並確認 nginx 已存在。", "16A34A")

    card(slide, Inches(4.95), Inches(1.92), Inches(7.55), Inches(1.35), "最常見的原因", "安裝包把 OpenClaw gateway 指到錯的 port，或使用了 127.0.0.1。", "2563EB")
    card(slide, Inches(4.95), Inches(3.48), Inches(7.55), Inches(1.35), "最快的人工解法", "把 `.env` 與 `install-state.env` 內的 gateway 改成主機 IP + 18790，然後重啟 web / nginx。", "047857")
    card(slide, Inches(4.95), Inches(5.04), Inches(7.55), Inches(1.35), "給客戶的話術", "先不要改太多，只要照檢查順序做：服務 -> 連線 -> 資料 -> 驗證。", "1D4ED8")


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = [
    add_cover,
    add_autofill_slide,
    add_connect_slide,
    add_nginx_slide,
    add_data_slide,
    add_verify_slide,
    add_troubleshoot_slide,
]

for builder in slides:
    builder(prs.slides.add_slide(prs.slide_layouts[6]))

prs.save(OUT)
print(f"saved {OUT}")
