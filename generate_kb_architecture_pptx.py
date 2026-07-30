from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


OUT = Path("<project-root>/knowledge-base/kb_architecture_slide.pptx")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", "").upper())


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0.0, radius=False):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def set_text(frame, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Sans CJK TC"):
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return frame


def textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    set_text(box.text_frame, text, size, color, bold=bold, align=align)
    return box


def add_badge(slide, left, top, width, label, fill, text_color="FFFFFF"):
    badge = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        Inches(0.28),
        fill=fill,
        radius=True,
        transparency=0.10,
    )
    set_text(badge.text_frame, label, 8.5, text_color, align=PP_ALIGN.CENTER)
    return badge


def add_card(slide, left, top, width, height, title, subtitle, fill, accent):
    card = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill=fill,
        line="CBD5E1",
        radius=True,
    )
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(7)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = rgb(accent)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.name = "Noto Sans CJK TC"
    r2.font.size = Pt(10.2)
    r2.font.color.rgb = rgb("334155")
    return card


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    for left, top, width, height, color in [
        (Inches(-0.4), Inches(-0.3), Inches(3.8), Inches(3.8), "DBEAFE"),
        (Inches(9.3), Inches(-0.3), Inches(4.3), Inches(3.8), "D1FAE5"),
        (Inches(8.7), Inches(4.8), Inches(4.4), Inches(2.7), "FEF3C7"),
    ]:
        ellipse = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        ellipse.fill.solid()
        ellipse.fill.fore_color.rgb = rgb(color)
        ellipse.fill.transparency = 0.76
        ellipse.line.fill.background()

    border = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.22),
        Inches(0.18),
        Inches(12.89),
        Inches(7.14),
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    border.fill.transparency = 1


def connector(slide, x1, y1, x2, y2, color="8EA2BE"):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(2)
    return line


def add_cover(slide):
    add_background(slide)
    textbox(slide, Inches(0.6), Inches(0.42), Inches(4.2), Inches(0.22), "Knowledge Base System", 11, "2563EB")
    textbox(slide, Inches(0.6), Inches(0.72), Inches(7.2), Inches(0.8), "知識庫系統架構", 28, "0F172A", bold=True)
    textbox(slide, Inches(0.62), Inches(1.62), Inches(8.2), Inches(0.78), "本簡報說明知識庫系統的主要組成、資料流向與各元件職責，適合用於對客戶進行架構介紹。", 15, "475569")

    add_badge(slide, Inches(0.62), Inches(2.30), Inches(1.24), "簡化版", "DBEAFE", "1D4ED8")
    add_badge(slide, Inches(1.94), Inches(2.30), Inches(1.56), "架構說明", "D1FAE5", "047857")
    add_badge(slide, Inches(3.60), Inches(2.30), Inches(1.95), "正式簡報版", "FEF3C7", "B45309")

    note = add_card(
        slide,
        Inches(7.55),
        Inches(0.88),
        Inches(4.95),
        Inches(2.45),
        "簡報重點",
        "只講系統怎麼串起來，不講程式細節。",
        fill="FFFFFF",
        accent="2563EB",
    )
    tf = note.text_frame
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    r = p.add_run()
    r.text = "1. 使用者透過瀏覽器進入系統，使用聊天與查詢功能。"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(11)
    r.font.color.rgb = rgb("334155")
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "2. 前端頁面負責提供統一入口，並將請求送往後端。"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(11)
    r.font.color.rgb = rgb("334155")
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "3. 後端根據問題內容查詢知識庫資料，並整理成可閱讀的回答。"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(11)
    r.font.color.rgb = rgb("334155")
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "4. 背景任務則負責文件處理、索引更新與資料同步。"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(11)
    r.font.color.rgb = rgb("334155")

    flow_cards = [
        ("前端入口", "小幫手卡片盒與聊天入口", "提供統一的使用者操作介面"),
        ("後端核心", "知識庫查詢與回答", "負責理解問題、檢索資料與生成回覆"),
        ("資料庫", "Neo4j 與 Qdrant", "分別管理關聯脈絡與內容召回"),
    ]
    xs = [0.62, 4.42, 8.22]
    widths = [3.35, 3.35, 3.05]
    for (title, subtitle, tag), x, w in zip(flow_cards, xs, widths):
        card = add_card(slide, Inches(x), Inches(4.10), Inches(w), Inches(1.20), title, subtitle, fill="FFFFFF", accent="2563EB")
        add_badge(slide, Inches(x + 0.10), Inches(5.10), Inches(1.55 if title != "資料庫" else 1.35), tag, "EFF6FF", "1D4ED8")

    footer = textbox(slide, Inches(0.62), Inches(6.92), Inches(12.0), Inches(0.18), "這一頁只說明三個核心角色：前端入口、後端核心與資料庫。", 9, "64748B", align=PP_ALIGN.RIGHT)


def add_architecture_slide(slide):
    add_background(slide)
    textbox(slide, Inches(0.52), Inches(0.32), Inches(2.8), Inches(0.20), "System Overview", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.55), Inches(7.5), Inches(0.45), "知識庫系統架構", 24, "0F172A", bold=True)
    textbox(slide, Inches(0.52), Inches(0.98), Inches(9.4), Inches(0.42), "系統以使用者入口、知識庫處理層與資料庫層組成，透過清楚分工維持查詢效率與資料可維護性。", 12, "475569")

    # Three simple blocks
    left = add_card(slide, Inches(0.82), Inches(1.95), Inches(3.35), Inches(1.32), "前端入口", "小幫手卡片盒與聊天頁統一承接使用者操作。", fill="FFFFFF", accent="2563EB")
    mid = add_card(slide, Inches(4.98), Inches(1.88), Inches(3.35), Inches(1.46), "知識庫核心", "後端負責理解問題、查找資料並整理成正式回覆。", fill="FFFFFF", accent="16A34A")
    right = add_card(slide, Inches(9.10), Inches(1.95), Inches(3.05), Inches(1.32), "資料庫", "Neo4j 與 Qdrant 分別保存關聯脈絡與內容向量。", fill="FFFFFF", accent="D97706")

    connector(slide, Inches(4.20), Inches(2.63), Inches(4.98), Inches(2.63), "8EA2BE")
    connector(slide, Inches(8.33), Inches(2.63), Inches(9.10), Inches(2.63), "8EA2BE")

    add_badge(slide, Inches(1.02), Inches(3.58), Inches(1.72), "使用者操作入口", "EFF6FF", "1D4ED8")
    add_badge(slide, Inches(5.70), Inches(3.58), Inches(1.80), "問題理解與回覆", "ECFDF5", "047857")
    add_badge(slide, Inches(9.62), Inches(3.58), Inches(1.72), "知識資料儲存", "FEF3C7", "B45309")

    # Very small supporting note
    note = add_card(slide, Inches(1.02), Inches(4.65), Inches(11.40), Inches(1.38), "正式說明", "Neo4j 用於保存文件之間的關聯脈絡，適合回答跨文件、跨專案或需要關係推理的查詢；Qdrant 用於保存內容向量，適合做相似度搜尋與語意召回。兩者搭配後，可同時兼顧結構化關聯與內容檢索。", fill="FFFFFF", accent="2563EB")
    textbox(slide, Inches(0.62), Inches(6.90), Inches(12.0), Inches(0.18), "這張圖的重點是讓客戶理解：前端負責入口，後端負責處理，Neo4j 與 Qdrant 負責資料。", 9, "64748B", align=PP_ALIGN.RIGHT)


def add_roles_slide(slide):
    add_background(slide)
    textbox(slide, Inches(0.52), Inches(0.32), Inches(2.8), Inches(0.20), "Responsibilities", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.55), Inches(8.0), Inches(0.45), "各元件的正式說明", 24, "0F172A", bold=True)
    textbox(slide, Inches(0.52), Inches(0.98), Inches(9.5), Inches(0.36), "以下內容可直接作為客戶簡報時的口語說明，重點是讓對方知道每個元件負責什麼。", 12, "475569")

    roles = [
        ("前端入口", "前端提供統一的使用者入口，包含小幫手卡片盒、聊天頁與相關操作介面。使用者透過這一層發出查詢，並接收後端回傳的結果。", "2563EB"),
        ("知識庫核心", "後端負責理解使用者問題、決定查詢方式、整合檢索結果，並將資料整理成可閱讀的正式回覆。", "16A34A"),
        ("Neo4j", "Neo4j 保存文件、專案與章節之間的關聯，適合描述資料脈絡、交叉引用與結構化關係。", "D97706"),
        ("Qdrant", "Qdrant 保存文件內容向量，主要用於語意搜尋與相似內容召回，讓系統能更快找到相關資料段落。", "0EA5E9"),
    ]
    x_positions = [0.64, 6.88]
    y_positions = [1.55, 3.50]
    card_w = Inches(5.95)
    card_h = Inches(1.55)

    idx = 0
    for y in y_positions:
        for x in x_positions:
            title, body, accent = roles[idx]
            card = add_card(slide, x, Inches(y), card_w, card_h, title, body, fill="FFFFFF", accent=accent)
            idx += 1

    summary = add_card(
        slide,
        Inches(0.80),
        Inches(5.55),
        Inches(11.90),
        Inches(1.00),
        "整體結論",
        "前端負責入口與體驗，後端負責處理與整合，Neo4j 與 Qdrant 則分別支援關聯查詢與內容搜尋。這樣的分工可以讓系統在可維護性、擴充性與查詢效果之間取得平衡。",
        fill="FFFFFF",
        accent="2563EB",
    )


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

add_cover(prs.slides.add_slide(prs.slide_layouts[6]))
add_architecture_slide(prs.slides.add_slide(prs.slide_layouts[6]))
add_roles_slide(prs.slides.add_slide(prs.slide_layouts[6]))

prs.save(OUT)
print(f"saved {OUT}")
