from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = Path("/home/da40_ai_gb10/knowledge-base/query_examples_slide.pptx")


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value.replace("#", "").upper())


def add_shape(slide, shape_type, left, top, width, height, fill, line=None, transparency=0.0, radius=False):
    if radius:
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if transparency:
        shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def set_text(frame, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Sans CJK TC"):
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def textbox(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    set_text(shape.text_frame, text, size, color, bold=bold, align=align)
    return shape


def bullet_card(slide, left, top, width, height, index, category, text, accent):
    card = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"{index:02d}"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = rgb(accent)

    r = p.add_run()
    r.text = f"  {category}"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = rgb("64748B")

    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = text
    r2.font.name = "Noto Sans CJK TC"
    r2.font.size = Pt(12.2)
    r2.font.color.rgb = rgb("334155")
    return card


def add_background(slide):
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    for left, top, width, height, color in [
        (Inches(-0.5), Inches(-0.4), Inches(4.3), Inches(4.3), "DBEAFE"),
        (Inches(9.0), Inches(-0.2), Inches(4.8), Inches(4.2), "D1FAE5"),
        (Inches(8.9), Inches(4.7), Inches(4.1), Inches(2.8), "FEF3C7"),
    ]:
        deco = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        deco.fill.solid()
        deco.fill.fore_color.rgb = rgb(color)
        deco.fill.transparency = 0.78
        deco.line.fill.background()

    border = add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.22),
        Inches(0.18),
        Inches(12.89),
        Inches(7.14),
        fill="FFFFFF",
        line="CBD5E1",
        radius=True,
    )
    border.fill.transparency = 1


def add_cover(slide):
    add_background(slide)

    # Brand row
    logo = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.44), Inches(0.34), Inches(0.42), Inches(0.42), fill="2563EB", radius=True)
    set_text(logo.text_frame, "KB", 12, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)

    textbox(slide, Inches(0.95), Inches(0.31), Inches(2.8), Inches(0.18), "Knowledge Base Demo", 11, "2563EB")
    textbox(slide, Inches(0.95), Inches(0.50), Inches(3.2), Inches(0.18), "PPTX Title Slide Style", 9, "64748B")

    chips = [
        ("Query Examples", 9.28, 1.06),
        ("4G/5G + WiFi", 10.20, 1.20),
        ("Single Slide Cover", 11.15, 1.35),
    ]
    for label, x, width in chips:
        chip = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.35), Inches(width), Inches(0.30), fill="EFF6FF", line="BFDBFE", transparency=0.0, radius=True)
        set_text(chip.text_frame, label, 8.5, "1D4ED8", align=PP_ALIGN.CENTER)

    textbox(slide, Inches(0.62), Inches(0.88), Inches(1.8), Inches(0.18), "Opening Slide", 9, "2563EB")

    title = slide.shapes.add_textbox(Inches(0.60), Inches(1.08), Inches(7.1), Inches(1.40))
    set_text(title.text_frame, "4G/5G 與 WiFi 範例題型整理", 31, "0F172A", bold=True)

    subtitle = slide.shapes.add_textbox(Inches(0.62), Inches(2.30), Inches(6.4), Inches(1.15))
    set_text(
        subtitle.text_frame,
        "這份簡報會直接展示實際 query 範例，讓你在報告時先說清楚測試範圍，再逐頁帶出 4G/5G 與 WiFi 的代表題目。",
        15,
        "475569",
    )

    for i, (label, color) in enumerate([
        ("直接查數據", "2563EB"),
        ("完整 / 詳細", "16A34A"),
        ("比較 / 差異", "D97706"),
        ("泛問 / 相關文件", "64748B"),
    ]):
        left = 0.62 + i * 1.58
        width = 1.38 if i < 3 else 1.75
        pill_fill = {"2563EB": "DBEAFE", "16A34A": "DCFCE7", "D97706": "FEF3C7", "64748B": "E2E8F0"}[color]
        pill_line = {"2563EB": "BFDBFE", "16A34A": "BBF7D0", "D97706": "FDE68A", "64748B": "CBD5E1"}[color]
        pill_text = {"2563EB": "1D4ED8", "16A34A": "047857", "D97706": "B45309", "64748B": "475569"}[color]
        pill = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(3.54), Inches(width), Inches(0.34), fill=pill_fill, line=pill_line, radius=True)
        set_text(pill.text_frame, label, 10, pill_text, align=PP_ALIGN.CENTER)

    note = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.52), Inches(0.92), Inches(5.1), Inches(3.18), fill="FFFFFF", line="CBD5E1", radius=True)
    tf = note.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Presentation Note"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = rgb("2563EB")
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    r = p.add_run()
    r.text = "建議報告時先講這頁的定位：這不是搜尋結果頁，而是用來讓聽眾先理解接下來會怎麼測語意路由、文件命中與數據保留。"
    r.font.name = "Noto Sans CJK TC"
    r.font.size = Pt(13)
    r.font.color.rgb = rgb("475569")

    for idx, (kicker, head, body) in enumerate([
        ("4G/5G", "專案 + 指標 + 明細", "例如 Throughput、Handover、Case 對照與完整數值。"),
        ("WiFi", "型號 + 頻段 + 頻寬", "例如 2.4GHz / 5GHz / 6GHz 與原文保留。"),
    ]):
        x = 7.72 + (idx * 2.45)
        card = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.55), Inches(2.18), Inches(1.28), fill="FFFFFF", line="CBD5E1", radius=True)
        ctf = card.text_frame
        ctf.clear()
        ctf.word_wrap = True
        p = ctf.paragraphs[0]
        r = p.add_run()
        r.text = kicker
        r.font.name = "Noto Sans CJK TC"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = rgb("2563EB")
        p = ctf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = head
        r.font.name = "Noto Sans CJK TC"
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = rgb("0F172A")
        p = ctf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = body
        r.font.name = "Noto Sans CJK TC"
        r.font.size = Pt(10.5)
        r.font.color.rgb = rgb("475569")

    foot = slide.shapes.add_textbox(Inches(0.62), Inches(6.05), Inches(12.0), Inches(0.25))
    set_text(foot.text_frame, "下一頁開始就是實際的 10 條 4G/5G 範例題目與 10 條 WiFi 範例題目。", 9, "64748B", align=PP_ALIGN.RIGHT)


def add_examples_slide(slide, title, subtitle, items, accent, left_label, right_label):
    add_background(slide)

    textbox(slide, Inches(0.52), Inches(0.38), Inches(2.8), Inches(0.20), "Query Examples", 10, "2563EB")
    textbox(slide, Inches(0.52), Inches(0.62), Inches(7.0), Inches(0.55), title, 24, "0F172A", bold=True)
    textbox(slide, Inches(0.52), Inches(1.12), Inches(9.4), Inches(0.35), subtitle, 12.5, "475569")

    tag_fill = {"2E6DFF": "DBEAFE", "2FA05C": "DCFCE7"}[accent]
    tag_line = {"2E6DFF": "BFDBFE", "2FA05C": "BBF7D0"}[accent]
    tag_text = {"2E6DFF": "1D4ED8", "2FA05C": "047857"}[accent]
    tag1 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.48), Inches(0.48), Inches(1.20), Inches(0.32), fill=tag_fill, line=tag_line, radius=True)
    set_text(tag1.text_frame, left_label, 8.5, tag_text, align=PP_ALIGN.CENTER)
    tag2 = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.76), Inches(0.48), Inches(1.05), Inches(0.32), fill="EFF6FF", line="BFDBFE", radius=True)
    set_text(tag2.text_frame, right_label, 8.5, "1D4ED8", align=PP_ALIGN.CENTER)

    # two-column grid of 10 items
    col_left = Inches(0.58)
    col_right = Inches(6.78)
    y0 = Inches(1.68)
    card_w = Inches(5.92)
    card_h = Inches(0.70)
    gap = Inches(0.10)

    for idx, item in enumerate(items):
        col = 0 if idx < 5 else 1
        row = idx if idx < 5 else idx - 5
        left = col_left if col == 0 else col_right
        top = y0 + row * (card_h + gap)
        bullet_card(slide, left, top, card_w, card_h, idx + 1, item["type"], item["q"], accent)

    bottom = slide.shapes.add_textbox(Inches(0.58), Inches(6.92), Inches(12.0), Inches(0.20))
    set_text(bottom.text_frame, "這些題目可直接用於測試語意分類、文件命中與資料保留效果。", 8.8, "64748B", align=PP_ALIGN.RIGHT)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: cover
add_cover(prs.slides.add_slide(prs.slide_layouts[6]))

fourg_items = [
    {"type": "直接查數據", "q": "請查詢 SCU2140 的 Throughput 測試數據"},
    {"type": "完整 / 詳細", "q": "請顯示 SCU2060 詳細的 Throughput 測試數據"},
    {"type": "比較 / 差異", "q": "SCU2060、SCU2140、SCU5050 的 Throughput 有什麼差異？"},
    {"type": "案例列表", "q": "請列出 Throughput 底下有哪些 Case"},
    {"type": "單一 Case", "q": "SCU2140 的 Case 15 Throughput 數據是多少？"},
    {"type": "相關文件", "q": "請查詢 SCU2050 的相關報告數據"},
    {"type": "測試項目", "q": "請找出所有有 Latency 測試項目的報告"},
    {"type": "共通項目", "q": "SCU2140 和 SCU5050 共通的測試項目有哪些？"},
    {"type": "Performance", "q": "請查詢 SCU5050 的 Performance Test 數據"},
    {"type": "Handover", "q": "請查詢 SCU2060 的 Handover 測試結果如何？"},
]

wifi_items = [
    {"type": "直接查數據", "q": "請查詢 TP-Link Archer BE805 的 Throughput 測試數據"},
    {"type": "完整 / 詳細", "q": "請顯示 TP-Link Archer BE805 詳細的 Throughput 測試數據"},
    {"type": "2.4GHz", "q": "請查詢 TP-Link Archer BE805 的 2.4GHz Throughput 測試數據"},
    {"type": "5GHz", "q": "請查詢 TP-Link Archer BE805 的 5GHz Throughput 測試數據"},
    {"type": "6GHz", "q": "請查詢 TP-Link Archer BE805 的 6GHz Throughput 測試數據"},
    {"type": "80MHz", "q": "請查詢 TP-Link Archer BE805 的 5GHz 80MHz Throughput 測試數據"},
    {"type": "比較 / 差異", "q": "TP-Link Archer BE805 的 2.4GHz、5GHz、6GHz Throughput 有什麼差異？"},
    {"type": "知識問題", "q": "WiFi 7 和 WiFi 6 有什麼差別？"},
    {"type": "相關文件", "q": "TP-Link Archer BE805 的相關文件有哪些？"},
    {"type": "實務建議", "q": "企業環境 WiFi 頻道規劃要注意什麼？"},
]

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_examples_slide(
    slide2,
    "4G/5G 範例題目類型",
    "以下 10 條是實際可拿來測試 4G/5G 查詢路由、數值保留、Case 定位與相關文件召回的題目。",
    fourg_items,
    accent="2E6DFF",
    left_label="4G/5G",
    right_label="10 examples",
)

slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_examples_slide(
    slide3,
    "WiFi 範例題目類型",
    "以下 10 條是實際可拿來測試 WiFi 型號查詢、頻段抽取、頻寬保留與概念性問答的題目。",
    wifi_items,
    accent="2FA05C",
    left_label="WiFi",
    right_label="10 examples",
)

prs.save(OUT)
print(f"saved {OUT}")
